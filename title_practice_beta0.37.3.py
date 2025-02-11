# _*_ coding: utf-8 _*_
import pandas as pd
import pptx
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from openpyxl import load_workbook
from docx import Document
from docx.shared import Pt,Cm
import time
import os, re,datetime
from docx.oxml import parse_xml, register_element_cls
from docx.oxml.ns import nsdecls
from docx.oxml.shape import CT_Picture
from docx.oxml.xmlchemy import BaseOxmlElement, OneAndOnlyOne
import random
from tkinter import ttk
from tkinter import messagebox
import sys
import tkinter as tk
from PIL import ImageTk, Image
from win32com.client import Dispatch
import time
import shutil
import fitz
from PIL import Image

from pptx import Presentation
from pptx.util import Inches
from tkinter import simpledialog
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import logging
from functools import wraps
'''
version 0.1 2016/9
尝试利用python编程实现错题编号自动生成word文档
version 0.6 2018/6---2018/7
错一练三题目查询
一、确定该题的难度,查找相似难度的题目。
    练习题目难度在错题难度基础上浮动1
二、确定该题的研究对象，并在相似题目中进行比较，
    若研究对象完全相同，则确定为第一优先选择。
    若研究对象比错题多一个研究对象，则确定为第二选择；
    若研究对象比错题少一个研究对象，则确定为第三选择；
    其它则不符合筛选要求。将所有符合要求的题目编号和优先级别列入一个列表。
三、进行进一步筛选，如果所比较题目中有三个以上符合要求，则在符合要求
    题目中进行二次筛选，以优先级别作为首选条件，方法作为第二条件，过程作为
    第三条件
    
version 0.7 2018/8/24--2018/8/31
去除训练题目中重复的题目 
delete the repetitive titles that in practice titles.

version 0.8 2018/8/31--
给错题添加错误原因登记
一、在错题登记列表中提取个人信息，查找个人的错误题目
二、将题目的题号，图片呈现在界面中
三、查找每一题目的可能错误原因，由教师提前安排好（此处需要一个独立的功能）
    将该错误原因呈现在相应位置
四、学生选择该题的错误原因，
五、将学生选择的错误原因及题号存入excel表格中，每个学生成立一个单独的文件。
version 0.10 2018/11/9---2018/11/9
题型分类第二次，每章分成若干题型，并找出每种题型对应的训练题

version 0.11 2018/11/15---2018//11/19
1/题型分类第二次第二步，将每种题型按照难度排序，若总题目大于9个，则每类3个，
若再大于10个，则每类3个。
2/并将答案附于每个学生的题目后面。
3/每一类题型指定一个题目作为错题，训练题目从这个题目寻找，学生只登记题型编号。

version 0.12 2018/11/20---2018//11/20
1/将原题重做和错一练三合并到一个功能中，而不必采用两次操作。
2/添加整理人和整理时间
3/为原题重做添加答案

version 0.13 2018/12/19---2018/12/21
1/将所有的题目登记到一个文件中，在程序内实现题目分类并编号，并且可以多次登记

version 0.14 2018/12/27---2018/1/2
1/将考试试卷的选择题答案从excel文件中直接复制到错题word文档中。
2/去掉word文档的第一页（空白页）
3/去掉错题和训练题目中的重复题目，
4/可以自动确定是否要三维设计的题目出现在训练题目中，而不需要通过excel表格单列。
5/将生成的word文件名改为2016级X班的样式。
6/如果没有考试卷子，将生成的两个文件变成一个文件。

version 0.15  2019/1/15-2019/1/23
3/如果没有训练题，添加其原始录入题号
4/在错题中标注题型，目录就可以只写错题，便于题目分类查看。
5/将每一种题型的训练题进行储存，而不是每个人均进行寻找。
6/实现学生任务完成情况统计，并根据完成度安排下一次打印。

version 0.16 2020/9/7--2020/9/30
1/将所有训练题目进行统一编号，方便学生查询答案
2/利用邮件进行自动化处理，首先自动收取邮件得到错题编号，然后分析，再自动发送含有word文档附件的邮件。
3/多个邮件时，对不同的邮件采取不同的文件夹接收文件，再单独处理，然后发送word文档附件的邮件
4/由于是邮箱发送，将学生的完成任务情况筛选功能剔除。

version 0.17  2021/5/21--2021/5/21
1/将题目的题头更改为题号+难度两项内容，题目的类型单独列出,如果没有类型，则暂时忽略。

version 0.18  2021/12/19--2021
1/将每种题型的训练题目编号置于excel数据库中，需要的时候直接从数据库中提取。而不需要每次都去查找。
未完成

version 0.19 2022/10/13-2022/10/29
1.插入的题目图片放在题目结尾，并置于右侧。
2.修改选择题的选项顺序，并更改答案。
3.去掉人员筛选功能，
4.答案编写时的格式微调，更节约纸张。

version 0.20 2022/10/31-
1.实现程序的简单封装

version 0.21 2022/11/1-2022/11/3
1.实现简单的可视化和进度条显示 

version 0.22 2022/11/22-2022/12/2
1、为word文档添加页码
2.实现word文档题目图片的格式转换，如果有两张及以上，用嵌入式，如果只有一张，则用浮动式

version 0.23 2022/12/2
1/实现简单的封装，可以到其他电脑上使用，解决封装时的路径问题。（未解决文档合并）


version 0.24 2022/12/29-2023/1/2
1/将word文档格式插入的图片以嵌入式格式插入

version 0.25 2023/2/10-
1/将word文档合并调整为复制word文字(含格式)和图片.

version 0.26 2023/4/6
1/给答案二字添加标题级别，可以在视图窗格中显示。
2/可以添加试卷，题目2023304061。

version 0.28 2023/12/29--2023/1/3
1/复制word文档时，采用xml内部结构复制，即一个字一个字复制。
2/选择题选项随机排列后，答案进行相应变化，word文档内部字块的替换。
3/word文档自动分栏，以节约纸张

version 0.30 2023/1/4-2023/1/21
1/学生登记界面，学生登记自己的错题，并生成word文档至网络移动云盘，方便打印
2/生成快速登记界面，学生勾选题目，提升效率
3/计算题自动识别，题目处留五行空白
4/实验题的选择题同样更改选项

version 0.35 2023/1/21-2023/1/28
1/生成word题库的自动化程序，

version 0.36 2023/1/28-2023/3/15
1/加入知识点编号，学生可以直接选择知识点编号（不能勾选），生成相应的图片训练题目。
# 2/直接调用知识点对应的题目编号，不用每次都要查找，如果更新，则删掉原excel就可以了。
3/勾选题号时可以预览图片，3秒后自动关闭图片。
4/生成word文档的同时，生成ppt课件，供学生在课堂进行限时训练，培养时间的感觉

version 0.36.13 2023/3/16-2025/1/6
2/自动更改选择题的文字答案
1/按题目从中选择，进行组卷
3/组卷时，选择题长度小于半行的，几个选项合在一行。
5/BUG修补，故A选项正确，答案更正
6/页码位置放到最底部，自动分成两栏，
7/分节失败，采用加两个分页符区分学生
8/添加组卷功能，学生可以通过下拉框选择章节，也可以直接输入编号，并从图片预览中添加题目至试卷中，生成word文档
9/生成word文档添加两个课时的补充功能
10/记录错误日志，第一次尝试用函数装饰器

version0.37.1 2025/1/6-2025/1/21
1/下拉框选择时，自动更新下级下拉框内容，
2/组卷功能中，初始状态下，自动加载所选章节的图片

version0.37.2--
1/组卷和班级统一登记的错题 分位置存放
2/将所有文件存放到安装必备文件中
'''

def handle_exceptions(logger):
    def decorator(func):
        @wraps(func)
        def wrapper(*args, **kwargs):

            try:
                return func(*args, **kwargs)
            except Exception as e:
                logging.basicConfig(filename='error.log', level=logging.ERROR,
                                    format='%(asctime)s - %(levelname)s - %(message)s')
                logger.exception("异常发生: %s", e)
        return wrapper
    return decorator

logger = logging.getLogger(__name__)
def function_print(name):
    """提示某个函数运行中"""
    print("%s函数运行中"%name)
def excel(math):
    """
    #获取excel表的所有信息，按‘行’存入content中，
    :param math:
    :return:
    """
    from openpyxl import load_workbook
    content=[]
    try:
        wb=load_workbook(math)
    except Exception as e:
        print("加载文件失败:",e,math)
        tk.messagebox.showinfo(title="HI", message="加载文件失败:%s"%math)
    sheet=wb.get_sheet_by_name('Sheet1')
    for i in range(sheet.max_row):
        # print(sheet.max_row)
        content.append([])
        for j in range(sheet.max_column):
            content[i].append(sheet.cell(row=i+1,column=j+1).value)
    return content
def findcol(name):
    """#查找指定内容所在的列，返回列号"""
    # print("查找指定内容所在的列，返回列号")
    for i in range(colmax):
        k = chr(ord('A') + i)
        # print(name,"i,k:",i,k)
        if name==sheet['%s1'%k].value:
            # print("'%s'该列存在于第%s列"%(name,k))
            return k
    # print("'%s'该列不存在"%name)
    return None
def makedir(path):
    """尝试创建文件夹，有则不创建"""
    try:
        os.mkdir(path)
    except:
        pass
def findtitle(title):
    """#查找题目所在的行"""
    # print("#查找题目所在的行")
    for i in range(rowmax):

        if title==sheet['%s%d'%(bianhao,i+1)].value:
            # print("'%s'该列存在于第%s行" % (title,i+1))
            return i+1
    print("'%s'该题目不存在（类型查找）" % title)
    return None
def finddif(titlewrong):
    """#查找题目的难度"""
    # print("#查找题目的难度")
    row0 = findtitle(titlewrong)
    col0 = findcol('合计')
    # print("titlewrong,row0,col0:",titlewrong,row0,col0)
    try:
        dif = sheet['%s%d' % (col0, row0)].value
    except:
        dif=0
    # print('题目%s的难度为：%s'%(titlewrong,dif))
    return dif
def findobj(titlewrong,name):
    """#查找题目的研究目标"""
    # print("#########查找题目的:",name)
    row=findtitle(titlewrong)
    col=findcol(name)
    try:
        obj=sheet['%s%d' % (col,row)].value
    except:
        obj='10000'
    # print('题目%s的研究%s为：%s'%(titlewrong,name,obj))
    # obj2=obj
    chinese = eng_chinese(name,obj)
    return chinese
    # return dict.keys()#'''#以字典形式返回题目的研究对象
def find_class_and_unit(name):
    """#查找题目名称所含有的班级和单元"""
    # print("name:",name)
    names=['年','月','日',"级",'班']
    position=[]
    for i in names:
        if i in name:
            # print('name.find(%s)'%i,name.find(i),)
            position.append(name.find(i))
            # print('%s,position:' % i,position)
    # print('%s,position:'%i,position)
    information = {}
    information["year"]=name[:position[0]]
    information["month"]=name[position[0] + 1:position[1]]
    information["day"]=name[position[1] + 1:position[2]]
    information["grade"]=name[position[2] + 1:position[3]]
    information["classs"]=name[position[3] + 1:position[4]]
    # print("information:",information)
    return information
def eng_chinese(name,eng):
    """#查找题目研究对象、运动过程等题目信息中字母所对应的中文"""
    chinese='类型暂不清楚'
    for j in range(len(content[0])):
        if name == content[0][j]:#在第一行查找类别
            for i in range(len(content)):
                # print(type(eng))
                # print(eng)
                # print(type(content[i][j + 1]))
                # print(content[i][j + 1])
                # try:
                #     # eng=int(eng)
                #     content[i][j + 1]=int(content[i][j + 1])
                # except:
                #     pass
                # print("type(eng)",type(eng))
                # print("type(content[i][j + 1])",type(content[i][j + 1]))
                if eng==content[i][j+1]:

                    chinese=content[i][j]
                    # print('chinese:',chinese)
                    return chinese
    if chinese=='类型暂不清楚':
        # print('name and eng is',name,eng,'未找到对应中文名称')
        pass

    return chinese
def que_type_title(name,que_type):
    """#查找错一练三的类型所对应的题目"""
    # print("name,que_type",name,que_type)
    # print("函数que_type_title运行中")
    chinese = que_type
    for j in range(len(content[0])):
        if name == content[0][j]:  # 在第一行查找类别
            for i in range(len(content)):
                # print("content[i][j + 1]",content[i][j + 1])
               if  str(que_type)== str(content[i][j + 1]):
                    titl = content[i][j+2]
                    # print('quetion_type and title is',name,titl,)
                    return titl
    if chinese == que_type:
        print(name,que_type,'未找到对应题目')
        pass

    return chinese
def priority(obj,obj2):
    """#比较两个题目研究对象的优先级次序"""
    #print('obj:',obj,'obj2:',obj2)
    pri=1
    for i in obj:#遍历obj的所有研究对象
        #print('i:',i)
        if i in obj2:
            #print('obj[%s]:%d'%(i,obj[i]),'obj2[%s]:%d'%(i,obj2[i]))
            #print('i in obj2')
            if obj[i]==obj2[i]:
                pri+=0
            elif obj[i]==obj2[i]-1:
                pri+=1
            elif obj[i]==obj2[i]-2:
                pri+=2
            else:
                pri=0
                return pri
        else:
            pri=0
            return pri
    for i in obj2:
        if i in obj:
            pass
        else:
            pri+=obj2[i]
    return pri
def titlepractice(titlewrong):
    """#查找错题的训练题目，并返回合适训练题目的编号及其优先级。"""
    # print("titlepractice 运行中")
    # print("查找错题的训练题目，并返回合适训练题目的编号及其优先级。")
    # row0 = findtitle(titlewrong)

    # #判断是否存在titlepractice库
    # try:
    #     question_type_train_titles=pd.DataFrame(
    #         pd.read_excel(r"安装必备文件/知识点训练题.xlsx"))
    # except:
    #     question_type_train_titles=pd.DataFrame()
    #     question_type_train_titles.to_excel(r"安装必备文件/知识点训练题.xlsx")
    #
    # #在知识点训练题.xlsx中查找训练题，如果没有，则添加一个，如果有则直接返回
    col0=findcol('合计')
    # print("1")
    # print("titlewrong",titlewrong)
    # dif=int(finddif(titlewrong))    #print('错题%s的难度为：%d'%(titlewrong,dif))
    obj=findobj(titlewrong,'对象')#查询该题目对应的多种指标
    # pro=findobj(titlewrong,'过程')
    # sta=findobj(titlewrong,'状态')
    # met=findobj(titlewrong,'方法')
    # print('错题%s的研究对象为:%s,'%(titlewrong,obj))
    # print('错题%s的运动过程为:%s,' % (titlewrong, obj))
    sumtitle = 0
    #titlepras = []
    titlepras={}
    #titlepras[titlewrong]=0
    # titlepras['错题'] = [titlewrong]
    titlepras['错题'] = [titlewrong]
    # for i in range(12):
        # titlepras['训练题第%d优先'%i]=[]
    titlepras['训练题']=[]
    # print('titlepras:',titlepras)

    #查找合适难度的题目,一级判断，对研究对象的优先级进行判断
    for i in range(rowmax):#在所有题目中查找所需要的题目
        dif2 = sheet['%s%d' % (col0, i+1)].value
        # title=sheet['%s%d' % (bianhao, i+1)].value

        try:
            # if (dif - 2)<dif2 < (dif + 3):  # 查找合适难度的题目'''
            #if 0<dif2 < 4:  # 查找合适难度的题目，难度设定为中档题和难题'''
            if 0<dif2 < 15:  # 查找合适难度的题目，难度按照序列从1到15排布'''
                titlepra=sheet['%s%d' % (bianhao,i+1)].value
                # print('符合题目难度要求：题目%s,难度%s'%(titlepra,dif2))
                obj2=findobj(titlepra,"对象")#查询本题目的研究对象
                # print('研究对象%s'%obj2)
                if obj==obj2 and titlewrong!=titlepra and obj2!=10000:
                # if obj==obj2 and titlewrong!=titlepra and ("201903" in titlepra):
                    # titlepras['训练题第1优先' ].append(titlepra)
                    titlepras['训练题' ].append(titlepra)
        except:
            pass
    # print('titlepras:',titlepras)

    def titlepras_del(titlepras,number):
        """将每一难度的题目数量控制在number个"""
        number+=1
        number1 = number2=number3=0
        dif1 = dif2=dif3=1
        # print("titlepras_del运行中")
        titlepras2=[]
        for i in titlepras['训练题']:
            # print("i",i)
            dif = finddif(i)
            if dif == 1 :
                # pass
                number1+=1
                if number1<number:
                    titlepras2.append(i)


            elif dif == 2 :
                number2+=1
                if number2<number:
                    titlepras2.append(i)

            elif dif == 3 :
                number3+=1
                if number3<number:
                    titlepras2.append(i)
        titlepras['训练题']=titlepras2
    # print("file_name",file_name)

    title_pra_number = 9
    try:
        # if "模拟" in file_name or "模拟" :
        if "模拟" in file_name :
        # 判断是否为题库建立，如果有模拟，则是建立题库，没有则是实际使用
            title_pra_number = 100
    except:pass

    #print("titlepras['训练题']:",titlepras['训练题'],)
    #print("len(titlepras['训练题']):",len(titlepras['训练题']),)

    if len(titlepras['训练题']) > title_pra_number:
        print("训练题目数量超过9题")
        titlepras_del(titlepras,3)


    #print('删除多余题目以后titlepras:',titlepras)

    return titlepras
    #print('题目%s的难度为：%s'%(titlewrong,dif))
def title_replace(title):
    # print("title_replace 函数运行中")
    # print(title_replace_path)
    content=excel(title_replace_path)
    # print("content",content)
    for i in range(len(content)):
        # print(content[i][1])
        if title == content[i][1]:
            title = content[i][0]
            print(content[i][0],"成功替换",content[i][1])
            return title

#CT_Anchor new_pic_anchor add_float_picture的目的为添加word文档题目并且将图片转化四周环绕型
class CT_Anchor(BaseOxmlElement):
    """
    ``<w:anchor>`` element, container for a floating image.
    """
    extent = OneAndOnlyOne('wp:extent')
    docPr = OneAndOnlyOne('wp:docPr')
    graphic = OneAndOnlyOne('a:graphic')

    @classmethod
    def new(cls, cx, cy, shape_id, pic, pos_x, pos_y):
        """
        Return a new ``<wp:anchor>`` element populated with the values passed
        as parameters.
        """
        anchor = parse_xml(cls._anchor_xml(pos_x, pos_y))
        anchor.extent.cx = cx
        anchor.extent.cy = cy
        anchor.docPr.id = shape_id
        anchor.docPr.name = 'Picture %d' % shape_id
        anchor.graphic.graphicData.uri = (
            'http://schemas.openxmlformats.org/drawingml/2006/picture'
        )
        anchor.graphic.graphicData._insert_pic(pic)
        return anchor

    @classmethod
    def new_pic_anchor(cls, shape_id, rId, filename, cx, cy, pos_x, pos_y):
        """
        Return a new `wp:anchor` element containing the `pic:pic` element
        specified by the argument values.
        """
        pic_id = 0  # Word doesn't seem to use this, but does not omit it
        pic = CT_Picture.new(pic_id, filename, rId, cx, cy)
        anchor = cls.new(cx, cy, shape_id, pic, pos_x, pos_y)
        anchor.graphic.graphicData._insert_pic(pic)
        return anchor

    @classmethod
    def _anchor_xml(cls, pos_x, pos_y):
        return (
                '<wp:anchor distT="100" distB="100" distL="0" distR="0" simplePos="0" relativeHeight="100" \n'
                # '           square="1" locked="1" layoutInCell="1" allowOverlap="1" \n'
                '           behindDoc="1" locked="1" layoutInCell="1" allowOverlap="1" \n'
                '           %s>\n'
                '  <wp:simplePos x="0" y="0"/>\n'
                '<wp:positionH relativeFrom="margin">\n'  # 个人添加
                '<wp:align>right</wp:align>\n'  # 个人添加
                '</wp:positionH>\n'  # 个人添加
                '<wp:positionV relativeFrom="line">\n'  # 个人添加
                '<wp:align>top</wp:align>\n'  # 个人添加
                '</wp:positionV>\n'  # 个人添加

                '  <wp:extent cx="914400" cy="914400"/>\n'
                # '  <wp:wrapNone/>\n'
                '  <wp:wrapSquare wrapText="bothSides"/>\n'
                ' <wp:effectExtent r="4445" b="0" t="0" l="0"/>'  # 个人添加
                '  <wp:docPr id="666" name="unnamed"/>\n'
                '  <wp:cNvGraphicFramePr>\n'
                '    <a:graphicFrameLocks noChangeAspect="1"/>\n'
                '  </wp:cNvGraphicFramePr>\n'
                '  <a:graphic>\n'
                '    <a:graphicData uri="URI not set"/>\n'
                '  </a:graphic>\n'
                '</wp:anchor>' % (nsdecls('wp', 'a', 'pic', 'r'),)
        )# refer to docx.parts.story.BaseStoryPart.new_pic_inline

# def new_pic_anchor(part, image_descriptor, width, height, pos_x, pos_y):
#     # refer to docx.text.run.add_picture
#     """Return a newly-created `w:anchor` element.
#     The element contains the image specified by *image_descriptor* and is scaled
#     based on the values of *width* and *height*.
#     """
#     rId, image = part.get_or_add_image(image_descriptor)
#     cx, cy = image.scaled_dimensions(width, height)
#     shape_id, filename = part.next_id, image.filename
#     return CT_Anchor.new_pic_anchor(shape_id, rId, filename, cx, cy, pos_x, pos_y)

# def add_float_picture(p, image_path_or_stream, width=None, height=None, pos_x=0, pos_y=0):
#     """Add float picture at fixed position `pos_x` and `pos_y` to the top-left point of page.
#     """
#     run = p.add_run()
#     anchor = new_pic_anchor(run.part, image_path_or_stream, width, height, pos_x, pos_y)
#     run._r.add_drawing(anchor)

# class CT_Anchor2(BaseOxmlElement):
#     """
#     ``<w:anchor>`` element, container for a floating image.
#     """
#     extent = OneAndOnlyOne('wp:extent')
#     docPr = OneAndOnlyOne('wp:docPr')
#     graphic = OneAndOnlyOne('a:graphic')
#
#     @classmethod
#     def new(cls, cx, cy, shape_id, pic, pos_x, pos_y):
#         """
#         Return a new ``<wp:anchor>`` element populated with the values passed
#         as parameters.
#         """
#         anchor = parse_xml(cls._anchor_xml(pos_x, pos_y))
#         anchor.extent.cx = cx
#         anchor.extent.cy = cy
#         anchor.docPr.id = shape_id
#         anchor.docPr.name = 'Picture %d' % shape_id
#         anchor.graphic.graphicData.uri = (
#             'http://schemas.openxmlformats.org/drawingml/2006/picture'
#         )
#         anchor.graphic.graphicData._insert_pic(pic)
#         return anchor
#
#     @classmethod
#     def new_pic_anchor(cls, shape_id, rId, filename, cx, cy, pos_x, pos_y):
#         """
#         Return a new `wp:anchor` element containing the `pic:pic` element
#         specified by the argument values.
#         """
#         pic_id = 0  # Word doesn't seem to use this, but does not omit it
#         pic = CT_Picture.new(pic_id, filename, rId, cx, cy)
#         anchor = cls.new(cx, cy, shape_id, pic, pos_x, pos_y)
#         anchor.graphic.graphicData._insert_pic(pic)
#         return anchor
#
#     @classmethod
#     def _anchor_xml(cls, pos_x, pos_y):
#         return (
#                 '<wp:inline distT="100" distB="100" distL="0" distR="0" simplePos="0" relativeHeight="100" \n'
#                 # '           square="1" locked="1" layoutInCell="1" allowOverlap="1" \n'
#                 '           behindDoc="1" locked="1" layoutInCell="1" allowOverlap="1" \n'
#                 '           %s>\n'
#                 '  <wp:simplePos x="0" y="0"/>\n'
#                 # '<wp:positionH relativeFrom="margin">\n'  # 个人添加
#                 # '<wp:align>right</wp:align>\n'  # 个人添加
#                 # '</wp:positionH>\n'  # 个人添加
#                 # '<wp:positionV relativeFrom="line">\n'  # 个人添加
#                 # '<wp:align>top</wp:align>\n'  # 个人添加
#                 # '</wp:positionV>\n'  # 个人添加
#
#                 '  <wp:extent cx="914400" cy="914400"/>\n'
#                 # '  <wp:wrapNone/>\n'
#                 '  <wp:wrapSquare wrapText="bothSides"/>\n'
#                 ' <wp:effectExtent r="4445" b="0" t="0" l="0"/>'  # 个人添加
#                 '  <wp:docPr id="666" name="unnamed"/>\n'
#                 '  <wp:cNvGraphicFramePr>\n'
#                 '    <a:graphicFrameLocks noChangeAspect="1"/>\n'
#                 '  </wp:cNvGraphicFramePr>\n'
#                 '  <a:graphic>\n'
#                 '    <a:graphicData uri="URI not set"/>\n'
#                 '  </a:graphic>\n'
#                 '</wp:inline>' % (nsdecls('wp', 'a', 'pic', 'r'),)
#         )# refer to docx.parts.story.BaseStoryPart.new_pic_inline

# def new_pic_anchor2(part, image_descriptor, width, height, pos_x, pos_y):
#     # refer to docx.text.run.add_picture
#     """Return a newly-created `w:anchor` element.
#     The element contains the image specified by *image_descriptor* and is scaled
#     based on the values of *width* and *height*.
#     """
#     rId, image = part.get_or_add_image(image_descriptor)
#     cx, cy = image.scaled_dimensions(width, height)
#     shape_id, filename = part.next_id, image.filename
#     return CT_Anchor2.new_pic_anchor(shape_id, rId, filename, cx, cy, pos_x, pos_y)

# def add_float_picture2(p, image_path_or_stream, width=None, height=None, pos_x=0, pos_y=0):
#     """Add float picture at fixed position `pos_x` and `pos_y` to the top-left point of page.
#     """
#     run = p.add_run()
#     anchor = new_pic_anchor2(run.part, image_path_or_stream, width, height, pos_x, pos_y)
#     run._r.add_drawing(anchor)
# refer to docx.oxml.__init__.py
register_element_cls('wp:anchor', CT_Anchor)

#插入题目的word文档的相关处理办法，删除空的段落，添加文档，更改选项等
def wordtitle_answer_find_excel(title):
    """在excel文件中查找考试选择题的答案，以文本的形式存入word文档中"""
    # print("wordtitle_answer_find(title):")
    content = excel(r"%s\sys_files\试卷选择题答案.xlsx" % main_path)
    # print(content)
    # answer = 0
    for i in range(len(content)):  # 在第一列中查找题号
        # print("title[0:9]", title[0:10], title)
        if title[1:10] == content[i][0]:
            for j in range(len(content[0])):
                if title[10:] == content[0][j]:
                    answer = content[i][j]
                    if answer == None:
                        answer = "未找到答案"
                        print("%s 未找到答案" % title)
    return answer
def copy_docx_para_math(paragraph):
    """复制整段文字中的公式"""
    paragraph.omaths=[]
    string1 = paragraph._element.xml
    # print(string1.find("<m:r>"))
    # print(string1.find("<m:oMath>") + len("<m:oMath>"))
    # print(string1[string1.find("<m:oMath>")+len("<m:oMath>")+1:string1.find("</m:oMath>")])
    # print(re.findall("\n", string1))
    string_split = re.split("\n", string1)
    math = 0

    def omath(space_number, ij):
        """在公式中加入各级元素"""
        # print("begin to copy omath")
        i = space_number
        if "</" not in ij:
            if "/>" not in ij:
                oMath = OxmlElement(ij[i + 1:-1])  # creates a new element
                return oMath
            elif "w:rFonts" in ij:
                oMath = OxmlElement('w:rFonts')  # creates a new element
                oMath.set(qn('w:ascii'), u'Times New Roman')  # sets attribute on element
                oMath.set(qn('w:eastAsia'), u'宋体')
                return oMath
            elif "m:degHide" in ij:
                oMath = OxmlElement('m:degHide')  # creates a new element
                oMath.set(qn('m:val'), '1')  # sets attribute on element
                return oMath
            elif "=" in ij:
                string_split2 = re.split(" ", ij)

                # print("string_split2[space_number+1,-1]:",string_split2)
                # print("string_split2[space_number+1,-1]:",string_split2[space_number+1,-1])
                for string_split2_ij in string_split2:

                    if len(string_split2_ij)>2:
                        # print("string_split2_ij:",string_split2_ij)
                        if "=" not in string_split2_ij:

                            oMath = OxmlElement(string_split2[space_number][1:])  # creates a new element  "w:sz"
                        else:

                            position=string_split2_ij.find("=")
                            # print("string_split2_ij[0:position]:",string_split2_ij[0:position])
                            # print(string_split2_ij[position+2:-3])
                            oMath.set(qn(string_split2_ij[0:position]),string_split2_ij[position+2:-3])  # sets attribute on element

                return oMath
        elif "<m:t>" in ij:
            # print("ij[i + 1:-1]:", ij[i + 5:-6])
            oMath = OxmlElement("m:t")  # creates a new element
            oMath.text = ij[i + 5:-6]
            return oMath


    for ij in string_split:
        # print("ij:",ij)
        if "<m:oMath>" in ij:
            math = 1
            omath2 = omath(2, ij)
        if "</m:oMath>" in ij:
            math = 0
            paragraph.omaths.append(omath2)
        if math == 1 and ("</" not in ij or "<m:t>" in ij):
            space_number = len(re.match(r'^(\s*)(.*)', ij).group(1))  # 行首空格的数量
            # print(space_number, ij)
            # print(space_number)

            # 按照逻辑层次加入所有的元素
            #
            # if space_number==4 and "</" not in ij:
            #     omath4=omath(space_number,ij)
            #     omath2.append(omath4)
            # if space_number == 6 and "</" not in ij:
            #     omath6 = omath(space_number, ij)
            #     omath4.append(omath6)
            # if space_number == 8 and "</" not in ij:
            #     omath8 = omath(space_number, ij)
            #     omath6.append(omath8)
            # if space_number == 10 and "</" not in ij:
            #     omath10 = omath(space_number, ij)
            #     omath8.append(omath10)
            # if space_number == 12 and "</" not in ij:
            #     omath12 = omath(space_number, ij)
            #     omath10.append(omath12)
            # if space_number == 14 and "</" not in ij:
            #     omath14 = omath(space_number, ij)
            #     omath12.append(omath14)
            # if space_number == 16 and "</" not in ij:
            #     omath16 = omath(space_number, ij)
            #     omath14.append(omath16)
            try:

                if space_number == 4:
                    omath4 = omath(space_number, ij)
                    omath2.append(omath4)
                if space_number == 6:
                    omath6 = omath(space_number, ij)
                    omath4.append(omath6)
                if space_number == 8:
                    omath8 = omath(space_number, ij)
                    omath6.append(omath8)
                if space_number == 10:
                    omath10 = omath(space_number, ij)
                    omath8.append(omath10)
                if space_number == 12:
                    omath12 = omath(space_number, ij)
                    omath10.append(omath12)
                if space_number == 14:
                    omath14 = omath(space_number, ij)
                    omath12.append(omath14)
                if space_number == 16:
                    omath16 = omath(space_number, ij)
                    omath14.append(omath16)
            except Exception as e:
                print("omath对象插入失败：", e.__traceback__.tb_lineno, e)
            # if "<m:t>" in ij and "</m:t" in ij:
            # if "</m:oMath>" in ij:

            #     if "m:t" not in ij:
            #         print(ij[5:-1])
            #         oMath4 = OxmlElement(ij[5:-1])  # creates a new element
            #
            #     else:
            #         oMath4 = OxmlElement(ij[5:-1])  # creates a new element
            #         para_new._element.append(oMath4)
            #         oMath4.text=ij[5:-6]
            #     para_new._element.append(oMath4)
            # if "</" not in ij:
            #     if "=" not in ij:
            #         print(ij)
            # print()
            # print(ij)
    return paragraph.omaths
# def copy_docx_para_math(paragraph):
#     """复制整段文字中的公式"""
#     paragraph.omaths=[]
#     string1 = paragraph._element.xml
#     # print(string1.find("<m:r>"))
#     # print(string1.find("<m:oMath>") + len("<m:oMath>"))
#     # print(string1[string1.find("<m:oMath>")+len("<m:oMath>")+1:string1.find("</m:oMath>")])
#     # print(re.findall("\n", string1))
#     string_split = re.split("\n", string1)
#     math = 0
#
#     def omath(space_number, ij):
#         """在公式中加入各级元素"""
#         # print("begin to copy omath")
#         i = space_number
#         if "</" not in ij:
#             if "/>" not in ij:
#                 oMath = OxmlElement(ij[i + 1:-1])  # creates a new element
#                 return oMath
#             elif "w:rFonts" in ij:
#                 oMath = OxmlElement('w:rFonts')  # creates a new element
#                 oMath.set(qn('w:ascii'), u'Times New Roman')  # sets attribute on element
#                 oMath.set(qn('w:eastAsia'), u'宋体')
#                 return oMath
#             elif "m:degHide" in ij:
#                 oMath = OxmlElement('m:degHide')  # creates a new element
#                 oMath.set(qn('m:val'), '1')  # sets attribute on element
#                 return oMath
#             elif "=" in ij:
#                 string_split2 = re.split(" ", ij)
#
#                 # print("string_split2[space_number+1,-1]:",string_split2)
#                 # print("string_split2[space_number+1,-1]:",string_split2[space_number+1,-1])
#                 for string_split2_ij in string_split2:
#
#                     if len(string_split2_ij)>2:
#                         # print("string_split2_ij:",string_split2_ij)
#                         if "=" not in string_split2_ij:
#
#                             oMath = OxmlElement(string_split2[space_number][1:])  # creates a new element  "w:sz"
#                         else:
#
#                             position=string_split2_ij.find("=")
#                             # print("string_split2_ij[0:position]:",string_split2_ij[0:position])
#                             # print(string_split2_ij[position+2:-3])
#                             oMath.set(qn(string_split2_ij[0:position]),string_split2_ij[position+2:-3])  # sets attribute on element
#
#                 return oMath
#         elif "<m:t>" in ij:
#             # print("ij[i + 1:-1]:", ij[i + 5:-6])
#             oMath = OxmlElement("m:t")  # creates a new element
#             oMath.text = ij[i + 5:-6]
#             return oMath
#
#
#     for ij in string_split:
#         print("ij:",ij)
#         if "<m:oMath>" in ij:
#             math = 1
#             omath2 = omath(2, ij)
#         if "</m:oMath>" in ij:
#             math = 0
#             paragraph.omaths.append(omath2)
#         if math == 1 and ("</" not in ij or "<m:t>" in ij):
#             space_number = len(re.match(r'^(\s*)(.*)', ij).group(1))  # 行首空格的数量
#             # print(space_number, ij)
#             # print(space_number)
#
#             # 按照逻辑层次加入所有的元素
#             #
#             # if space_number==4 and "</" not in ij:
#             #     omath4=omath(space_number,ij)
#             #     omath2.append(omath4)
#             # if space_number == 6 and "</" not in ij:
#             #     omath6 = omath(space_number, ij)
#             #     omath4.append(omath6)
#             # if space_number == 8 and "</" not in ij:
#             #     omath8 = omath(space_number, ij)
#             #     omath6.append(omath8)
#             # if space_number == 10 and "</" not in ij:
#             #     omath10 = omath(space_number, ij)
#             #     omath8.append(omath10)
#             # if space_number == 12 and "</" not in ij:
#             #     omath12 = omath(space_number, ij)
#             #     omath10.append(omath12)
#             # if space_number == 14 and "</" not in ij:
#             #     omath14 = omath(space_number, ij)
#             #     omath12.append(omath14)
#             # if space_number == 16 and "</" not in ij:
#             #     omath16 = omath(space_number, ij)
#             #     omath14.append(omath16)
#             try:
#
#                 if space_number == 4:
#                     omath4 = omath(space_number, ij)
#                     omath2.append(omath4)
#                 if space_number == 6:
#                     omath6 = omath(space_number, ij)
#                     omath4.append(omath6)
#                 if space_number == 8:
#                     omath8 = omath(space_number, ij)
#                     omath6.append(omath8)
#                 if space_number == 10:
#                     omath10 = omath(space_number, ij)
#                     omath8.append(omath10)
#                 if space_number == 12:
#                     omath12 = omath(space_number, ij)
#                     omath10.append(omath12)
#                 if space_number == 14:
#                     omath14 = omath(space_number, ij)
#                     omath12.append(omath14)
#                 if space_number == 16:
#                     omath16 = omath(space_number, ij)
#                     omath14.append(omath16)
#             except Exception as e:
#                 print("omath对象插入失败：", e.__traceback__.tb_lineno, e)
#             # if "<m:t>" in ij and "</m:t" in ij:
#             # if "</m:oMath>" in ij:
#
#             #     if "m:t" not in ij:
#             #         print(ij[5:-1])
#             #         oMath4 = OxmlElement(ij[5:-1])  # creates a new element
#             #
#             #     else:
#             #         oMath4 = OxmlElement(ij[5:-1])  # creates a new element
#             #         para_new._element.append(oMath4)
#             #         oMath4.text=ij[5:-6]
#             #     para_new._element.append(oMath4)
#             # if "</" not in ij:
#             #     if "=" not in ij:
#             #         print(ij)
#             # print()
#             # print(ij)
#     return paragraph.omaths
def copy_docx_run(run,output_para):
    """复制字块，为omath公式的段落准备"""
    output_run = output_para.add_run(run.text)
    output_run.style.name = run.style.name
    # print("style_name:",run.style.name)
    # output_run.bold = run.bold
    output_run.bold = False#所有的字体都不用粗体
    # print("run.bold:",run.bold)
    # Run's italic data
    output_run.italic = run.italic
    output_run.font.size = run.font.size
    # print("字体：", run.font.name, run.text)
    # output_run.font.name = run.font.name
    output_run.font.name = "Times New Roman"
    # 设置中文字体使用字体2->宋体,且要先执行上一句定义字体
    output_run._element.rPr.rFonts.set(qn('w:eastAsia'), u'宋体')
    # print("run.size:",run.size)
    # Run's underline data
    output_run.underline = run.underline
    output_run.font.subscript = run.font.subscript
    output_run.font.superscript = run.font.superscript

    # Run's color data
    output_run.font.color.rgb = run.font.color.rgb
    # Run's font data
    output_run.style.name = run.style.name
    # print("run.style.name:",run.style.name)
    # output_run.style.name = run.style.name

    # 开始复制公式

    try:
        # if field_number == 0 and run._element.xpath('w:fldChar')[0].attrib[qn("w:fldCharType")] == "begin":
        if run._element.xpath('w:fldChar')[0].attrib[qn("w:fldCharType")] == "begin":
            # print("开始复制公式：")
            fldChar1 = OxmlElement('w:fldChar')  # creates a new element
            fldChar1.set(qn('w:fldCharType'), 'begin')  # sets attribute on element
            output_run._element.append(fldChar1)
            # print("复制后：",output_run._element.xpath('w:fldChar')[0].attrib[qn("w:fldCharType")])
            # field_number = 1
    except:
        pass

    try:
        # 复制公式内容
        # print("run._element.xpath('w:instrText')[0].text:::",run._element.xpath('w:instrText')[0].text)
        # print("run._element.xpath('w:instrText')[0].text:::",type(run._element.xpath('w:instrText')[0].text),field_number)
        # if field_number == 1 and run._element.xpath('w:instrText')[0].text:
        if run._element.xpath('w:instrText')[0].text:
            # print(11111)
            # field_text=run._element.xpath('w:instrText')[0].text
            # fldChar3 = OxmlElement('w:instrText')
            instrText = OxmlElement('w:instrText')
            instrText.set(qn('xml:space'), 'preserve')  # sets attribute on element
            instrText.text = "%s" % (run._element.xpath('w:instrText')[0].text)
            output_run._element.append(instrText)

    except:
        pass
    try:
        # if field_number == 1 and run._element.xpath('w:fldChar')[0].attrib[qn("w:fldCharType")] == "end":
        if run._element.xpath('w:fldChar')[0].attrib[qn("w:fldCharType")] == "end":
            fldChar2 = OxmlElement('w:fldChar')  # creates a new element
            fldChar2.set(qn('w:fldCharType'), 'end')  # sets attribute on element
            output_run._element.append(fldChar2)
            # 结束公式复制
            # field_number = 0
    except:
        pass
    return output_para
def copy_docx_para(paragraph, output_para):
    """复制整段文字，含omath公式"""
    # print("para.text:",paragraph.text)
    # print("para.text:",paragraph._element.xml)
    # i=0
    # print("len(paragraph.runs):",len(paragraph.runs))
    # print("len(paragraph.runs):",len(paragraph.runs))
    if "oMath" in paragraph._element.xml:#加载段落的所有omath公式
        paragraph.omaths=[]
        paragraph.omaths=copy_docx_para_math(paragraph)
    run_num=omath_num=0
    for j in paragraph._element:
        # print("paragraph:",j,j.tag)
        if "}r" in j.tag:
            copy_docx_run(paragraph.runs[run_num], output_para)
            run_num+=1
            # print("run_num:",run_num)

        if "oMath" in j.tag:
            # print("omath_num:",omath_num,paragraph.omaths[omath_num])
            output_para._element.append(paragraph.omaths[omath_num])
            omath_num+=1
    # print(len(paragraph.runs))
    # print(len(paragraph._element))
    # for run in paragraph.runs:
    #     print(i,run,run.text)
    #
    #     copy_docx_run(run,output_para)
    #     i+=1

    return output_para
# def copy_docx_para(input_para,output_para):
#     """将段落的内容复制到新文档"""
#     for run in input_para.runs:
#         copy_docx_run(run, run.text, output_para)
#     output_para.paragraph_format.alignment = input_para.paragraph_format.alignment
def copy_docx_run2(run,run_text_new,output_para):
    output_run = output_para.add_run("%s" % run_text_new)
    output_run.style.name = run.style.name
    # print("style_name:",run.style.name)
    # output_run.bold = run.bold
    output_run.bold = False#所有的字体都不用粗体
    # print("run.bold:",run.bold)
    # Run's italic data
    output_run.italic = run.italic
    output_run.font.size = run.font.size
    # print("字体：", run.font.name, run.text)
    # output_run.font.name = run.font.name
    output_run.font.name = "Times New Roman"
    # 设置中文字体使用字体2->宋体,且要先执行上一句定义字体
    output_run._element.rPr.rFonts.set(qn('w:eastAsia'), u'宋体')
    # print("run.size:",run.size)
    # Run's underline data
    output_run.underline = run.underline
    output_run.font.subscript = run.font.subscript
    output_run.font.superscript = run.font.superscript

    # Run's color data
    output_run.font.color.rgb = run.font.color.rgb
    # Run's font data
    output_run.style.name = run.style.name
    # print("run.style.name:",run.style.name)
    # output_run.style.name = run.style.name

    # 开始复制公式

    try:
        # if field_number == 0 and run._element.xpath('w:fldChar')[0].attrib[qn("w:fldCharType")] == "begin":
        if run._element.xpath('w:fldChar')[0].attrib[qn("w:fldCharType")] == "begin":
            # print("开始复制公式：")
            fldChar1 = OxmlElement('w:fldChar')  # creates a new element
            fldChar1.set(qn('w:fldCharType'), 'begin')  # sets attribute on element
            output_run._element.append(fldChar1)
            # print("复制后：",output_run._element.xpath('w:fldChar')[0].attrib[qn("w:fldCharType")])
            # field_number = 1
    except:
        pass

    try:
        # 复制公式内容
        # print("run._element.xpath('w:instrText')[0].text:::",run._element.xpath('w:instrText')[0].text)
        # print("run._element.xpath('w:instrText')[0].text:::",type(run._element.xpath('w:instrText')[0].text),field_number)
        # if field_number == 1 and run._element.xpath('w:instrText')[0].text:
        if run._element.xpath('w:instrText')[0].text:
            # print(11111)
            # field_text=run._element.xpath('w:instrText')[0].text
            # fldChar3 = OxmlElement('w:instrText')
            instrText = OxmlElement('w:instrText')
            instrText.set(qn('xml:space'), 'preserve')  # sets attribute on element
            instrText.text = "%s" % (run._element.xpath('w:instrText')[0].text)
            output_run._element.append(instrText)

    except:
        pass
    try:
        # if field_number == 1 and run._element.xpath('w:fldChar')[0].attrib[qn("w:fldCharType")] == "end":
        if run._element.xpath('w:fldChar')[0].attrib[qn("w:fldCharType")] == "end":
            fldChar2 = OxmlElement('w:fldChar')  # creates a new element
            fldChar2.set(qn('w:fldCharType'), 'end')  # sets attribute on element
            output_run._element.append(fldChar2)
            # 结束公式复制
            # field_number = 0
    except:
        pass
    return output_para


def copy_picture(paragraph, output_doc):
    # 通过rid属性查询图片，并完成复制
    import xml.etree.cElementTree as ET
    namespace = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
                 'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
                 'v': "urn:schemas-microsoft-com:vml",
                 'wp': "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing",
                 'a': "http://schemas.openxmlformats.org/drawingml/2006/main",
                 'pic': "http://schemas.openxmlformats.org/drawingml/2006/picture"}

    # 定义了namespace，建议参看xml.etree.cElementTree中关于namespace的部分，添加了较多，可以自定义修改

    # 打开了对应的word文档，而不是上面讲到的.xml，但是本质是一样的。
    def get_img(root_element, target_tag, target_attribute, out_list):
        # 遍历得到目标属性target_attribute，得到属性value,写入out_list
        # 通过解读word的document.xml，得知与图片相关的rId属性值会出现在两个标签中
        # v:imagedata和a:blip
        # <v:imagedata r:id="rId8" o:title=""/>，<a:blip r:embed="rId8">
        for child in root_element:
            tag = child.tag
            attribute = child.attrib
            if tag in target_tag and target_attribute in child.attrib.keys():
                target_value = child.attrib[target_attribute]
                # print(target_value)
                out_list.append(target_value)
            else:
                get_img(child, target_tag, target_attribute, out_list)

    # xml_element = []

    # doc.paragraphs Proxy object wrapping <w:p> element.
    # xml_element.append(paragraph._element.xml)

    # for tbl in doc.tables:
    #     # print(tbl)
    #     # doc.tables Proxy class for a WordprocessingML <w:tbl> element.
    #     xml_element.append(tbl._element.xml)

    rId = []
    id = []

    # for element in xml_element:
    root = ET.fromstring(paragraph._element.xml)

    target_tag = ['{urn:schemas-microsoft-com:vml}imagedata',
                  '{http://schemas.openxmlformats.org/drawingml/2006/main}blip']
    # 即v:imagedata和a:blip
    target_attribute1 = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
    target_attribute2 = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    # 即:r:embed和r:id
    get_img(root, target_tag, target_attribute1, rId)
    get_img(root, target_tag, target_attribute2, id)

    dd = id + rId
    imgs = []

    for rid in dd:
        imgs.append(paragraph.part.related_parts[rid])
        # print(paragraph.part.related_parts[rid])

    # 1.4 保存图片到本地

    i = 1
    #
    for img in imgs:
        makedir(r"test")
        f = open(r"test/img%d.jpg" % i, 'wb')
        # help(f)
        # print("type(f):",type(f))

        f.write(img.blob)

        f.close()
        output_doc.add_picture(r"test/img%d.jpg" % i)
        # os.remove(r"test/img%d.jpg" % i)

        i += 1
        # print("i:", i)
def copy_docx(document,add_doc):
    """复制所在路径input_file的word文档内容至output_file"""

    def get_para_data(output_doc_name, paragraph):
        """将段落的内容复制到新文档"""
        output_para = output_doc_name.add_paragraph()
        for run in paragraph.runs:
            copy_docx_run2(run,run.text,output_para)
            # output_run = output_para.add_run("%s" % run.text)
            # output_run.style.name = run.style.name
            # # print("style_name:",run.style.name)
            # output_run.bold = run.bold
            # # print("run.bold:",run.bold)
            # # Run's italic data
            # output_run.italic = run.italic
            # output_run.font.size = run.font.size
            # print("字体：", run.font.name, run.text)
            # # output_run.font.name = run.font.name
            # output_run.font.name = "Times New Roman"
            # # 设置中文字体使用字体2->宋体,且要先执行上一句定义字体
            # output_run._element.rPr.rFonts.set(qn('w:eastAsia'), u'宋体')
            # # print("run.size:",run.size)
            # # Run's underline data
            # output_run.underline = run.underline
            # output_run.font.subscript = run.font.subscript
            # output_run.font.superscript = run.font.superscript
            #
            # # Run's color data
            # output_run.font.color.rgb = run.font.color.rgb
            # # Run's font data
            # output_run.style.name = run.style.name
            # # print("run.style.name:",run.style.name)
            # # output_run.style.name = run.style.name
            #
            # # 开始复制公式
            #
            # try:
            #     if field_number == 0 and run._element.xpath('w:fldChar')[0].attrib[qn("w:fldCharType")] == "begin":
            #         # print("开始复制公式：")
            #         fldChar1 = OxmlElement('w:fldChar')  # creates a new element
            #         fldChar1.set(qn('w:fldCharType'), 'begin')  # sets attribute on element
            #         output_run._element.append(fldChar1)
            #         # print("复制后：",output_run._element.xpath('w:fldChar')[0].attrib[qn("w:fldCharType")])
            #         field_number = 1
            # except:
            #     pass
            #
            # try:
            #     # 复制公式内容
            #     # print("run._element.xpath('w:instrText')[0].text:::",run._element.xpath('w:instrText')[0].text)
            #     # print("run._element.xpath('w:instrText')[0].text:::",type(run._element.xpath('w:instrText')[0].text),field_number)
            #     if field_number == 1 and run._element.xpath('w:instrText')[0].text:
            #         # print(11111)
            #         # field_text=run._element.xpath('w:instrText')[0].text
            #         # fldChar3 = OxmlElement('w:instrText')
            #         instrText = OxmlElement('w:instrText')
            #         instrText.set(qn('xml:space'), 'preserve')  # sets attribute on element
            #         instrText.text = "%s" % (run._element.xpath('w:instrText')[0].text)
            #         output_run._element.append(instrText)
            #
            # except:
            #     pass
            # try:
            #     if field_number == 1 and run._element.xpath('w:fldChar')[0].attrib[qn("w:fldCharType")] == "end":
            #         fldChar2 = OxmlElement('w:fldChar')  # creates a new element
            #         fldChar2.set(qn('w:fldCharType'), 'end')  # sets attribute on element
            #         output_run._element.append(fldChar2)
            #         # 结束公式复制
            #         field_number = 0
            # except:
            #     pass

        # Paragraph's alignment data
        output_para.paragraph_format.alignment = paragraph.paragraph_format.alignment


    # input_doc = Document(r'C:\Users\Administrator\Desktop\test2.docx')
    # input_doc = Document(r'C:\Users\Administrator\Desktop\test5.docx')
    # input_doc = Document(r"word题目汇总/%s.docx"%title)
    input_doc = add_doc
    #导出word文档中的图片链接地址
    # docx_document= Document(r'C:\Users\赵安宁\Desktop\test7.docx')
    # export_files = []
    # # 遍历所有附件
    # index = 0
    # docx_related_parts = input_doc.part.related_parts
    # for part in docx_related_parts:
    #     part = docx_related_parts[part]
    #     part_name = str(part.partname)
    #     # 附件路径（partname）
    #     if part_name.startswith('/word/media/') or \
    #             part_name.startswith('/word/embeddings/'):
    #         # 只导出这两个目录下的# 构建导出路径
    #         index += 1
    #         save_dir = os.path.dirname(os.path.abspath(__file__))
    #         # 获取当前py脚本路径
    #         index_str = str(index).rjust(2, '0')
    #         save_path = save_dir + '\\' + index_str + ' - ' + os.path.basename(part.partname)
    #         # 拼接路径
    #         print('导出路径：', save_path)
    #         # 写入文件
    #         with open(save_path, 'wb') as f:f.write(part.blob)
    #         # 记录文件export_files.append(save_path)
    # print('导出的所有文件：', export_files)
    # output_doc = Document(r'C:\Users\Administrator\Desktop\A20202122Z4.docx')
    # output_doc = Document(r'C:\Users\Administrator\Desktop\test2.docx')
    output_doc = document
    # output_doc.add_paragraph()
    i = 1
    title_type="calculation"
    # 复制图片到指定文件夹
    # image_number_all = get_pictures(input_doc, r"test/image")
    # image_number_new = 1
    for para in input_doc.paragraphs:
        # print("第%d段" % i)
        i += 1
        if len(para.text) > 1 :  # 有文本时才复制
            get_para_data(output_doc, para)
        copy_picture(para, output_doc)

    # output_doc.save(output_file)
    # print("生成文档，请在以下位置查看",output_file )
def delete_paragraph(paragraph):
    """删除段落word文档中的paragraph"""
    p = paragraph._element
    p.getparent().remove(p)
    # p._p = p._element = None
    paragraph._p = paragraph._element = None
    # print("delete成功")
    return paragraph
def delete_blank_row(para):
    """删除空行"""
    if len(para.text) <= 1 and len(para.runs) != 1:  # 删除空行,等于1为图片
        p = para._element
        p.getparent().remove(p)
        p._p = p._element = None
    return para
def docx_to_picture(title):
    """将docx文档转换为image"""
    # input_doc=r"F:\个性化题库20180423\题不二错1.8\题不二错1.8测试版\test\2023125Z1K1.docx"
    # input_doc_path=r"F:\个性化题库20180423\题不二错1.8\题不二错1.8测试版\test"
    makedir(r"cach")
    makedir(r"cach/word")
    makedir(r"cach/pdf")
    makedir(r"cach/pic")
    makedir(r"cach/word_to_images")
    # makedir(r"cach/word_to_images")
    import os

    main_path = os.path.dirname(os.path.realpath(sys.argv[0]))
    # doc_for_pic.save(r"%s\cach\word\%s.docx"%(main_path,title))
    # input_doc_path=r"%s\cach\word"%(main_path
    input_doc_path=r"%s\cach\word\%s.docx"%(main_path,title)
    output_pdf_path=r"%s\cach\pdf"%(main_path)
    # image_path=r"%s\cach\word_to_images"%(main_path)
    image_path=r"%s\cach\pic"%(main_path)

    def clear_folder(path):
        """清空文件夹下所有文件"""
        path=main_path+path
        print("path:",path)

        # print("path:",os.listdir(path))
        import os
        for file in os.listdir(path):
            file_path = os.path.join(path, file)
            if os.path.isfile(file_path):
                os.remove(file_path)
            elif os.path.isdir(file_path):
                clear_folder(file_path)

    # clear_folder(input_doc_path)

    # shutil.copy(r"%s\word题目汇总\%s.docx"%(main_path,title),
    #             r"%s\%s.docx"%(input_doc_path,title))
    # 定义PDF文件路径和输出区间路径
    # 待转换pdf文件路径
    # path = input_doc_path
    # doc/docx转换后的路径
    # path_convert = output_path
    # 转换后的图片路径
    # path_images = r"cach/images"
    # print("-----doc开始转换为docx-----")

    # for i in os.listdir(path):
    #     file_name, file_suffix = i.split(".")
    #     if file_suffix == "doc":
    #         word = Dispatch('Word.Application')
    #         doc = word.Documents.Open(path + f"{i}")
    #         doc.SaveAs(path + f"{file_name}.docx", FileFormat=12)
    #         print(i, "转换完成")
    #         doc.Close()
    #         word.Quit()
    #         os.remove(path+f"{i}")
    #         sleep(3)


    print("-----开始转换为pdf-----",title)
    word = Dispatch('Word.Application')
    doc = word.Documents.Open(r"%s\cach\word\%s.docx"%(main_path,title))
    doc.SaveAs(output_pdf_path + "\\" + f"{title}.pdf", FileFormat=17)
    print(title, "...转换完成")
    doc.Close()
    word.Quit()
    # sleep(3)
    # for i in os.listdir(input_doc_path):
    #     file_name, file_suffix = i.split(".")
    #     if file_suffix == "docx":
    #         word = Dispatch('Word.Application')
    #         doc = word.Documents.Open(input_doc_path +"\\"+ f"{i}")
    #         doc.SaveAs(output_pdf_path +"\\"+ f"{file_name}.pdf", FileFormat=17)
    #         print(i, "...转换完成")
    #         doc.Close()
    #         word.Quit()
    #         sleep(3)
    #     else:
    #         pass

    # 循环遍历PDF文件，并转换为图片
    # for filename in os.listdir(path):
    #     if filename.endswith(".pdf"):
    #         # 获取当前 PDF 文件的总页数
    #         doc = fitz.open(path + filename)
    #         total_pages = doc.page_count
    #         doc.close()
    #
    #         print(f"正在转换 {filename}，共 {total_pages} 页...")
    #         # 按照原图输出pdf文件为word
    #         # 若想黑白输出pdf为图片格式，修改grayscale=False为grayscale=True
    #         for i, page in enumerate(convert_from_path(path + filename, grayscale=False), start=1):
    #             # 构造文件名
    #             output_filename = os.path.splitext(filename)[0] + "_" + str(i) + ".png"
    #             # 保存图片
    #             page.save(path_images + output_filename, "png")
    #             # 输出转换进度
    #             print(f"已完成第 {i}/{total_pages} 页的转换")
    def covert2pic(file_path, zoom, png_path):
        doc = fitz.open(file_path)
        total = doc.page_count
        for pg in range(total):
            page = doc[pg]
            zoom = int(zoom)  # 值越大，分辨率越高，文件越清晰
            rotate = int(0)

            trans = fitz.Matrix(zoom / 100.0, zoom / 100.0).prerotate(rotate)
            pm = page.get_pixmap(matrix=trans, alpha=False)
            if not os.path.exists(png_path):
                os.mkdir(png_path)
            # save = os.path.join(png_path, '%s.png' % (title))
            save_path = os.path.join(png_path, '%s.jpg' % (title))
            pm.save(save_path)

            #裁剪图片大小
            # print("图片路径：",save_path)
            image = Image.open(save_path)
            # print("原图大小：{}".format(image.size))
            # box = (25, 150, 625, 750)  # 裁剪大小
            # NewIm = Im.crop(box)
            # print("新图大小：{}".format(NewIm.size))
            # NewIm.save('crop1.jpg')  # 保存图片

            #图片大小范围，左侧100起，至750，上侧90起，至600
            cropped_image = image.crop((240,260,1600,1200))
            # cropped_image.show()

            try:
                #查看是否存在，如果存在，则提示是否替换
                Image.open(save_path)
                yesno = tk.messagebox.askyesno(title="覆盖提示", message=
                "您即将保存的题目：%s,图片已经存在，是否替换%s"%title)
                if yesno==True:
                    cropped_image.save(save_path)

            except:
                cropped_image.save(save_path)

            print("succes")

        doc.close()

    # pdfPath = r'F:\个性化题库20180423\题不二错1.8\题不二错1.8测试版\test\test4\2023126Z7K9.pdf'
    # imagePath = r'test4'
    covert2pic(r"%s\%s.pdf"%(output_pdf_path,title), 300, image_path)

    print("-----已完成图片转换%s-----"%title)

@handle_exceptions(logger)
def docx_add(document, title,title_option_changes):
    """以word形式添加题目文档"""
    from docx import Document
    # from docxcompose.composer import Composer
    # target_doc = Document(r"test/2020271K1.docx")
    # print("file1,type_number:", document, title)
    # target_doc = Document("%s" % file1)
    # target_composer = Composer(document)
    # target_composer.append(page_break_doc)
    def change_option(title,title_option_changes):
        """更改选择题的选项顺序"""
        print("开始更改选择题顺序",title)
        target_doc_old=Document(r"word题目汇总\%s.docx" % title)
        doc_for_pic=Document()#为ppt准备的word文档，此处没有直接生成ppt，还未详细学习ppt
        doc_for_pic_select=False#选择题判断
        doc_for_pic.add_paragraph(target_doc_old.paragraphs[0].text)#添加题号
        delete_paragraph(target_doc_old.paragraphs[0])  # 删除题号
        paragraphs=target_doc_old.paragraphs
        title_option_changes[str(title)] = None
        select_number = 0
        for i in range(len(paragraphs)):
            # print("段落：",i)
            para=paragraphs[i]
            if len(para.text) <= 1 and len(para.runs) < 1:  # 删除空行
                p = para._element
                p.getparent().remove(p)
                p._p = p._element = None
                continue
            #判断该题是否为选择题？则必有ABCD四个选项。或者ABC三个选项。
            # print("paragraphs[i][0:2]:",paragraphs[i].text,type(paragraphs[i].text))
            #bug修复，如果四个选项在一行，则会出现错误，无下一行，假设下一行为空值
            try:
                paragraphs[i + 1].text[0:2]
            except:
                target_doc_old.add_paragraph()
                paragraphs=target_doc_old.paragraphs


            if (paragraphs[i].text[0:2]=="A．" or paragraphs[i].text[0:2]=="A.")\
                    and (paragraphs[i+1].text[0:2]=="B．" or paragraphs[i+1].text[0:2]=="B."):
                if paragraphs[i+2].text[0:2]=="C．" or paragraphs[i+2].text[0:2]=="C.":
                    options=[0,1,2,]
                    if paragraphs[i+3].text[0:2]=="D．" or paragraphs[i+3].text[0:2]=="D.":
                        options=[0,1,2,3,]
                        if paragraphs[i+3].text[0:2]=="E．" or paragraphs[i+3].text[0:2]=="E.":
                            options=[0,1,2,3,4]
                opti_number=len(options)
                # print("该题是选择题，",title)
                options2=options
                options4=options[0:opti_number]
                result=1
                while result==1:
                    result=0
                    random.shuffle(options2)
                    for option_i in range(opti_number):
                        if options4[option_i]==options[option_i]:
                            result=1
                            continue
                # print("options,options4:",options,options4)
                # j=i
                title_option_changes[str(title)]=options
                options3=["A","B","C","D","E"]
                #选择题选项，一次完成四个选项的复制
                for j in range(len(options)):
                    # print(i,j,options[j],len)
                    # print("paragraphs[i+options[j]]:",paragraphs[i+options[j]].text)
                    #如果一段的文本内容少于一行的一半，则不换行。一行为46个英文字符（2栏）
                    len_A=len(paragraphs[i+options[0]].runs)
                    len_B=len(paragraphs[i+options[1]].runs)
                    len_C=len(paragraphs[i+options[2]].runs)
                    len_D=len(paragraphs[i+options[3]].runs)
                    # print("len:",title,len_A,len_B,len_C,len_D)
                    if j==1 and len_A+len_B<23:
                        run = paragraphs[i + options[j]].runs[0]
                        copy_docx_run2(run, "  ", paragraph_new)
                    elif j==2 and len_A+len_B+len_C<20:
                        run = paragraphs[i + options[j]].runs[0]
                        copy_docx_run2(run, "  ", paragraph_new)
                    elif j==3 and (len_A+len_B+len_C+len_D<18 or len_C+len_D<20):
                        # print("paragraphs:")
                        # print(paragraphs)
                        # print(paragraphs[i + options[j]])
                        # print(paragraphs[i + options[j]].runs)
                        # print(paragraphs[i + options[j]].runs[k])
                        run = paragraphs[i + options[j]].runs[0]
                        copy_docx_run2(run, "  ", paragraph_new)
                    else:
                        paragraph_new = document.add_paragraph()

                    paragraph_new2=doc_for_pic.add_paragraph()

                    #复制段落的所有内容包含文字和公式，并对改变的选项进行变更
                    #在无法加载omath公式又找不到原因的情况下，调用功能重新编程，结果成功
                    paragraph=paragraphs[i+options[j]]
                    if "oMath" in paragraph._element.xml:  # 加载段落的所有omath公式
                        paragraph.omaths = []
                        paragraph.omaths = copy_docx_para_math(paragraph)
                    run_num = omath_num = 0
                    for k in paragraph._element:
                        # print("paragraph:",j,j.tag)
                        if "}r" in k.tag:
                            if run_num == 0:
                                # print("options3:",options3)
                                text = options3[j]
                            else:
                                text = paragraph.runs[run_num].text
                            copy_docx_run2(paragraph.runs[run_num],text, paragraph_new)
                            # copy_docx_run(paragraph.runs[run_num], paragraph_new)
                            run_num += 1
                            # print("run_num:", run_num)

                        if "oMath" in k.tag:
                            # print("omath_num:",omath_num,paragraph.omaths[omath_num])
                            paragraph_new._element.append(paragraph.omaths[omath_num])
                            omath_num += 1
                    # print(len(paragraph.runs))
                    # print(len(paragraph._element))
                    # for run in paragraph.runs:
                    #     print(i,run,run.text)
                    #
                    #     copy_docx_run(run,output_para)
                    #     i+=1

                    # return output_para

                    # for k in range(len(paragraphs[i+options[j]].runs)):
                    #     run=paragraphs[i+options[j]].runs[k]
                    #     if k==0:
                    #         text=options3[j]
                    #     else:
                    #         text=run.text
                    #     copy_docx_run(run,text,paragraph_new)
                    #     # print("text:",text)
                    #     copy_docx_run(run,text,paragraph_new2)
                doc_for_pic_select=True
                print("22222222222222")
                select_number=len(options)

            elif select_number!=0: # 其它选项的段落，直接跳过
                select_number-=1
            elif select_number==0 and len(paragraphs[i].runs) > 1:#复制其它段落(不含图片段落)
                paragraph_new = document.add_paragraph()
                paragraph_new2 = doc_for_pic.add_paragraph()
                copy_docx_para(paragraphs[i],paragraph_new)
                copy_docx_para(paragraphs[i],paragraph_new2)
                # for run in paragraphs[i].runs:
                #     copy_docx_run2(run, run.text, paragraph_new)
                #     copy_docx_run2(run, run.text, paragraph_new2)
            # print("3333333")
            copy_picture(para,document)
            copy_picture(para,doc_for_pic)
        if doc_for_pic_select==True:
            #选择题需要生成新图片供ppt使用
            for para in doc_for_pic.paragraphs:
                para.paragraph_format.space_after = Pt(0)  # 段后0磅

            doc_for_pic.save(r"cach/word/%s.docx"%title)
            # print("选择题，保存文档")
        # if title_option_changes[str(title)]==None:
        #     copy_docx(document, target_doc_old)  # 复制非选择题
        # for i in target_doc_old.paragraphs:
            # print("i.textttt:",i.text)
        # print("更改成功")
        # print("11111111111111")

        return target_doc_old,title_option_changes
    target_doc2,title_option_changes=change_option(title,title_option_changes)
    # print("title_option_changes:",title_option_changes)
    # delete_paragraph(target_doc.paragraphs[0])
    # print("delete1：：")
    # print("i.inlieshape:", len(target_doc2.inline_shapes),target_doc2.inline_shapes)

    # target_doc4 =target_doc3 = target_doc2
    # number=0
    # if len(para.text) <= 1 and len(para.runs) < 1:  # 删除空行
    #     p = para._element
    #     p.getparent().remove(p)
    #     p._p = p._element = None
    for para in document.paragraphs:
        para.paragraph_format.space_after = Pt(0)  # 段后0磅
        # print("title,para.text:",title,para.text)
        # print(len(para.runs))
        delete_blank_row(para)

    return document,title_option_changes


def wordtitle(document,title,title_option_changes):
    """# 在word文档中输入题目和图片,如果图片不存在，则用一张预备图片(题目不存在)替代"""
    title_option_changes[title]=None
    # document ,title_option_changes= docx_add(document, title,title_option_changes)
    # title_option_changes[title]=None
    try:#插入文字题目
        # document.save("test1.docx")
        # print("错误1")
        # document = docx_add("test1.docx", "%s.docx" % title)
        # print(":::::", help(document.sections))
        document,title_option_changes = docx_add(document, title,title_option_changes)
        title_docx = Document(r"word题目汇总\%s.docx" % title)
        # 计算题儿留五行空白。
        doc_text = ""
        title_type = "calculation"
        for para in title_docx.paragraphs:
            doc_text +=para.text
        if "A" in doc_text and "B" in doc_text and "C" in doc_text and "D" in doc_text:
            title_type = ""
        elif "实验" in doc_text:
            title_type = ""
        if title_type == "calculation":
            for i in range(5):
                document.add_paragraph()


        # print("正确2")
        # document.save("%s.docx" % title)

    except Exception as f:#插入图片题目
        # tk.messagebox.showinfo(title="提示",message="%s%s"%(f.__traceback__.tb_lineno,f))
        print("插入题目word文档失败",title,"原因:",f.__traceback__.tb_lineno,f,)
        # document.add_picture(r'%s\图片题目汇总\%s.jpg' % (main_path, title))  # 题目内容
        try:
            # title=title_replace(title)
            print(":::",r'%s\图片题目汇总\%s.jpg' % (main_path, title))
            image = Image.open(r'%s\图片题目汇总\%s.jpg' % (main_path, title))
            width, height = image.size
            print(f"图片宽度: {width / 28}, 图片高度: {height}")
            if width/28>12:#分为一栏
                if width/28>24:
                    document.add_picture(r'%s\图片题目汇总\%s.jpg' % (main_path, title), width=Inches(7.0))#等比例压缩缩放
                else:
                    document.add_picture(r'%s\图片题目汇总\%s.jpg' % (main_path, title))  # 题目内容
                # section = document.add_section(WD_SECTION_START.NEW_COLUMN)
                # section = document.add_section(WD_SECTION.CONTINUOUS)
                section = document.sections[len(document.sections)-1]
                sectPr = section._sectPr
                cols = sectPr.xpath('./w:cols')[0]
                cols.set(qn('w:num'), '1')
                # section_types=sectPr.xpath('./w:type')[0]
                # section_types.set(qn('w:val'), 'continuous')
                # print("len(document.sections):",len(document.sections))
                # print("section.start_type:",section.start_type)


        except:
            print("插入题目图片失败", title)
            document.add_picture(r'%s\图片题目汇总\00.jpg'%main_path)  # 题目内容
            title_not_exit.append(title)


    return title_option_changes

@handle_exceptions(logger)
def generate_answer(text):
    import re
    options = re.findall(r'([A-Z|、|.]+)(正确|错误|选项正确|选项错误|符合题意)', text)
    # answer = ''
    answer = []  # 所有选项
    answer_right = ""  # 所有正确答案
    for option, correctness in options:
        if correctness == '正确' or correctness == '错误' or \
                correctness == '选项正确' or correctness == '选项错误':
            # answer += option
            answer.append(option)
        if correctness == '正确' or   correctness == '选项正确' or   correctness == '符合题意':
            answer_right += option
            # answer_right.append(option)
            # answer_right =answer_right.join(option)
    # return ''.join(sorted(answer))

    options = re.findall(r'(故选+)([A-Z|、|.])', text)
    # answer = ''
    for correctness, option in options:
        if correctness == '故选':
            # answer += option
            answer.append(option)
            answer_right += option
            # answer_right.append(option)
            # answer_right =answer_right.join(option)
    # return ''.join(sorted(answer))
    answer_right = answer_right.replace("、", "")  # 去除正确答案中的顿号
    answer_right = answer_right.replace(".", "")
    return answer, answer_right
def wordtitle_answer(document,title,title_option_changes,pra_number):
    """# 在word文档中输入题目和图片,如果图片不存在，则用一张预备图片(题目不存在)替代"""


    try:#插入文字题目答案

        from docx import Document
        # from docxcompose.composer import Composer
        # target_doc = Document(r"test/2020271K1.docx")
        # print("file1,type_number:", document, title)
        # target_doc = Document("%s" % file1)
        # target_composer = Composer(document)
        # target_composer.append(page_break_doc)
        answer_docx_title=Document(r"word题目汇总\%s.docx"%title)
        # delete_paragraph(target_doc.paragraphs[0])
        doc_answer_for_pic = Document()  # 为ppt准备的中间word文档
        doc_answer_for_pic.add_paragraph(answer_docx_title.paragraphs[0].text)
        paragraph_new2 = doc_answer_for_pic.add_paragraph()
        paragraph_new3 = doc_answer_for_pic.add_paragraph()
        delete_paragraph(answer_docx_title.paragraphs[0])  # 删除题号
        # 未进行选项变更的题目直接复制
        # print("title_option_changes::",title_option_changes,title_option_changes[title[1:]])
        if title_option_changes[title[1:]]==None:
            for paragraph in answer_docx_title.paragraphs:
                copy_docx_para(paragraph,document.paragraphs[-1])
            # copy_docx(document,answer_docx_title)
            return
        text=""
        for paragraph in answer_docx_title.paragraphs:
            text += paragraph.text
        answer, answer_right = generate_answer(text)

        for paragraph in answer_docx_title.paragraphs:
            # print("answer and answer_ringt is: ", title,answer, answer_right)
            # print("选择题更改顺序，",title_option_changes)
            # middle_doc=Document()
            # middle_para=middle_doc.add_paragraph()#中间量备用
            # copy_docx_para(paragraph, document.paragraphs[-1])
            if len(paragraph.text)<7 and "答案" in paragraph.text:
                pass
            else:
                copy_docx_para(paragraph, paragraph_new2)
            # for run in paragraph.runs:
            #     #将段落的字块打散,每个字符构建一个单独的块
            #     if run.text != "":
            #         for text in run.text:
            #             # copy_docx_run(run,text,middle_para)
            #             copy_docx_run(run,text,paragraph_new2)
            #     else:
            #         # copy_docx_run(run, run.text, middle_para)
            #         copy_docx_run(run, run.text, paragraph_new2)
            # if paragraph.text == "":  # 如果是空段落，则自动进入下一段，图片答案不寻找答案
            #     continue
        #删除答案中的"答案"二字


        #text="BCD　[物体在两极，万有引力等于重力，有GMmR2＝mg0，故A正确；物体在地球赤道上，"
        #一开始的几个字母就是答案，直接替换
        """查找更正后的正确答案"""
        options = ["A", "B", "C", "D", "E"]
        answer_list = []
        for answer_option in answer_right:
            for number in range(len(options)):
                if options[number] == answer_option:
                    # print("find answer:", answer_option)
                    answer_list.append(number)
        answer_list.sort()
        # print("answer_list:", answer_list,)


        #新旧选项替换
        old_option={"A":0,"B":1,"C":2,"D":3,"E":4}
        new_option={}
        print("title_option_changes:",title_option_changes)

        for number in range(len(title_option_changes[title[1:]])):
            new_option[title_option_changes[title[1:]][number]]=options[number]
        print("new_option:",new_option)
        number = 0
        answer_text = ""

        for answer_number in title_option_changes[title[1:]]:
            if int(answer_number) in answer_list:
                answer_text += new_option[answer_number]
                number += 1
        # print("answer_text:", answer_text)



        for number in range(len(answer_text)):
            # middle_para.runs[number].text=answer_text[number]
            # print("number,answer_text,paragraph.runs[number]:",number,answer_text,paragraph.runs[number],paragraph_new2.runs[number])
            # print("paragraph.runs[number].text:",paragraph.runs[number].text,number)
            # paragraph.runs[number].text=answer_text[number]
            try:
                paragraph_new2.runs[number].text=answer_text[number]
            except:
                paragraph_new2.add_run( answer_text[number])
        # paragraph_new2.text[2:4]="ab"
        runs=paragraph_new2.runs
        # print("paragraph_new2.text:",paragraph_new2.text)

        #更改答案正文里面的内容，按照特定格式更改，故选A， A正确， A、B、C正确
        answer_text_list=["A","B","C","D","、",".","正","错","确","误","选","项",
                          "符","合","题","意"," ","　"]
        #答案只有选项，则长度小于5
        if len(paragraph_new2.text)<5:
            yesno=True
            answer_list2=[]
            answer_text=""
            for run in runs:
                print("title,run.text:",title,run.text)
                if run.text not in answer_text_list:
                    yesno=False
                    break
            if yesno==True:
                for number in range(len(runs)):
                    # print("runs::",title, runs, new_option[old_option[runs[number].text]])
                    try:
                        answer_list2.append(new_option[old_option[runs[number].text]])

                        # runs[number].text = new_option[old_option[runs[number].text]]
                    except Exception as e:
                        print("选择题答案选项更改失败1：", e.__traceback__.tb_lineno, e)
                answer_list2.sort()
                print("answer_list2:",answer_list2)
                for number in range(len(answer_list2)):
                    runs[number].text=answer_list2[number]
                    answer_text += answer_list2[number]
        else:
            for number in range(len(paragraph_new2.runs)):
                # print("paragraph_new23.text:", paragraph_new2.text)
                try:
                    # print("runs[number].text:",number,runs[number].text,runs[number+1].text)
                    if runs[number].text in options and runs[number-2].text=="故" and runs[number-1].text=="选":
                           runs[number].text = new_option[old_option[runs[number].text]]
                    elif runs[number].text in options and runs[number+1].text in ["正","错",""
                            ] and runs[number+2].text in ["确","误"]:
                        # print("A正确：",runs[number].text,runs[number+1].text)
                        runs[number].text=new_option[old_option[runs[number].text]]
                        # print("替换为B正确",runs[number].text,runs[number+1].text)
                    elif runs[number].text in options and runs[number+1
                    ].text in answer_text_list and runs[number+2
                    ].text in answer_text_list and runs[number+3
                    ].text in answer_text_list and number>2:#去除最开始的几个选项答案被重复修改
                           runs[number].text = new_option[old_option[runs[number].text]]


                except Exception as e:
                    print("选择题答案选项更改失败：",e.__traceback__.tb_lineno,e)
        # middle_doc.save("test.docx")
        # print("保存成功")
        # print("title:::::",title,title_option_changes)
        # for run in paragraph_new2.runs:
            # print("run::",run.text)
        copy_docx_para(paragraph_new2,document.paragraphs[-1])

        #将题目的答案保存到中间excel表格中供pptx使用
        if name=="赵安宁":
            try:
                #插入行
                title_answer_text=pd.DataFrame(pd.read_excel(r"cach/excel/select_answer.xlsx"))
                # print("title_answer_text2:",title_answer_text)
                # print("pranumber,answer_text:",pra_number,answer_text)
                new_row = pd.DataFrame([[pra_number,answer_text]],
                                       columns=["题号","答案"])
                title_answer_text = pd.concat([title_answer_text, new_row], ignore_index=True)
                title_answer_text=title_answer_text.set_index("题号")
                # print("title_answer_text3:",title_answer_text)
            except Exception as e:
                print("插入行失败",e.__traceback__.tb_lineno,e)
                makedir(r"cach/excel")
                title_answer_text = pd.DataFrame([[pra_number,answer_text]],
                                       columns=["题号","答案"])
                title_answer_text=title_answer_text.set_index("题号")
            title_answer_text.to_excel(r"cach/excel/select_answer.xlsx")
            # print("title_answer_text:",title_answer_text)

        for i in document.paragraphs:
            i.paragraph_format.space_after = Pt(0)  # 段后0磅
        for i in doc_answer_for_pic.paragraphs:
            i.paragraph_format.space_after = Pt(0)  # 段后0磅
        makedir(r"cach")
        makedir(r"cach/word")
        doc_answer_for_pic.save(r"cach/word/%s.docx"%title)




    except Exception as f:#插入图片题目
        print("插入文字答案失败，失败原因：",f.__traceback__.tb_lineno,f,)#问题行数
        try:
            print("插入文字题目失败：",title)
            # title=title_replace(title)
            # print("错误33")
            # print("title：",title)
            document.add_picture(r'%s\图片题目汇总\%s.jpg' % (main_path, title))  # 题目内容
            # print("错误4")
        except:

            try:
                print("插入图片题目答案失败：", title)

                answer = wordtitle_answer_find_excel(title)  # 考试情况下插入选择题答案
                # document.add_heading(answer, 3)
                paragraph.add_run(answer)
            except:
                print("插入选择题答案失败，添加不成功图片", title)
                document.add_picture(r'%s\图片题目汇总\00.jpg'%main_path)  # 题目内容
                title_not_exit.append(title)
            # answer = wordtitle_answer_find_excel(title)  # 考试情况下插入选择题答案


def student_wrong4(path,exce,team,main_number):
    """#识别学生的错题，如果是错一练三，查找训练题，并一起编入title_wrong中"""
    print("函数student_wrong4运行中")
    # print(path)
    content=excel(path+'\\'+exce)
    # print(content)
    titlewrong={}
    titlewrong2 = {}  # 用于存放考试题目，图片长度不同
    name_finish={}
    rowmax=len(content)
    colmax=len(content[0])
    titlewrong['team'] = team
    titlewrong2['team'] = team
    # main_number = "2019"

    progressbarOne['maximum'] = colmax
    # 进度值初始值
    progressbarOne['value'] = 0
    win.update()


    # titlewrong2['team'] = team
    for j in range(colmax):#从列开始，登记姓名和错题
        progressbarOne['value'] += 1
        win.update()

        name=content[0][j]
        # tk.Label(win,text="%s:当前学生数%s,总学生数：%s"%(team,j,colmax)).pack
        label_2.set("%s:获取学生--%s错题信息，当前学生数%s,总学生数：%s"%(team,name,j+1,colmax))
        win.update()
        print("当前学生数和总学生数：",j,colmax)
        # text.insert(tk.INSERT, "第%s个学生信息获取中，本班共%s个学生\n" % (j + 1, colmax))
        # print("第%个学生错题信息处理中，姓名："%(j),name)
        # print("第%个（共%个）学生错题信息处理中，姓名："%(j,colmax),name)

        # print("name",name)
        # print("finish_name_all_dict[name]:",finish_name_all_dict)
        # finish=finish_name_all_dict[name][-1]
        # name_is_exist=0
        
        # try:#检测学生登记姓名是否存在，如果不存在，则不打印
        #     name_is_exist=(finish_name_all_dict[name][-1]==1)
        # except:
        #     print("——-----------错题登记姓名：%s不存在---------------"%name)
        # if name!=None and name_is_exist==1:
        if name!=None :
            # if len(name)==3:
            #     name+='  '
            # elif len(name)==2:
            #     name+='    '
            titlewrong[name] = []
            titlewrong2[name]=[]#用于存放考试题目，图片长度不同

            for i in range(rowmax)[1:]:
                title = str(content[i][j])
                # print("name:",name)
                if title != None and title!="None" :
                    # print("title[0]",title[0])
                    # print("title[0:4]",title[0:4])
                    if title[0] == "3" and title[0:4]!="3019" or title[0]=="5":
                        # 判断是否为错一练三,为知识点专练
                        # print("titlewrong:",titlewrong)
                        # print("name:",name)
                        # print("titlewrong[name]:",titlewrong[name])
                        title3 = que_type_title("对象",title)
                        title4=titlepractice(title3)
                        # print("title,title3:",title,title3)
                        # print("错一练三，title,title4:",title,title4)
                        titlewrong[name].append(title4)
                        # titlewrong=class_title_pra2(titlewrong,name,title)
                    elif title[0] == "2" or title[0] == "4" or title[0:4]=="3019" :
                        if title[0:4] == "2019" or title[0:4] == "2016" or \
                                title[0:4]=="3019" or title[0:4]=="2020" or title[0:4]=="2022"\
                                or title[0:4]=="2023":# 判断是否为二次登记

                            title_registor = {}  # 将错题登记到个人信息中{'金世语': [{'错题': [2019030802], '训练题': ['20190308187']}],
                            title_registor["错题"] = [title]
                            titlewrong[name].append(title_registor)
                            # print("二次登记，title_registor",title_registor)
                        else:

                            title = main_number + "0" + title
                            title_registor2 = {}  # 将错题登记到个人信息中{'金世语': [{'错题': [2019030802], '训练题': ['20190308187']}],
                            title_registor2["错题"] = [title]
                            titlewrong2[name].append(title_registor2)
                            # print("试卷，title_registor2",title_registor2)
                    else:
                        # print("number:",title)
                        # print("number[0:4]:",title[0:4])
                        # title = main_number + "0" + title
                        title = main_number  + title
                        title_registor={}#将错题登记到个人信息中{'金世语': [{'错题': [2019030802], '训练题': ['20190308187']}],
                        title_registor["错题"]=[title]
                        titlewrong[name].append(title_registor)
                        # print("三维设计，title_registor1",title_registor)
                    # finish=3
            for key in list(titlewrong.keys()):
                if not titlewrong.get(key):
                    del titlewrong[key]
            for key in list(titlewrong2.keys()):
                if not titlewrong2.get(key):
                    del titlewrong2[key]
        # name_finish[name] = finish

    # titlewrong['unit']=unit
    # print('titlewrong',titlewrong)
    # print('titlewrong2',titlewrong2)
    return [titlewrong,titlewrong2]
    # return [sorted(titlewrong.keys()),sorted(titlewrong2.keys())]
# def student_wrong5(path,exce,team,main_number):
#     """#识别学生的错题，如果是错一练三，查找训练题，并一起编入title_wrong中,去除筛选功能"""
#     print("函数student_wrong5运行中")
#     # print(path)
#     content=excel(path+'\\'+exce)
#     # print("exceexce:",exce)
#     # print(content)
#     titlewrong={}
#     titlewrong2 = {}  # 用于存放考试题目，图片长度不同
#     name_finish={}
#     rowmax=len(content)
#     colmax=len(content[0])
#     titlewrong['team'] = team
#     titlewrong2['team'] = team
#     # main_number = "2019"
#
#     # titlewrong2['team'] = team
#     for j in range(colmax):#从列开始，登记姓名和错题
#         print("当前学生数和总学生数：",j,colmax)
#         # print("第%个学生错题信息处理中，姓名："%(j),name)
#         # print("第%个（共%个）学生错题信息处理中，姓名："%(j,colmax),name)
#         name=content[0][j]
#         # print("name",name)
#         # print("finish_name_all_dict[name]:",finish_name_all_dict)
#         # finish=finish_name_all_dict[name][-1]
#         # name_is_exist=0
#         # try:#检测学生登记姓名是否存在，如果不存在，则不打印
#         #     name_is_exist=(finish_name_all_dict[name][-1]==1)
#         # except:
#         #     print("——-----------错题登记姓名：%s不存在---------------"%name)
#         if name!=None :
#             # if len(name)==3:
#             #     name+='  '
#             # elif len(name)==2:
#             #     name+='    '
#             titlewrong[name] = []
#             titlewrong2[name]=[]#用于存放考试题目，图片长度不同
#
#             for i in range(rowmax)[1:]:
#                 title = str(content[i][j])
#                 # print("name:",name)
#                 if title != None and title!="None" :
#                     print("title[0]：",title[0])
#                     print("title[0:4]：",title[0:4])
#                     if title[0] == "3" and title[0:4]!="3019" and len(title)==5:
#                         # 判断是否为错一练三
#                         # print("titlewrong:",titlewrong)
#                         # print("name:",name)
#                         # print("titlewrong[name]:",titlewrong[name])
#                         title3 = que_type_title("对象",title)
#                         title4=titlepractice(title3)
#                         # print("title,title3:",title,title3)
#                         # print("错一练三，title,title4:",title,title4)
#                         titlewrong[name].append(title4)
#                         # titlewrong=class_title_pra2(titlewrong,name,title)
#                     elif title[0] == "2" or title[0] == "4" or title[0:4]=="3019":
#                         print("题目判断：",title,)
#                         if title[0:4] == "2019" or title[0:4] == "2016" or \
#                                 title[0:4]=="3019" or title[0:4]=="2022":# 判断是否为二次登记
#                             print("题目判断2：", title, )
#
#                             title_registor = {}  # 将错题登记到个人信息中{'金世语': [{'错题': [2019030802], '训练题': ['20190308187']}],
#                             title_registor["错题"] = [title]
#                             titlewrong[name].append(title_registor)
#                             print("二次登记，title_registor",title_registor)
#                         else:#不为二次登记，为第一次登记题目
#
#                             title = main_number + "0" + title
#                             title_registor2 = {}  # 将错题登记到个人信息中{'金世语': [{'错题': [2019030802], '训练题': ['20190308187']}],
#                             title_registor2["错题"] = [title]
#                             titlewrong2[name].append(title_registor2)
#                             print("试卷，title_registor2",title_registor2)
#                     else:
#                         # print("number:",title)
#                         # print("number[0:4]:",title[0:4])
#                         # title = main_number + "0" + title
#                         title = main_number  + title
#                         title_registor={}#将错题登记到个人信息中{'金世语': [{'错题': [2019030802], '训练题': ['20190308187']}],
#                         title_registor["错题"]=[title]
#                         titlewrong[name].append(title_registor)
#                         print("三维设计，title_registor1",title_registor)
#                     # finish=3
#             for key in list(titlewrong.keys()):
#                 if not titlewrong.get(key):
#                     del titlewrong[key]
#             for key in list(titlewrong2.keys()):
#                 if not titlewrong2.get(key):
#                     del titlewrong2[key]
#         # name_finish[name] = finish
#
#     # titlewrong['unit']=unit
#     # print('titlewrong',titlewrong)
#     # print('titlewrong2',titlewrong2)
#     return [titlewrong,titlewrong2]
#     # return [sorted(titlewrong.keys()),sorted(titlewrong2.keys())]
# def stu_title_pra(student_wrongs):
#     """#将一个学生的错题题目和训练题插入word文档"""
#     titlewrongs = []
#     k1=0#为保存进度提醒
#     try:
#         os.mkdir(r'%s\学生错题word文档' % (main_path))
#     except:
#         pass
#     for i in student_wrongs:  # 将学生的错题导入，寻找相应的训练题
#         # print('i:', i)
#         document = Document()
#         k1+=1
#         document.add_heading("%s的错一练三"%i, 0)
#         for l in student_wrongs[i]:
#             # print('l', l)
#             titlewrong = l
#             # print('--------------------titlewrong:', titlewrong)
#             try:
#                 a = titlepractice(titlewrong)
#                 # print('a:',a)
#             except:
#                 continue
#             # print('a:', a)
#             optitlepractice(a)
#             k = 0
#             for j in a:
#                 # print('j,a[j]', j, a[j])
#                 if a[j] != []:
#                     k += 1  # 判断是否具有练习题
#                     # p = document.add_paragraph('题目的优先级：%s'% i)  # 题号
#                     # print("在word中输入训练题目的优先级")
#                     # print('训练题目的优先级：', i)
#             # print('k:----------------',k)
#
#             if k == 1:
#                 # print("没有与该题相对应的训练题")
#                 # wordtitle(titlewrong)
#                 p = document.add_paragraph("titlewrong: %s 没有与该题相对应的训练题" % titlewrong)
#                 titlewrongs.append(titlewrong)
#                 continue
#                 # print("没有与该题相对应的训练题2")
#
#             else:
#                 for j in a:
#                     # print("开始打印题目",'a[j]:',a[j])
#                     if a[j] != []:
#                         k += 1
#                         for m in a[j]:
#                             # print('j,a[i]',j,a[j])
#                             output(m)
#                             p = document.add_paragraph()  # 题号
#                             p_format=p.paragraph_format
#                             p_format.line_spacing =Pt(15)
#                             p_format.space_before = Pt(5)
#                             p_format.space_after=Pt(5)
#                             run = p.add_run('题目类别：%s 题号：%s' % (j, m))
#                             run.space_after = Pt(1400)
#                             run.space_before = Pt(1400)
#                             run.font.name = u'宋体'
#                             run.font.size = Pt(15)
#                             # run.font.
#                             # p = document.add_paragraph('题目的对象：%s 过程:%s  状态:%s  方法:%s  难度:%s'
#                             #                            % (findobj(m, '对象'), findobj(m, '过程'),
#                             #                               findobj(m, '状态'), findobj(m, '方法'), finddif(m)))  # 题号
#                             p = document.add_paragraph('题目的类型：%s '
#                                                        % (findobj(m, '对象'),))  # 题号
#                             p_format = p.paragraph_format
#                             p_format.line_spacing = Pt(15)
#                             p_format.space_before = Pt(5)
#                             p_format.space_after = Pt(5)
#                             wordtitle(document,m)
#         document.save(r'%s\%s%s.docx' % (title_wrong_pra_path,i,student_wrongs['unit']))
#         print('第%d个学生错题保存完成'%k1,i)
#     print("总共有%d题没有对应的训练题分别是%s" % (len(titlewrongs), titlewrongs))
#     return document
# def class_title_wrong(student_wrongs):
#     """#将全班学生的错题打印出来"""
#     # print("class_title_wrong:student_wrongs",student_wrongs)
#     document = Document()
#     # document.add_heading('%s错题.docx' % (student_wrongs['team']), 0)
#     for i in student_wrongs:  # 将学生的错题导入
#         # print('i:', i,'student_wrongs[i]:',student_wrongs[i])
#         def word_add_title_wrong(document):
#             from docx.enum.section import WD_SECTION
#             new_section=document.add_section(WD_SECTION.ODD_PAGE)#增加分页，将每个学生隔开
#             new_section.left_margin = Cm(1.27)
#             new_section.right_margin = Cm(1.27)
#             new_section.top_margin = Cm(1.27)
#             new_section.bottom_margin = Cm(1.27)
#             document.add_heading("%s的“题不二错”" % i,3)
#             for l in student_wrongs[i]:
#                 print('l', l)
#                 titlewrong = l
#                 # print('--------------------titlewrong:', titlewrong)
#                 document.add_heading(' 题号：%s' % (titlewrong),4)
#                 wordtitle(document, titlewrong)
#             # document.add_page_break()  # 增加分页
#         def word_add_title_wrong_answer(document):
#             from docx.enum.section import WD_SECTION
#             new_section = document.add_section(WD_SECTION.ODD_PAGE)  # 增加分页，将每个学生隔开
#             new_section.left_margin = Cm(1.27)
#             new_section.right_margin = Cm(1.27)
#             new_section.top_margin = Cm(1.27)
#             new_section.bottom_margin = Cm(1.27)
#             document.add_heading("%s的“题不二错” 答案" % (i),3)
#             for l in student_wrongs[i]:
#                 # print('l', l)
#                 titlewrong_answer = "A"+l
#                 # print('--------------------titlewrong:', titlewrong)
#                 document.add_heading(' 题号：%s' % (titlewrong_answer),4)
#                 wordtitle(document,titlewrong_answer)
#                 # document.add_page_break()  # 增加分页
#         if i!='team' and i!='unit' and student_wrongs[i]!=[]:
#             print('i:', i, 'student_wrongs[i]:', student_wrongs[i])
#             word_add_title_wrong(document)
#             # word_add_title_wrong_answer(document)# 将学生的错题答案导入
#     # print("总共有%d题没有对应的训练题分别是%s" % (len(titlewrongs), titlewrongs))
#     # print("2class_title_wrong:student_wrongs", student_wrongs)
#     document.save(r'%s\%s错题(含答案).docx'%(title_wrong_path,student_wrongs['team']+student_wrongs['unit']))  # 题目登记中某一表格学生的原题重做
#
#     if i != 'team' and i != 'unit' and student_wrongs[i] != []:
#         document.add_heading('%s错题.docx' % (student_wrongs['team'] ), 0)
#         for i in student_wrongs:  # 将学生的错题导入,生成单个学生的错题
#             document = Document()
#             word_add_title_wrong(document)
#             document.save(
#                 r'%s\%s错题.docx' % (title_wrong_path, i))  # 题目登记中某一表格学生的原题重做
#
#     return document
def class_title_pra(student_wrongs,classs,time_now,author,title_wrong_path,title_wrong_pra_path):
    """#将一个班级的错题题目和训练题插入word文档"""
    # print("student_wrongs:",student_wrongs)
    # print("student_wrongs['team']",student_wrongs['team'])
    team=student_wrongs['team']
    document = Document()
    # section = document.sections[0]
    # #* 设置栏目数为2
    # sectPr = section._sectPr
    # cols = sectPr.xpath('./w:cols')[0]
    # cols.set(qn('w:num'), '2')
    #document.add_heading("%s 的“题不二错”\n时间：%s" % (classs,time_now),1)
    # document = class_title_pra(document, student_wrongs)
    titlewrongs = []

    name_number=1
    len_student_wrongs=len(student_wrongs)
    # print('student_wrongs',student_wrongs)
    for name in student_wrongs:  # 将学生的错题导入，寻找相应的训练题
        progressbarOne['value'] += 1
        label_2.set("%s:将学生--%s题目添加到word文档，%s/%s"%(team,name,name_number,len(student_wrongs)))
        win.update()
        title_answer_types = {}
        title_option_changes = {}

        print("第%s个学生，共%s个"%(name_number,len_student_wrongs),name)
        student_practices={}
        student_practices[name] = []
        # if name!='team' and name!='unit' and student_wrongs[name]!=[] and finish_name_all_dict[name][-1]==1:
        if name!='team' and name!='unit' and student_wrongs[name]!=[]:#不需要筛选功能，上一行为添加筛选功能
            #去除team和unit两个姓名和没有登记错题的学生,以及验证其上一次打印题目是否完成
            from docx.enum.section import WD_SECTION
            from docx.enum.section import WD_SECTION_START
            print("name_number:",name_number)

            if name_number==2:
                new_section=document.sections[0]
            else:
                document.add_page_break()# 增加分页，将每个学生隔开
                document.add_page_break()# 增加分页，将每个学生隔开
                # new_section = document.add_section(WD_SECTION_START.ODD_PAGE)
                # new_section = document.add_section(WD_SECTION.ODD_PAGE)
                # new_section = document.add_section(WD_SECTION.ODD_PAGE)
                # runs=document.sections[-1].footer[0].paragraph[0].runs

            # print("new_section:",   new_section.__dict__,             help(new_section))
            # if len(student_wrongs)==2:
            #     document.add_heading("%s" % notice4, 4)
            # else:
            document.add_heading("%s 的“题不二错”%s" % (name,time_now2), 3)
            # document.add_heading("%s" % author, 4)
            number=0#记录该 学生的题目打印进度
            pra_number=1
            not_exit_number=0
            # wrong_number=0
            sum_number=len(student_wrongs[name])
            print(" 将%s的题目（含训练题）输出到word文档中" % name,end=" ")
            for titlepras in student_wrongs[name]:  # 将每一个学生的题目（含训练题）输出到word文档中
                # print("1")
                # print("student_practices[name]",student_wrongs[name])
                for titlepra in titlepras:
                    # print("2")
                    # print("titlepras",titlepras)
                    if titlepras[titlepra] != [] :
                        # print("3")
                        new_section.left_margin = Cm(1.27)
                        new_section.right_margin = Cm(1.27)
                        new_section.top_margin = Cm(1.0)
                        new_section.bottom_margin = Cm(1.0)
                        # print("titlepras[0]:",titlepras)
                        # print("titlepra:",titlepra)
                        # print("titlepras[titlepra]:",titlepras[titlepra])
                        # print("titlepras[‘错题’]:",titlepras["错题"])
                        title_type=findobj(titlepras[titlepra][0], '对象')
                        if titlepra=="错题" and title_type!="类型暂不清楚":#如果是错题，则打印题目类型
                            document.add_heading('%s' % (findobj(titlepras[titlepra][0], '对象')), 2)
                            # print("title_type:",title_type)
                        for title in titlepras[titlepra]:
                            # print("4")
                            new_section.left_margin = Cm(1.27)
                            new_section.right_margin = Cm(1.27)
                            new_section.top_margin = Cm(1.0)
                            new_section.bottom_margin = Cm(1.0)
                            # heading=document.add_heading('（%d）题库编号：%s 难度：%s' % (pra_number, m,finddif(m)), 3)
                            heading=document.add_heading('%d、题库编号：%s' % (pra_number, title), 3)
                            heading.paragraph_format.line_spacing =1
                            heading.paragraph_format.space_before = Pt(0)
                            heading.paragraph_format.space_after = Pt(0)
                            # if '错题' in titlepra:
                            #     # print("5")
                            #     wrong_number += 1
                            #     # if wrong_number%3==0:
                            #     #     print("第%s题"%wrong_number,end=" ")
                            #     # document.add_heading('%d、题目类别：%s 题号：%s' % (wrong_number,titlepra,m),2)
                            #     document.add_heading('%d、%s %s' % (wrong_number,m,findobj(m,'对象')),2)
                            # elif '训练题' in titlepra:
                            #     # document.add_heading('(%d)%s 题号：%s' % (pra_number,titlepra,m),findobj(m,'对象'),3)
                            #     document.add_heading('(%d)%s,%s,难度:%s' % (pra_number,m,findobj(m,'对象'),finddif(m)),3)
                            pra_number += 1

                            title_option_changes=wordtitle(document,title,title_option_changes)

            pra_number = 1
            # wrong_number = 0
            # print("")
            # def word_title_answer(student_practices,name):
            print(" 将%s的题目答案（含训练题）输出到word文档中" % name,end="")
            # if "模拟" in file_name:
            #     break
            for i in range(len(student_wrongs[name])):
                # for titlepras in student_wrongs[name]:  # 将每一个学生的题目（含训练题）的答案输出到word文档中
                print("开始输入题目的答案",name,titlepras)
                titlepras = student_wrongs[name][i]


                for titlepra in titlepras:
                    if titlepras[titlepra] != [] :
                        #统计题目答案的类型选项更改的word,未更改的word,试卷的选择题，试卷的图片

                        # def wordtitle_answer_find(title):
                        #
                        #     # print("title_option_changes[title]:", title[1:], title_option_changes)
                        #     try:
                        #         title2 = title[1:]
                        #
                        #         # print("title2,title_option_changes:", title2, title_option_changes)
                        #         # print(type(title_option_changes[title2]), title_option_changes[title2])
                        #         if title_option_changes[title2] != None:
                        #
                        #             # def find_change_option_answer(title):
                        #             #     #查找选择题目原来的答案
                        #             #     document_answer = Document(r"word题目汇总\%s.docx" % title)
                        #             #     for paragraph in document_answer.paragraphs:
                        #             #         answer_options = re.findall(r'([A-Z]+)(正确|错误)', paragraph.text)
                        #             #         answer_old = []
                        #             #         for answer_option, correctness in answer_options:
                        #             #             if correctness == '正确':
                        #             #                 answer_old .append( answer_option)
                        #             #     print("选择题的答案是：",answer_old,title)
                        #             #     return answer_old
                        #             # answer_old_change=find_change_option_answer(title)
                        #
                        #
                        #             # 查找选择题目的正确答案
                        #                 #更改答案
                        #             # print("选项已经更改")
                        #             def find_change_option_answer(title):
                        #                 document_answer = Document(r"word题目汇总\%s.docx" % title)
                        #                 options = ["A", "B", "C", "D", "E"]
                        #                 # print("answer:", title_option_changes)
                        #                 title_answer_option = []#题目答案选项，通过答案AB直接得到
                        #                 for paragraph in document_answer.paragraphs:
                        #                     # print(paragraph.text, paragraph.text != None)
                        #                     # print(paragraph.runs[0].text)
                        #                     if "答案" in paragraph.text and paragraph.runs[0].text[0:2] == "答案":
                        #                         for option in options:
                        #                             #查找选择题的答案，未成功
                        #                             if option in paragraph.runs[1].text[0:4]:
                        #                                 title_answer_option.append(option)
                        #                             text2 = paragraph.runs[2].text
                        #                             # print("title_answer_option:", option, title_answer_option, text2)
                        #                             if (paragraph.runs[2].text[0:4] != "解析") and len(text2) < 5:
                        #                                 text3 = paragraph.runs[2].text
                        #                                 if (option in paragraph.runs[2].text[0:4]) and len(text3) < 5:
                        #                                     title_answer_option.append(option)
                        #                                     text4 = paragraph.runs[3].text
                        #                                     if (paragraph.runs[3].text[0:4] != "解析") and len(text4) < 5:
                        #                                         if option in paragraph.runs[3].text[0:4]:
                        #                                             title_answer_option.append(option)
                        #                 # print("title_answer_option2:", title_answer_option)
                        #
                        #                 return title_answer_option
                        #
                        #             title_answer_option = find_change_option_answer(title)
                        #             answers = ""
                        #
                        #             optionss = ""
                        #             options = ["A", "B", "C", "D", "E", ]
                        #
                        #             for i in title_option_changes[title2]:
                        #                 optionss += options[i]
                        #             # print("options:",options)
                        #             title_answer_option_number = []
                        #             for k in title_answer_option:
                        #                 for j in range(len(options)):
                        #                     if options[j] == k:
                        #                         title_answer_option_number.append(j)
                        #             # print("title_answer_option_number:",title_answer_option_number)
                        #             # print("title_option_changes[title2]:",title_option_changes[title2])
                        #             # title_option_changes[title2].sort()
                        #             # print("title_option_changes[title2]:", title_option_changes[title2])
                        #
                        #             for number in title_answer_option_number:
                        #                 for l in range(len(title_option_changes[title2])):
                        #                     # print(number,l,type(number),type(l))
                        #                     if number == title_option_changes[title2][l]:
                        #                         # print(number, l)
                        #                         answers += options[l]
                        #             l = list(answers)  # 答案重新排序
                        #             l.sort()
                        #             answers = "".join(l)
                        #             change_text = "新答案：%s 已经将原题目选项%s更改为ABCD" % (answers, optionss)
                        #             # print("新答案：%s 已经将原题目选项%s更改为ABCD" % (answers, optionss))
                        #         else:
                        #             # print("选项未更改")
                        #             change_text = ""
                        #
                        #     except:
                        #         change_text = ""
                        #
                        #
                        #     return change_text
                        # for title in titlepras[titlepra]:
                        #     title="A"+str(title)
                        #     try:
                        #         if wordtitle_answer_find(title)[0:3] == "新答案":
                        #             title_answer_types[title] = "word_choice"
                        #
                        #         else:
                        #             Document(r"word题目汇总\%s.docx" % title)
                        #             title_answer_types[title] = "word_calculation"
                        #     except:
                        #         try:
                        #             if wordtitle_answer_find_excel(title)!="未找到答案":
                        #                 title_answer_types[title] = "text"
                        #             else:
                        #                 title_answer_types[title] = "img1"
                        #         except:
                        #             title_answer_types[title] = "img2"

                            # print("title_answer_types[title]:",title,title_answer_types[title] )
                        # print("titlepras[titlepra]:",titlepras)


                        for m in titlepras[titlepra]:
                            # print('j,a[i]',i)
                            m='A'+str(m)
                            # print("开始输入题目的答案",m)
                            # output(m)
                            new_section.left_margin = Cm(1.27)
                            new_section.right_margin = Cm(1.27)
                            new_section.top_margin = Cm(1.0)
                            new_section.bottom_margin = Cm(1.0)
                            # document.add_heading('题目：%s答案' % (wrong_number,titlepra,m),2)

                            # answer_change_text=wordtitle_answer_find(m)
                            # if i>0:
                            #     # print("i:",i,titlepras[i-1])
                            #     #题号前加一个A表示答案Answer
                            #     previous_title = "A" + str(student_wrongs[name][i-1]["错题"][0])
                            #     # print("previous_title:",previous_title)
                            #     # print(title_answer_types)
                            #     # print("类型：",title_answer_types[previous_title])
                            #     # print("类型：",title_answer_types[previous_title],title_answer_types[m])
                            #     if title_answer_types[previous_title] == "text" and \
                            #             title_answer_types[m] == "text":
                            #         # print("本题和上题答案都是text类型:", title,wordtitle_answer_find_excel(m))
                            #         paragraph.add_run('%d、答案：%s' % (pra_number,wordtitle_answer_find_excel(m)))
                            #         # print("AA")
                            #         pra_number += 1
                            #         continue
                            paragraph=document.add_paragraph( '       %d、' % (pra_number))#行首空几格
                            # paragraph=document.add_paragraph( '%d、答案：' % (pra_number))
                            # if answer_change_text!="":
                            #     paragraph = document.add_paragraph(text='%s' % ( answer_change_text))
                            #     paragraph.paragraph_format.space_after = Pt(0)
                            pra_number += 1
                            wordtitle_answer(document, m,  title_option_changes,pra_number)
                            # wordtitle_answer(document, m, title_answer_types, paragraph,title_option_changes)
                            # titlepras=student_practices[name]document,title,title_option_changes,pra_number
                    else:
                        titlewrongs.append(titlepras[titlepra])
            # finish_name_all_dict[name][-1]=3
            # print("")
        # if name=="赵安宁":
        #     # 生成题目的ppt文档
        #     yesno = tk.messagebox.askyesno("是否生成pptx文件")
        #     if yesno:
        #         document_pptx = pptx.Presentation()
        #         slide = document_pptx.slides.add_slide(document_pptx.slide_layouts[0])
        #
        #         # 在幻灯片上加入标题
        #         title = slide.shapes.title
        #         title.text = "错题重做"
        #
        #         # 在幻灯片上加入副标题
        #         subtitle = slide.placeholders[1]
        #         subtitle.text = "Python PPTX subtitle"
        #
        #         # 在幻灯片上添加一个文本框
        #         body_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
        #         tf = body_shape.text_frame
        #         tf.text = "Python PPTX textbox"
        #
        #         # 在幻灯片上插入图片
        #         for title in
        #             slide = document_pptx.slides.add_slide(document_pptx.slide_layouts[1])
        #         left = Inches(1)
        #         top = Inches(2)
        #         pic = slide.shapes.add_picture('2.png', left, top)
        #
        #         # 将 PPT 保存到本地文件
        name_number+=1

    print("总共有%d题没有对应的训练题分别是%s" % (len(titlewrongs), titlewrongs))
    # document_name = r'%s\%s的“题不二错”%s.docx' % (title_wrong_pra_path, team, time_now)
    def AddFooterNumber(run):
        """在页脚中插入页码"""
        fldChar1 = OxmlElement('w:fldChar')  # creates a new element
        fldChar1.set(qn('w:fldCharType'), 'begin')  # sets attribute on element
        instrText = OxmlElement('w:instrText')
        instrText.set(qn('xml:space'), 'preserve')  # sets attribute on element
        instrText.text = 'Page'
        fldChar2 = OxmlElement('w:fldChar')
        fldChar2.set(qn('w:fldCharType'), 'separate')
        t = OxmlElement('w:t')
        t.text = "Seq"
        fldChar2.append(t)
        fldChar4 = OxmlElement('w:fldChar')
        fldChar4.set(qn('w:fldCharType'), 'end')
        r_element = run._r
        r_element.append(fldChar1)
        # print("r_element1:",r_element.__doc__)

        r_element.append(instrText)
        # print("r_element2:", r_element, type(r_element))
        # print("r_element2:", r_element, type(r_element[0]))
        # print("r_element2:", r_element, len(r_element[0]))

        r_element.append(fldChar2)
        # print("r_element3:", r_element)
        r_element.append(fldChar4)
        # print("r_element4:", r_element)
    def InsertPageNumber(Doc):
        """插入页脚内容"""
        footer = Doc.sections[0].footer  # 获取第一个节的页脚
        footer.is_linked_to_previous = True  # 编号续前一节
        paragraph = footer.paragraphs[0]  # 获取页脚的第一个段落
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER  # 页脚居中对齐
        run_footer = paragraph.add_run()  # 添加页脚内容
        AddFooterNumber(run_footer)
        font = run_footer.font
        font.name = 'Times New Roman'  # 新罗马字体
        font.size = Pt(10)  # 10号字体
        font.bold = True  # 加粗
    InsertPageNumber(document)
    # footer_text = document.sections[0].footer.paragraphs[0].text
    # help(document.pages)
    # print("footer_text:", footer_text)
    # if int(footer_text) % 2 == 0:
    #     print("int(footer_text)%2:", int(footer_text) % 2)
    #     document.add_page_break()  # 增加分页，将每个学生隔开
    # print("分页成功:", footer_text)
    # document.save(r'%s\%s的“题不二错”%s.docx'%(title_wrong_pra_path,team,time_now))  # 题目登记中某一表格学生的错一练三
    document.save(r'%s\错题重做%s.docx'%(title_wrong_pra_path,time_now))  # 题目登记中某一表格学生的错一练三
    # os.startfile(r'%s\%s的“题不二错”%s.docx'%(title_wrong_pra_path,team,time_now))
    os.startfile(r'%s\错题重做%s.docx'%(title_wrong_pra_path,time_now))
    #保存至云盘
    # try:
    #     print("main_path:",main_path)
    #     doc_network_disk=Document(r"安装必备文件\network_disk.docx")
    #     network_disk_path=doc_network_disk.paragraphs[0].text
    #     print("network_disk_path:",network_disk_path)
    #
    #     document.save(r'%s/%s的“题不二错”%s.docx'%(network_disk_path,team,time_now))
    #     tk.messagebox.showinfo(title="保存成功", message=        "文档保存至云盘成功")
    #     # os.startfile(r'%s/%s的“题不二错”%s.docx'%(network_disk_path,team,time_now))
    # except Exception as e:
    #     print("保存文件至云盘失败",e)
    #     yesno=tk.messagebox.askyesno(title="错误提示", message=
    #     "文档保存至云盘失败，是否打开云盘并登陆，账号为1124022，密码为12345678")
    #     if yesno==True:
    #         try:
    #             import subprocess
    #             # 定义要打开的软件路径或名称
    #             software_path = "C:\Program Files\Richtech\EDataStore\EmployeeLogin.exe"
    #             # 调用系统默认程序打开指定的文件或URL
    #             subprocess.Popen(software_path)
    #             yesno2 = tk.messagebox.askyesno(title="提示", message=
    #             "继续保存文档至云盘")
    #             if yesno2:
    #                 document.save(r'U:\%s的“题不二错”%s.docx' % ( team, time_now))
    #                 os.startfile(r'U:\%s的“题不二错”%s.docx' % ( team, time_now))
    #
    #         except FileNotFoundError:
    #             print("未能找到该软件")
    #             tk.messagebox.showinfo(title="错误提示", message=
    #             "文档保存至云盘失败，请前往“%s”查看"%(title_wrong_pra_path))
    #             os.startfile(r'%s\%s的“题不二错”%s.docx'%(title_wrong_pra_path,team,time_now))


    print(r'前往“%s\%s的“题不二错””查看'%(title_wrong_pra_path,team))
    document_name,doc_path=r'%s的“题不二错”%s.docx'%(team,time_now),r'%s\%s的“题不二错”%s.docx'%(title_wrong_path,team,time_now)
    # document_name,title_wrong_pra_path=r'%s\%s的“题不二错”%s.docx'%(team,time_now),title_wrong_path
    # return document_name,doc_path
    return document_name,team,title_wrong_path
def title_not_exit_print(title_not_exit):
    """#将题目输出过程中不存在的题目显示出来，以供查找是否出现题目未录入的情况"""
    set01 = set(title_not_exit)
    print(set01)
    dict01 = {}
    for item in set01:
        dict01.update({item: title_not_exit.count(item)})
    print(dict01)
    print('title_not_exit:', sorted(dict01.items(), key=lambda x: x[1], reverse=True))
def class_register(*args):
    """登陆班级"""

    grade=combobox_grade.get()
    classs=combobox_class.get()


    # print(combobox_grade.get(),combobox_class.get(),combobox_name.get())
    # print("names:",classs,Class(grade, classs).names,classs)
    names=Class(grade, classs).names
    combobox_name['values'] = (names)  # 设置下拉框的选项列表
    combobox_name.current(0)
    combobox_name.set(names[0])
    students=Student(grade,classs,names[0])
    # students.title_message

    # print("students::",students.__dict__)
@handle_exceptions(logger)
def register(*args):
    """登陆"""
    grade=combobox_grade.get()
    classs=combobox_class.get()
    name=combobox_name.get()

    # print(combobox_grade.get(),combobox_class.get(),combobox_name.get())
    # print("names:",Class(grade, classs).names,classs)
    names=Class(grade, classs).names
    combobox_name['values'] = (names)  # 设置下拉框的选项列表
    # combobox_name.current(0)
    # combobox_name.set(names[0])
    students=Student(grade,classs,name)
    # students.title_message

    # print("students::",students.__dict__)
    return students
def delete_title_show(title_show):
    """清空题目内容模块的展示内容"""
    try:
        title_show.delete("0.0", tk.END)
    except Exception as e:
        print("删除失败", e)
def input_title_test(*args):
    """输入题目，对题目是否存在进行判断，存在则将内容显示在屏幕上"""
    # all_titles=pd.read_excel(r"word题目汇总/题目清单.xlsx")
    #将表格中的所有数据全部转换为str类型
    # all_titles=pd.DataFrame(pd.read_excel(r"word题目汇总/题目清单.xlsx")).astype('str')
    all_titles=pd.DataFrame(pd.read_excel(r"word题目汇总/题目清单.xlsx"))
    title=entry_input.get()
    # print("title:",title)
    # print("all_titles:",all_titles,type(entry_input.get()),all_titles.iloc[0,0],type(all_titles.iloc[0,0]))
    # print(entry_input.get(),all_titles['题号'])
    # print(all_titles['题号'].isin([int(title)]))
    if (all_titles['题号'].isin([title])).any():#判断一个数据是否在表格中
        # print("题号正确")
        #预览题目
        title_doc=Document(r"word题目汇总/%s.docx"%title)
        label_title_text=tk.Label(win,text="%s信息如下"%title)
        text=""
        # title_show.delete(1)
        text="你已经登记的题号："
        students=register()
        for title in students.title_message.loc[:,students.name]:
            if str(title)!="nan":
                text+=str(title)
                text+="  "
        # print(text)
        # title_show.insert(tk.INSERT,text)
        delete_title_show(title_show)

        text+="\n你即将添加的题目：\n"

        for paragrah in title_doc.paragraphs:
            # print(paragrah.text)
            text+=paragrah.text
        # print("title_doc.text:",text)
        title_show.insert(tk.INSERT,  text)
        return True


    elif len(title)>3:
        print("输入题号错误1")
        delete_title_show(title_show)
        return False

def input_title_test2(title):
    """输入题目，对多选框题目是否存在进行判断，存在则将内容显示在屏幕上"""
    # all_titles=pd.read_excel(r"word题目汇总/题目清单.xlsx")
    #将表格中的所有数据全部转换为str类型
    print("title::::",title)
    all_titles=pd.DataFrame(pd.read_excel(r"word题目汇总/题目清单.xlsx")).astype(str)
    # print("title:",title)
    # print("all_titles:",all_titles,type(entry_input.get()),all_titles.iloc[0,0],type(all_titles.iloc[0,0]))
    # print(entry_input.get(),all_titles['题号'])
    # print(all_titles['题号'].isin([int(title)]))
    if (all_titles['题号'].isin([title])).any():#判断一个数据是否在表格中
        print("题号正确")
        # word_to_pic(title)
        #预览题目
        # title_doc=Document(r"word题目汇总/%s.docx"%title)
        # label_title_text=tk.Label(win,text="%s信息如下"%title)
        # text=""
        # title_show.delete(1)
        text="你已经登记的题号："
        students=register()
        for ti in students.title_message.loc[:,students.name]:
            if str(ti)!="nan":
                text+=str(ti)
                text+="  "
        print("text:",text)
        # title_show.insert(tk.INSERT,text)
        delete_title_show(title_show)

        # text+="\n你即将添加的题目：\n"

        # for paragrah in title_doc.paragraphs:
            # print(paragrah.text)
            # text+=paragrah.text
        # print("title_doc.text:",text)
        title_show.insert(tk.INSERT, text)


        student = Student(combobox_grade.get(), combobox_class.get(), combobox_name.get()
                          ).select_title(title)
        #
        # entry_input.delete(0, tk.END)
        # title_show.update()


    else:
        print("输入题号错误2",title)
        delete_title_show(title_show)
        tk.messagebox.showinfo(title="错误提示", message=
        "您登记的错题%s题库中暂时还没有收录，请检查是否录入错误，或者与管理员联系"%title)
def input_title_save(envent):
    """按回车键的时候保存题目，展示题目列表，并清空内容"""
    title = entry_input.get()
    all_titles=pd.DataFrame(pd.read_excel(r"word题目汇总/题目清单.xlsx")).astype('str')
    # print("title:",title)
    # print("all_titles:",all_titles,type(entry_input.get()),all_titles.iloc[0,0],type(all_titles.iloc[0,0]))
    # print(entry_input.get(),all_titles['题号'])
    # print(all_titles['题号'].isin([int(title)]))
    if not (all_titles['题号'].isin([title])).any():#判断一个数据是否在表格中
        tk.messagebox.showinfo(title="错误提示", message=
        "提示：您登记的错题%s题库中暂时还没有收录，请检查是否录入错误，或者与管理员联系"%title)
    else:
        student = Student(combobox_grade.get(), combobox_class.get(), combobox_name.get()
                      ).register_title(title)
        entry_input.delete(0,tk.END)
# def admin():
#     """管理员录入题目至题库"""
#     admin_window=tk.Toplevel(win)
#     admin_window.title("管理员将题目录入题库")
#     label1=tk.Label(text="当前的文件是：%s"file)
#     label2=tk.Label(text="当前的段落是：%s"file)
#     label2=
#     button_file=tk.Button()
#     for root, dirs, files in os.walk(r"word题目原始文档"):
#         print("[dirs,files]:",[dirs,files])
#         # return [dirs,files]
#         for file in files:
#             doc=Document(r"word题目原始文档/%s"%file)
#             # label=tk.Label(admin_window,text=file)
#             # label.pack()
#             len_para=len(doc.paragraphs)
#             text1 = tk.StringVar()
#             i=0
#             para=doc.paragraphs[i]
#
#             print("段落:")
#             para_text=""
#             for run in para.runs:
#                 # print(run.text,end="")
#                 para_text+=run.text
#             # if para_text!="":
#
#             print("para_text:",para_text)
#             text1.set(para_text)
#             print("text1:",text1)
#             label=tk.Label(admin_window, textvariable=text1,)
#             # label=tk.Label(admin_window, text=text,)
#
#             label.pack()
#             # text.set('new')
#
#
#     admin_window.mainloop()
class Class:
    """建立一个班级的信息"""

    def __init__(self, grade, classs):
        self.classs = classs
        self.grade = grade
        students_list = pd.DataFrame(pd.read_excel(r"sys_files/学生名单.xlsx"))

        students_list = students_list.set_index("班级")
        students_list = students_list.loc[self.classs, "姓名"]
        # print("students_list:",classs,type(students_list),students_list,classs)
        self.names = []
        if type(students_list)==str:
            self.names.append(students_list)
        else:
            for i in students_list:
                self.names.append(i)
        # Class(self.grade, self.classs)
        # self.title_message = pd.DataFrame(columns=names)
        # self.names=names
        self.students=pd.DataFrame(columns=self.names)
        # print("self.students:",self.students)

class Student:
    #建立一个学生的信息库，包含错题，登记时间，各次考试的错题，打印的次数
    def __init__(self,grade,classs,name):
        self.name=name
        self.grade=grade
        self.classs=classs
        delete_title_show(title_show)


        try:
            self.title_message=pd.DataFrame(pd.read_excel(r"%s\%s%s%s错题.xlsx" %(
                student_wrong_path,time_now2,grade,classs)))
        except:

            students_list=pd.DataFrame(pd.read_excel(r"sys_files/学生名单.xlsx"))
            print(students_list)
            students_list=students_list.set_index("班级")
            students_list=students_list.loc[self.classs,"姓名"]
            self.names=[]
            if type(students_list) == str:
                self.names.append(students_list)
            else:
                for i in students_list:
                    self.names.append(i)
            # Class(self.grade,self.classs)
            self.title_message=pd.DataFrame(columns=self.names)
            # self.title_message.to_excel(r"%s\%s%s%s错题.xlsx" %(
            #     student_wrong_path,time_now2,grade,classs))
            # os.startfile(r"%s\%s%s%s错题.xlsx" %(
            #     student_wrong_path,time_now2,grade,classs))
        # print("self.title_message:",self.title_message,self.name)
        print("self.title_message.loc[self.name]:",self.title_message.loc[:,self.name])
        text="你已经登记的题号："
        for title in self.title_message.loc[:,self.name]:
            # print("title:::",title,type(title),type(str(title)))
            if str(title)!="nan":
                text+=str(title)
                text+="  "
        # print("text::",text)
        title_show.insert(tk.INSERT,text)
    def register_title(self,title):
        """登记题目至excel表格"""
        # print("self.title_message:",self.title_message)
        all_titles = pd.DataFrame(pd.read_excel(r"word题目汇总/题目清单.xlsx"))
        title = entry_input.get()
        # print("title:",title)
        # print("all_titles:",all_titles,type(entry_input.get()),all_titles.iloc[0,0],type(all_titles.iloc[0,0]))
        # print(entry_input.get(),all_titles['题号'])
        # print(all_titles['题号'].isin([int(title)]))
        # self.title_message.loc[:, self.name]
        # print("111111",((self.title_message.loc[:, self.name]).isin([title])).any())
        # print(((self.title_message.loc[:, self.name]).isin([title])))
        # print("222222",((self.title_message.loc[:, self.name])))
        # print("title:",title)
        if ((self.title_message.loc[:, self.name]).isin([title])).any():  # 判断一个数据是否在表格中
            print("该题目已经存在",title)
        else:
            self.title_message.loc[len(self.title_message),self.name]=str(title)
            # print("self.title_message:",self.title_message)
            # print("save success:",title)
        # print(self.title_message)

        #显示窗口和excel表都要更新
        text = "你已经登记的题号："
        for title in self.title_message.loc[:, self.name]:
            # print("title:::",title,type(title),type(str(title)))
            if str(title) != "nan":
                text += str(title)
                text += "  "
        print("text::", text)
        title_show.insert(tk.INSERT, text)
        self.title_message=self.title_message.set_index(self.name)
        self.title_message.to_excel(r"%s\%s%s%s错题.xlsx" %(
                student_wrong_path,time_now2,self.grade,self.classs))
    def select_title(self,title):
        """登记题目至excel表格"""
        # print("self.title_message:",self.title_message)
        all_titles = pd.DataFrame(pd.read_excel(r"word题目汇总/题目清单.xlsx"))
        # title = entry_input.get()
        # print("title:",title)
        # print("all_titles:",all_titles,type(entry_input.get()),all_titles.iloc[0,0],type(all_titles.iloc[0,0]))
        # print(entry_input.get(),all_titles['题号'])
        # print(all_titles['题号'].isin([int(title)]))
        # self.title_message.loc[:, self.name]
        print("111111",((self.title_message.loc[:, self.name]).isin([title])).any())
        print(((self.title_message.loc[:, self.name]).isin([title])))
        print("222222",((self.title_message.loc[:, self.name])))
        if ((self.title_message.loc[:, self.name]).isin([title])).any():  # 判断一个数据是否在表格中
            print("该题目已经存在",title)
        else:
            self.title_message.loc[len(self.title_message),self.name]=str(title)
            print("self.title_message:",self.title_message)
            print("save success:",title)
        # print(self.title_message)

        #显示窗口和excel表都要更新
        text = "你已经登记的题号："
        for title in self.title_message.loc[:, self.name]:
            # print("title:::",title,type(title),type(str(title)))
            if str(title) != "nan":
                text += str(title)
                text += "  "
        print("text::", text)
        title_show.insert(tk.INSERT, text)
        self.title_message=self.title_message.set_index(self.name)
        self.title_message.to_excel(r"%s\%s%s%s错题.xlsx" %(
                student_wrong_path,time_now2,self.grade,self.classs))
    def delete_title(self,title):
        """删除excel表中的题目"""
        #pandas查找单元格
        self.title_message.loc[self.title_message[self.name]==title,self.name]=""#pandas替换
        self.title_message=self.title_message.set_index(self.name)
        self.title_message.to_excel(r"%s\%s%s%s错题.xlsx" %(
                student_wrong_path,time_now2,self.grade,self.classs))
    def wrong_season_next(self,title):
        """记录学生的错题的错误原因"""
        # title = wrong_sea[0]
        # i = wrong_sea[2]
        # title=self.title_all.iloc[0]
        wrong_season_text = combobox_wrong_season.get()
        try:
            stu_wro_season=pd.DataFrame(pd.read_excel(r"F:\个性化题库20180423\title practice"
                                r"\安装必备文件\学生\%s.xlsx"%self.name))
        except Exception as e:
            print("打开学生错题记录失败:",e)
            stu_wro_season = pd.DataFrame(columns=["题号","登记时间","错误原因"])
        print("stu_wro_season:", stu_wro_season)

        new_row = pd.DataFrame([[title,time_now,wrong_season_text]],
                               columns=stu_wro_season.columns)
        stu_wro_season = pd.concat([stu_wro_season, new_row], ignore_index=True)
        stu_wro_season=stu_wro_season.set_index("题号")
        print("stu_wro_season:", stu_wro_season)
        makedir(r"安装必备文件\学生")
        stu_wro_season.to_excel(r"安装必备文件\学生\%s.xlsx"%self.name)

        makedir(r"cach")
        makedir(r"cach\%s"%self.name)
        try:
            self.title_wrong =Document(r"cach\%s\wrong_season_cach.docx"%self.name)
            self.title_wrong.paragraphs[0].add_run(title)
            self.title_wrong.paragraphs[0].add_run(wrong_season_text)
            self.title_wrong.save(r"cach\%s\wrong_season_cach.docx"%self.name)
        except Exception as e:
            print("文件不存在:",e)
            self.title_wrong = Document()
            self.title_wrong.add_paragraph().add_run(title)
            self.title_wrong.add_paragraph().add_run(wrong_season_text)
            self.title_wrong.save(r"cach\%s\wrong_season_cach.docx"%self.name)
        return self.title_wrong.paragraphs[0].text

        # wrong_season(self.name)
        # print("self.title_all:",len(self.title_all),self.title_all)
        # self.title_all=self.title_all.iloc[1:]
        # title=self.title_all.iloc[0]
        # index = title.index("K")  # 寻找字符的索引位置
        # result = title[index + 1:]  # 使用切片截取字符后面的内容
        # # wrong_title_text.set("题目：第%s题,题库编号:%s" % (title[index + 1:] ,title))

def make_word_question_bank():
    """通过word文件制作word题库，测试资料为步步高"""
    input_doc = Document(r"test\第六章 1　圆周运动.docx")

    list_title = []
    for i in range(20):
        list_title.append("%d." % i)
        list_title.append("%d、" % i)
        list_title.append("%d．" % i)
        # list_title.append("i．")
    # print(list_title)
    title_para_number_all = {}
    title_number = 1
    title_pre = "20231261K"
    for para_number in range(len(input_doc.paragraphs)):

        # print(para.text[0:2] in list_title)
        para = input_doc.paragraphs[para_number]
        # print(para.text)
        # print(para.text[0:2])
        # print(para.text[0:3])
        # print(para.text[0:4])
        if para.text[0:2] in list_title or para.text[0:3] in list_title or para.text[0:4] in list_title:
            title_para_number_all[title_number] = para_number  # 定位每一题的起始段落位置
            print("title_para_number_all:", title_para_number_all)
            title_number += 1
    title_para_number_all[title_number] = len(input_doc.paragraphs)  # 定位每一题的起始段落位置
    #
    # def makedir(path):
    #     """尝试创建文件夹，有则不创建"""
    #     try:
    #         os.mkdir(path)
    #     except:
    #         pass
    #
    # # 开始复制
    # def copy_picture(paragraph, output_doc):
    #     # 通过rid属性查询图片，并完成复制
    #     import xml.etree.cElementTree as ET
    #     namespace = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
    #                  'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    #                  'v': "urn:schemas-microsoft-com:vml",
    #                  'wp': "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing",
    #                  'a': "http://schemas.openxmlformats.org/drawingml/2006/main",
    #                  'pic': "http://schemas.openxmlformats.org/drawingml/2006/picture"}
    #
    #     # 定义了namespace，建议参看xml.etree.cElementTree中关于namespace的部分，添加了较多，可以自定义修改
    #     # 打开了对应的word文档，而不是上面讲到的.xml，但是本质是一样的。
    #     def get_img(root_element, target_tag, target_attribute, out_list):
    #         # 遍历得到目标属性target_attribute，得到属性value,写入out_list
    #         # 通过解读word的document.xml，得知与图片相关的rId属性值会出现在两个标签中
    #         # v:imagedata和a:blip
    #         # <v:imagedata r:id="rId8" o:title=""/>，<a:blip r:embed="rId8">
    #         for child in root_element:
    #             tag = child.tag
    #             attribute = child.attrib
    #             if tag in target_tag and target_attribute in child.attrib.keys():
    #                 target_value = child.attrib[target_attribute]
    #                 # print(target_value)
    #                 out_list.append(target_value)
    #             else:
    #                 get_img(child, target_tag, target_attribute, out_list)
    #
    #     # xml_element = []
    #
    #     # doc.paragraphs Proxy object wrapping <w:p> element.
    #     # xml_element.append(paragraph._element.xml)
    #
    #     # for tbl in doc.tables:
    #     #     # print(tbl)
    #     #     # doc.tables Proxy class for a WordprocessingML <w:tbl> element.
    #     #     xml_element.append(tbl._element.xml)
    #     rId = []
    #     id = []
    #     # for element in xml_element:
    #     root = ET.fromstring(paragraph._element.xml)
    #     target_tag = ['{urn:schemas-microsoft-com:vml}imagedata',
    #                   '{http://schemas.openxmlformats.org/drawingml/2006/main}blip']
    #     # 即v:imagedata和a:blip
    #     target_attribute1 = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
    #     target_attribute2 = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    #     # 即:r:embed和r:id
    #     get_img(root, target_tag, target_attribute1, rId)
    #     get_img(root, target_tag, target_attribute2, id)
    #
    #     dd = id + rId
    #     imgs = []
    #
    #     for rid in dd:
    #         imgs.append(paragraph.part.related_parts[rid])
    #         # print(paragraph.part.related_parts[rid])
    #
    #     # 1.4 保存图片到本地
    #     i = 1
    #     #
    #     for img in imgs:
    #         makedir(r"test")
    #         # print("img:::",img,)
    #         f = open(r"test/img%d.jpg" % i, 'wb')
    #         # help(f)
    #         # print("type(f):",type(f))
    #
    #         f.write(img.blob)
    #
    #         f.close()
    #         output_doc.add_picture(r"test/img%d.jpg" % i)
    #         os.remove(r"test/img%d.jpg" % i)
    #         i += 1
    #         # print("i:", i)
    #
    # def copy_docx_run(run, run_text_new, output_para):
    #     output_run = output_para.add_run("%s" % run_text_new)
    #     output_run.style.name = run.style.name
    #     # print("style_name:",run.style.name)
    #     # output_run.bold = run.bold
    #     output_run.bold = False  # 所有的字体都不用粗体
    #     # print("run.bold:",run.bold)
    #     # Run's italic data
    #     output_run.italic = run.italic
    #     output_run.font.size = run.font.size
    #     # print("字体：", run.font.name, run.text)
    #     # output_run.font.name = run.font.name
    #     output_run.font.name = "Times New Roman"
    #     # 设置中文字体使用字体2->宋体,且要先执行上一句定义字体
    #     output_run._element.rPr.rFonts.set(qn('w:eastAsia'), u'宋体')
    #     # print("run.size:",run.size)
    #     # Run's underline data
    #     output_run.underline = run.underline
    #     output_run.font.subscript = run.font.subscript
    #     output_run.font.superscript = run.font.superscript
    #
    #     # Run's color data
    #     output_run.font.color.rgb = run.font.color.rgb
    #     # Run's font data
    #     output_run.style.name = run.style.name
    #     # print("run.style.name:",run.style.name)
    #     # output_run.style.name = run.style.name
    #
    #     # 开始复制公式
    #
    #     try:
    #         # if field_number == 0 and run._element.xpath('w:fldChar')[0].attrib[qn("w:fldCharType")] == "begin":
    #         if run._element.xpath('w:fldChar')[0].attrib[qn("w:fldCharType")] == "begin":
    #             # print("开始复制公式：")
    #             fldChar1 = OxmlElement('w:fldChar')  # creates a new element
    #             fldChar1.set(qn('w:fldCharType'), 'begin')  # sets attribute on element
    #             output_run._element.append(fldChar1)
    #             # print("复制后：",output_run._element.xpath('w:fldChar')[0].attrib[qn("w:fldCharType")])
    #             # field_number = 1
    #     except:
    #         pass
    #
    #     try:
    #         # 复制公式内容
    #         # print("run._element.xpath('w:instrText')[0].text:::",run._element.xpath('w:instrText')[0].text)
    #         # print("run._element.xpath('w:instrText')[0].text:::",type(run._element.xpath('w:instrText')[0].text),field_number)
    #         # if field_number == 1 and run._element.xpath('w:instrText')[0].text:
    #         if run._element.xpath('w:instrText')[0].text:
    #             # print(11111)
    #             # field_text=run._element.xpath('w:instrText')[0].text
    #             # fldChar3 = OxmlElement('w:instrText')
    #             instrText = OxmlElement('w:instrText')
    #             instrText.set(qn('xml:space'), 'preserve')  # sets attribute on element
    #             instrText.text = "%s" % (run._element.xpath('w:instrText')[0].text)
    #             output_run._element.append(instrText)
    #
    #     except:
    #         pass
    #     try:
    #         # if field_number == 1 and run._element.xpath('w:fldChar')[0].attrib[qn("w:fldCharType")] == "end":
    #         if run._element.xpath('w:fldChar')[0].attrib[qn("w:fldCharType")] == "end":
    #             fldChar2 = OxmlElement('w:fldChar')  # creates a new element
    #             fldChar2.set(qn('w:fldCharType'), 'end')  # sets attribute on element
    #             output_run._element.append(fldChar2)
    #             # 结束公式复制
    #             # field_number = 0
    #     except:
    #         pass
    #     return output_para

    def copy_paragraph_edit_first(input_para, output_doc):
        """复制段落至新文档，复制段落时去除题首的题号和题目来源 """
        para_new = output_doc.add_paragraph()
        print("input_para.text", input_para.text)
        first_text = ""
        begin_number = 0
        run_end_number = 0
        text_end_number = 0
        # 更改题首的序号
        for run_number in range(len(input_para.runs)):
            first_text += input_para.runs[run_number].text
            print("first_text:", first_text)
            if first_text in list_title:
                begin_number = run_number + 1
                print("input_para.runs[begin_number].text:", input_para.runs[begin_number].text)
                break
            else:
                begin_number = -1
        # 去除题目开始的题目来源，一般以（）形式存在。
        # print("input_para.runs[0].text:",input_para.runs[0].text)
        try:
            # if input_para.runs[begin_number].text[0] in ["（","("] and input_para.text[0] not in ["（","("]:
            if input_para.runs[begin_number].text[0] in ["（", "("]:
                # print("11")
                for run_number in range(begin_number, len(input_para.runs)):
                    # print("22",input_para.runs[run_number].text)
                    for text_number in range(len(input_para.runs[run_number].text)):
                        # print("333")
                        if input_para.runs[run_number].text[text_number] in [")", "）"]:
                            # print("44")
                            run_end_number = run_number
                            text_end_number = text_number
                            if text_end_number == len(input_para.runs[run_end_number].text):  # 块的最后一个字符
                                run_end_number += 1
                                text_end_number = 0
                            else:
                                text_end_number += 1  # 非块的最后一个字符
                            break
        except Exception as f:  # 插入图片题目
            print("插入文字答案失败，失败原因：", f.__traceback__.tb_lineno, f, )  # 问题行数
        # 根据第一个）的位置确定开始复制正文文本的块和文字，是否为块的最后一个字符

        print("begin_number,run_end_number,text_end_number:", begin_number, run_end_number, text_end_number)

        # 从正式开始的文字块开始复制，第一个run可能有部分不需要复制
        # 第一个run,含有）
        if text_end_number == 1 and begin_number > -1:
            for number in range(text_end_number + 1, len(input_para.runs[run_end_number].text)):
                print("input_para.runs[number].text:", input_para.runs[number].text)
                run = input_para.runs[number]
                copy_docx_run2(run, run.text[number], para_new)
            # 后续的run，不含有）
            for number in range(run_end_number + 1, len(input_para.runs)):
                print("number:", number)
                run = input_para.runs[number]
                copy_docx_run2(run, run.text, para_new)
        else:  # 复制整段
            for run in input_para.runs:
                copy_docx_run2(run, run.text, para_new)

            # copy_picture(input_para,output_doc)

    # print(title_para_number_all)
    for title_number in title_para_number_all:
        if title_number != len(title_para_number_all):
            print(title_para_number_all[title_number], title_number)
            # print(title_para_number_all[title_number+1],title_number)
            output_doc = Document()
            title = title_pre + str(title_number)
            output_doc.add_paragraph(title)
            for title_number_add in range(title_para_number_all[title_number], title_para_number_all[title_number + 1]):
                # 为考点的提示文字删除
                if "考点一" not in input_doc.paragraphs[title_number_add].text or (
                        "考点二" not in input_doc.paragraphs[title_number_add].text) or (
                        "考点三" not in input_doc.paragraphs[title_number_add].text) or (
                        "考点四" not in input_doc.paragraphs[title_number_add].text) or (
                        "考点五" not in input_doc.paragraphs[title_number_add].text):
                    copy_paragraph_edit_first(input_doc.paragraphs[title_number_add], output_doc, )
                # 最后一段为图片，则删除，不是题目的图片，比如能力提升，尖子生选练等
                if title_number_add != title_para_number_all[title_number + 1] - 1:
                    copy_picture(input_doc.paragraphs[title_number_add], output_doc)
            for para in output_doc.paragraphs:
                # print("para.text:", para.text, len(para.runs))
                if len(para.text) <= 1:
                    print("文本小于等于1")
                if len(para.text) <= 1 and len(para.runs) != 1:  # 删除空行,等于1为图片
                    print("空行")
                    p = para._element
                    p.getparent().remove(p)
                    p._p = p._element = None
                else:
                    print("文本大于1")
            output_doc.save(r"test\%s.docx" % title)


def register_wrong_season(*args):
    """错题原因登记"""
    student=register()
    title=entry_input.get()
    if input_title_test(title)==True:
        wrong_season_text=combobox_wrong_season.get()
        print("title,wrong_season:",title,wrong_season_text)
        makedir(r"安装必备文件\学生\%s" % (student.classs))
        try:
            stu_wro_season = pd.DataFrame(pd.read_excel(
                r"安装必备文件\学生\%s\%s.xlsx" %(student.classs, student.name)))
        except Exception as e:
            print("打开学生错题记录失败:", e)
            stu_wro_season = pd.DataFrame(columns=["题号", "登记时间", "错误原因"])
        print("stu_wro_season:", stu_wro_season)
        new_row = pd.DataFrame([[title, time_now, wrong_season_text]],
                               columns=stu_wro_season.columns)
        stu_wro_season = pd.concat([stu_wro_season, new_row], ignore_index=True)
        stu_wro_season = stu_wro_season.set_index("题号")
        print("stu_wro_season:", stu_wro_season)
        makedir(r"安装必备文件\学生\%s"%student.classs)
        stu_wro_season.to_excel(r"安装必备文件\学生\%s\%s.xlsx" %(student.classs, student.name))
        input_title_save(title)
        entry_input.delete(0,tk.END)
        text=student.wrong_season_next(title)
        print("student.title_wrong:",student.wrong_season_next(title))
        title_show.insert(tk.INSERT,text)
    else:
        print("输入题号错误3")
        delete_title_show(title_show)
        tk.messagebox.showinfo(title="错误提示",
                message="输入题号%s题库系统中不存在，请确认是否输入正确，字母只能用大写"%title)


    # wrong_season(student.name)
    # print("student.title_all:",len(student.title_all),student.title_all)
    # student.title_all = student.title_all.iloc[1:]
    # title = student.title_all.iloc[0]
    # index = title.index("K")  # 寻找字符的索引位置
    # result = title[index + 1:]  # 使用切片截取字符后面的内容
    # wrong_title_text.set("题目：第%s题,题库编号:%s" % (title[index + 1:], title))
def wrong_title_next(wrong_sea):
    """错题原因登记 下一题"""
    student=register()
    com

    student.wrong_season()
    wrong_title_text.set("题目：第%d题,题库编号:%s" % (i + 1, title))
# def question_type_train():
#     """题型针对性训练，包含知识点的专门训练"""
#     student=register()
#     from tkinter import Toplevel
#     question_type_train_window = tk.Tk()
#     question_type_train_window.geometry("250x200")
#     question_type_train_window.title("题型针对性训练")
#     students = register()
#     question_type_all=pd.DataFrame(pd.read_excel(r"安装必备文件/题目分类方法2.xlsx"))
#     # print("question_type_all:",question_type_all)
#     question_type_all=question_type_all.loc[:,"知识点"]
#     # yscroll = tk.Scrollbar(question_type_train_window, orient=tk.VERTICAL, command=btn.yview)
#     yscroll = tk.Scrollbar(question_type_train_window, orient=tk.VERTICAL)
#     yscroll.grid(row=0, column=5, sticky='NS')
#
#     i=9
#     for question_type in question_type_all:
#         btn = tk.Checkbutton(question_type_train_window, text=question_type,
#                              command=lambda c=question_type: input_title_test2(c),)
#                              # command=lambda c=question_type: input_title_test2(c), yscrollcommand = yscroll.set)
#         print(i%4,i//4)
#         btn.grid(row=i//4, column=i%4,sticky="w")
#
#         i+=1
#         print("question_type:",question_type)
#     yscroll.config(command=btn.yview)
#
#
#     # print("students:",students.__dict__)
#     # label = tk.Label(register_window, text="你要输入的题目的章节是:%s%s"
#     #                                        % (combobox_chapter.get(), combobox_section.get()))
#     # 布局，按行、列布局，不考虑宽度
#     # label.grid(row=0, column=0)
#     # for chapter in range(15):
#     #     if combobox_chapter.get() == "第%d章" % chapter:
#     #         if chapter < 5:
#     #             title_pre = "202311" + str(chapter)  # 必修1第1章
#     #         elif chapter < 9:
#     #             title_pre = "202312" + str(chapter)  # 必修2第5章
#     #         elif chapter > 8:
#     #             title_pre = "202313" + str(chapter)  # 必修3第5章
#     #         break
#     # for section in range(7):
#     #     if combobox_section.get() == "第%d节" % section:
#     #         title_pre += str(section)  # +第1节
#     #         break
#     # for special_subject in range(20):
#     #     if combobox_section.get() == "专题%d" % special_subject:
#     #         title_pre += str("Z%d" % special_subject)  # +专题1
#     #         break
#     #
#     # print("title_pre:", title_pre)
#     # register_ti(title_pre)
#     # 如果一节内容有两个课时，则双排展示
#
#     print("--------------------你需要手动设置具有两个课时的章节列表------------------")
#     two_train_in_one_section = ["20231254", "20231262",  "20231281","2023127Z10",
#                                 "20231131", "20231134", "20231135", "20231112"
#         , "20231113", "20231114", "2023114Z14", "2023114Z15","20231145" ]
#     three_train_in_one_section=["20231134"]
#
#     if title_pre in two_train_in_one_section:
#         # label=tk.Label(register_window,text="第一课时")
#         for i in range(1, 15):
#             title = str(title_pre) + "1KK" + str(i)
#             # button带参数传递，响应函数
#             print("title：str(title_pre)+1KK+str(i):", title)
#             btn = tk.Checkbutton(register_window, text="第一课时（训练一）第%d题" % i,
#                                  command=lambda c=title: input_title_test2(c))
#             btn.deselect()
#             btn.grid(row=i, column=0)
#         for i in range(1, 15):
#             title = str(title_pre) + "2KK" + str(i)
#             # button带参数传递，响应函数
#             btn = tk.Checkbutton(register_window, text="第二课时（训练二）第%d题" % i,
#                                  command=lambda c=title: input_title_test2(c))
#             btn.deselect()
#             btn.grid(row=i, column=1)
#         if title_pre in three_train_in_one_section:
#             for i in range(1, 15):
#                 title = str(title_pre) + "3KK" + str(i)
#                 # button带参数传递，响应函数
#                 btn = tk.Checkbutton(register_window, text="第二课时（训练二）第%d题" % i,
#                                      command=lambda c=title: input_title_test2(c))
#                 btn.deselect()
#                 btn.grid(row=i, column=2)
#     else:
#         for i in range(1, 15):
#             title = str(title_pre) + "K" + str(i)
#             # button带参数传递，响应函数
#             btn = tk.Checkbutton(register_window, text="第%d题" % i,
#                                  command=lambda c=title: input_title_test2(c))
#             btn.deselect()
#             btn.grid(row=i, column=0)
#     save = tk.Button(register_window, text="保存", bg='#BEBEBE',
#                      command=lambda c=register_window: clear(c))
#     save.grid(row=i + 1, column=1)
#     text = "你已经登记的题号："
#     students = register()
#     for title in students.title_message.loc[:, students.name]:
#         if str(title) != "nan":
#             text += str(title)
#             text += "  "
#     print(text)
#     register_window.mainloop()

class Test_paper():
    """组卷窗口"""
    def __init__(self,students):
        self.students=students
        self.gui = tk.Toplevel()  # create gui window
        self.gui.title("组卷")  # set the title of gui
        self.gui.geometry("1300x800")  # set the window size of gui
        self.title_all={"单选题":0,"多选题":0,"实验题":0,"计算题":0,}
        self.button_regester_title=tk.Button(self.gui,text="将常用错题录入错题库",
                                             command=self.regester_title)
        self.button_regester_title.grid(row=0,column=1)
        self.label3 = tk.Label(self.gui, text="输入学生易错题目的章节代码")
        self.label3.grid(row=1,column=1)

        self.combobox_book = ttk.Combobox(self.gui)
        self.combobox_book.bind('<<ComboboxSelected>>')
        # combobox_class['values'] = ("12班")
        # combobox_class.current(0)
        self.combobox_book['values'] = ("2023级考试","步步高2023级必修1","步步高2023级必修2","步步高2023级必修3",
                               "步步高2023级选择必修1","步步高2023级选择必修2","步步高2023级选择必修3")  # 设置下拉框的选项列表
        self.combobox_book.current(4)
        self.combobox_book.bind('<<ComboboxSelected>>', self.book_select)

        self.combobox_chapter = ttk.Combobox(self.gui)
        self.combobox_chapter.bind('<<ComboboxSelected>>',self.chapter_select)
        # combobox_class['values'] = ("12班")
        # combobox_class.current(0)
        # self.combobox_chapter['values'] = ("第1章", "第2章", "第3章", "第4章",
        #                               "第5章", "第6章", "第7章", "第8章", "第9章", "第10章", "第11章", "第12章", "第13章")  # 设置下拉框的选项列表
        # self.combobox_chapter.current(0)

        # combobox_name = ttk.Combobox(win,postcommand=register)
        self.combobox_section = ttk.Combobox(self.gui)
        self.combobox_section['values'] = ("章末素养提升","第1节", "第2节", "第3节", "第4节", "第5节", "第6节",
                                      '专题1', '专题2', '专题3', '专题4', '专题5', '专题6', '专题7',
                                      '专题8', '专题9', '专题10', '专题11', '专题12', '专题13',
                                      '专题14', '专题15', '专题16', '专题17', '专题18', '专题19',
                                      )  # 设置下拉框的选项列表
        self.combobox_section.current(0)
        self.combobox_section.bind('<<ComboboxSelected>>', self.title_init)

        self.combobox_train = ttk.Combobox(self.gui)
        self.combobox_train['values'] = ("课时选择","第1课时","第2课时",
                                      )  # 设置下拉框的选项列表
        self.combobox_train.current(0)
        self.combobox_train.bind('<<ComboboxSelected>>', self.title_init)

        self.combobox_book.grid(row=2,column=1)
        self.combobox_chapter.grid(row=3,column=1)
        self.combobox_section.grid(row=4,column=1)
        self.combobox_train.grid(row=5,column=1)

        self.sv3 = tk.StringVar()
        # sv3.trace("w", self.f)
        self.entry_input3 = tk.Entry(self.gui, textvariable=self.sv3)
        self.entry_input3.bind("<Return>", self.title_init)  # 回车响应
        # entry_input3.bind('<Key>', self.update)
        self.entry_input3.grid(row=6,column=1)

        # title=students.title_message.loc[0, students.name]
        # self.label4 = tk.Label(self.gui, text="题目:%s" % title)
        # self.label4.grid(row=3,column=1)

        # image = Image.open(r"图片题目汇总/%s.jpg" % title)  # open image from path
        image = Image.open(r"图片题目汇总/202312811KK1.jpg")  # open image from path
        image = image.resize((600, 450))  # 调整图像大小
        img = ImageTk.PhotoImage(image)  # read opened image
        self.label1 = tk.Label(self.gui, image=img)  # create a label to insert this image
        self.label1.grid(row=7, column=1,rowspan=9)  # set the label in the main window

        self.button_next = tk.Button(self.gui, text="上一题", command=self.previous_title)
        self.button_next.grid(row=16, column=1)

        self.button_next = tk.Button(self.gui, text="下一题", command=self.next_title)
        self.button_next.grid(row=17, column=1)


        self.button_join = tk.Button(self.gui, text="加入试卷", command=self.title_join_test_paper)
        self.button_join.grid(row=18, column=1)

        # self.label5=tk.Label(self.gui,text="已经加入试卷题目")
        # self.label5.grid(row=1,column=2)
        #
        # self.frame2=tk.Frame(self.gui)
        # self.label6=tk.Label(self.frame2,text="666")
        # self.label6.grid(row=1,column=1)
        # self.frame2.grid(row=2,column=3,rowspan=1500)

        # #使滚轮可以控制全部试卷显示界面
        # def on_mousewheel(event):
        #     # 鼠标滚轮事件处理函数
        #     canvas.yview_scroll(int(-1 * (event.delta / 120)), "units")
        # frame = tk.Frame(self.gui)
        # frame.grid(row=0,columnspan=4,column=2,rowspan=20)
        # canvas = tk.Canvas(frame)
        # scrollbar = tk.Scrollbar(frame, orient=tk.VERTICAL, command=canvas.yview)
        #
        # canvas.create_window((0, 0), window=text, anchor=tk.NW)
        # # 配置canvas和scrollbar之间的关联
        # canvas.configure(yscrollcommand=scrollbar.set)
        # # 使用bindtag将鼠标滚轮事件绑定到canvas上
        # canvas.bind("<MouseWheel>", on_mousewheel)
        # # 设置canvas的滚动区域
        # canvas.configure(scrollregion=canvas.bbox("all"))
        # # 放置小部件
        # scrollbar.grid(row=1,column=6,rowspan=19)
        # canvas.grid(row=1,column=2,columnspan=4,rowspan=19)

        self.canvas = tk.Canvas(self.gui)  # 创建画布
        # canvas.place(x=0, y=0, height=300, width=500)
        self.canvas.grid(row=0,column=3,rowspan=20)

        self.myscrollbar = tk.Scrollbar(self.gui, orient="vertical", command=self.canvas.yview)  # 创建滚动条
        # myscrollbar.place(x=500, y=0, height=300)
        self.myscrollbar.grid(row=0,column=4,rowspan=20,sticky = "ns")
        self.canvas.configure(yscrollcommand=self.myscrollbar.set)

        self.rollFrame = tk.Frame(self.canvas)  # 在画布上创建frame
        self.canvas.create_window((0, 0), window=self.rollFrame, anchor='nw')  # 要用create_window才能跟随画布滚动



        def myfunction(event):
            self.canvas.configure(scrollregion=self.canvas.bbox("all"), width=600, height=600)
        self.rollFrame.bind("<Configure>", myfunction)
        tk.Label(self.rollFrame,text="已经添加的试题").grid(row=0,column=2)
        # tk.Label(self.rollFrame,text="已经添加的试题").grid(row=0,column=2)

        # def data():
        #     for i in range(100):
        #         tk.Label(rollFrame, text=i).grid(row=i, column=0)
        #         tk.Label(rollFrame, text="my text" + str(i)).grid(row=i, column=1)
        #         tk.Label(rollFrame, text="..........").grid(row=i, column=2)
        # data()


        # self.gui.attributes('-fullscreen', True)#全屏

        # self.title_all_number = 1
        # print("title:::", title, type(title), type(str(title)))
        # if str(title) != "nan":
        #     # self.gui.refresh(title)
        #
        #     self.label3 = tk.Label(self.gui, text="设置本题%s的作答时间：(分钟)"%title)
        #     # self.label3.grid(row=1,column=1)
        #     self.sv3 = tk.StringVar()
        #     # sv3.trace("w", self.f)
        #     self.entry_input3 = tk.Entry(self.gui, textvariable=self.sv3)
        #     self.entry_input3.bind("<Return>", self.refresh)  # 回车响应
        #     # entry_input3.bind('<Key>', self.update)
        #     # self.entry_input3.grid(row=2,column=1)
        #
        #     # button=tk.Button(self.gui,text="下一题",command=Gui)
        #     # button.pack()
        # self.gui.state("zoomed")#最大化
        self.book_select()
        self.gui.mainloop()  # start mainloop

    def chapter_select(self,*args):
        """选择书籍，相应的刷新章节"""
        title_data = pd.DataFrame(pd.read_excel(r"word题目汇总/题目清单.xlsx")).set_index("章节")
        if "考试" in self.combobox_book.get():
            # print("combobox_book.get():", combobox_book.get())
            year = int(self.combobox_chapter.get()[0:4])
            month = int(self.combobox_chapter.get()[5:7])
            my_date = datetime.date(year, month, 1)  # 2023年10月1日

            # 将日期转换为整数，表示自1970年1月1日以来的总天数
            print("year,month:", year, month)
            date_as_int = (my_date - datetime.date(1970, 1, 1)).days + 25569  # 计算天数差
            print("date_as_int:", date_as_int)
            title_data = title_data.loc[date_as_int, "节"]
            chapter_sections = []
            for chapter_section in title_data:
                print("book_chapter:", chapter_section)
                try:
                    chapter_section = str(int(chapter_section))
                except:
                    pass
                if chapter_section not in chapter_sections:
                    chapter_sections.append(chapter_section)
            chapter_sections2 = []
            for chapter_section in chapter_sections:
                if chapter_section not in chapter_sections2:
                    chapter_sections2.append(chapter_section)
            # print("title_data", title_data)
            print("book_chapters2:", chapter_sections2, chapter_sections)
        else:
            title_data = title_data.loc[int(self.combobox_chapter.get()[1:-1]), "节"]
            chapter_sections = []
            for chapter_section in title_data:
                print("book_chapter_步步高:", chapter_section)
                try:
                    chapter_section = str(int(chapter_section))
                except:
                    pass
                if chapter_section not in chapter_sections:
                    chapter_sections.append(chapter_section)
            print("chapter_sections:", chapter_sections)
            chapter_sections2 = []
            for chapter_section in chapter_sections:
                try:
                    chapter_sections2.append("第" + str(int(chapter_section)) + "节")
                except:
                    chapter_sections2.append(chapter_section)
            # print("title_data", title_data)
            print("book_chapters3:", chapter_sections2)
        chapter_sections2.sort()
        self.combobox_section.config(values=chapter_sections2)  # 重新排序
        self.combobox_section.current(0)
        self.title_init()

    def book_select(self,*args):
        """选择书籍，相应的刷新章节"""
        title_data = pd.DataFrame(pd.read_excel(r"word题目汇总/题目清单.xlsx")).set_index("录入来源")
        print("combobox_book.get():", self.combobox_book.get())
        title_data = title_data.loc[self.combobox_book.get(), "章节"]
        book_chapters = []
        for book_chapter in title_data:
            print("book_chapter:", book_chapter)
            try:
                book_chapter = str(int(book_chapter))
            except:
                pass
            if book_chapter not in book_chapters:
                book_chapters.append(book_chapter)
        book_chapters_2 = []
        for book_chapter in book_chapters:
            if "考试" not in self.combobox_book.get():
                book_chapters_2.append("第" + book_chapter + "章")
            else:
                print("2345")
                # delta = pd.Timedelta(str(int(book_chapter)) + 'days')
                # time = pd.to_datetime('1899-12-30') + delta
                # 将数字 41957 转换为 Python 的日期时间对象
                import datetime
                python_datetime = datetime.datetime(1899, 12, 30) + datetime.timedelta(days=int(book_chapter))
                # time.strftime('%Y{y}%m{m}%d{d}').format(y='年', m='月', d='日')
                python_datetime = python_datetime.strftime('%Y{y}%m{m}').format(y='年', m='月')
                book_chapters_2.append(python_datetime)
        # print("title_data", title_data)
        print("book_chapters:", book_chapters_2)
        self.combobox_chapter.config(values=book_chapters_2)
        self.combobox_chapter.current(0)
        self.chapter_select()
    def next_title(self,*args):
        """组卷状态下，输入下一题"""
        self.next_title_number+=1
        self.refresh()
    def previous_title(self,*args):
        """组卷状态下，输入上一题"""
        self.next_title_number-=1
        self.refresh()
    def update(self,title,*args):
        #更新题目图片
        print("title,*args:",title,*args)
        # self.label1.destroy()#删除原图片
        self.label3.configure(text="设置本题%s的作答时间：(秒钟)"%title)
        len_entry=len(self.entry_input3.get())
        self.entry_input3.delete(0,len_entry)
        # self.label3 = tk.Label(self.gui, text="请输入本题的作答时间（分钟）:")
        # self.label3.pack()
        # self.sv3 = tk.StringVar()
        # # sv3.trace("w", self.f)
        # self.entry_input3 = tk.Entry(self.gui, textvariable=self.sv3)
        # self.entry_input3.bind("<Return>", self.refresh)  # 回车响应
        # # entry_input3.bind('<Key>', self.update)
        # self.entry_input3.pack()
        # switch_time = self.entry_input3.get()
        # self.label3 = tk.Label(self.gui, text="请输入本题的作答时间（分钟）:")
        # self.label3.pack()
        # self.sv3 = tk.StringVar()
        # # sv3.trace("w", self.f)
        # self.entry_input3 = tk.Entry(self.gui, textvariable=self.sv3)
        # self.entry_input3.bind("<Return>", self.refresh)  # 回车响应
        # # entry_input3.bind('<Key>', self.update)
        # self.entry_input3.pack()
        # switch_time = self.entry_input3.get()
        # self.title_all_number += 1
        # button=tk.Button(self.gui,text="下一题",command=Gui)
        # button.pack()

        try:
            image = Image.open(r"cach/pic/%s.jpg" % title)  # open image from path
        except Exception as e:
            print("插入图片失败：", e.__traceback__.tb_lineno, e)
            image = Image.open(r"图片题目汇总/%s.jpg" % title)  # open image from path
        image = image.resize((600,450))  # 调整图像大小
        img = ImageTk.PhotoImage(image)  # read opened image
        # new_image = Image.open("new_image.jpg")
        # new_tk_image = ImageTk.PhotoImage(new_image)
        # label.configure(image=new_tk_image)
        self.label1.config(image=img)#更新
        self.label1.image=img#图片更新要有这一行
        # tk.update()
        print("图片更新成功")
        # self.label1 = tk.Label(self.gui, image=img)  # create a label to insert this image
        # self.label1.pack()  # set the label in the main window

        # try:
        #     for element in self.prs.slides[self.title_all_number]._element[2][0]:
        #         element.set('advTm', "%s" % (switch_time * 1000 * 60))  # 设置幻灯片切换时间
        #         switch_time=0
        #         print("设置时间成功：：", switch_time)
        # except Exception as e:
        #     print("设置时间不成功：", e)

    def title_init(self,*args):
        """输入题号初始化"""
        self.next_title_number=0
        self.refresh()

    def refresh(self,*args):
        #识别章节代码


        if self.entry_input3.get():
            title_pre = self.entry_input3.get()
        else:
            title_data = pd.DataFrame(pd.read_excel(r"word题目汇总/题目清单.xlsx"))
            if "考试" in self.combobox_book.get():
                print("combobox_book.get()2:", self.combobox_book.get())
                year = int(self.combobox_chapter.get()[0:4])
                month = int(self.combobox_chapter.get()[5:7])
                kaoshitime = str(year) + "年"
                my_date = datetime.date(year, month, 1)  # 2023年10月1日

                # 将日期转换为整数，表示自1970年1月1日以来的总天数
                date_as_int = (my_date - datetime.date(1970, 1, 1)).days + 45658 - 20089  # 计算天数差

                # title_data = title_data.set_index("章节")
                title_data = title_data[(title_data["录入来源"] == self.combobox_book.get()) &
                                        (title_data["章节"] == date_as_int) & (
                                                    title_data["节"] == self.combobox_section.get())]

                print("title_data:",title_data)
                title_data = title_data.set_index("节")
                title_pre=title_data["代号"][0]
                print("title_pre:",title_pre)
            else:
                chapter = self.combobox_chapter.get()
                section = self.combobox_section.get()
                print("chapter,section:", chapter, section)
                if "第" in chapter and "章" in chapter:
                    chapter = int(chapter[1:len(chapter) - 1])
                if "第" in section and "节" in section:
                    section = str(section[1: -1])
                print("chapter,section:", chapter, section)
                title_data = title_data[(title_data["录入来源"] == self.combobox_book.get()) &
                                        (title_data["章节"] == chapter) & (title_data["节"] == section)]

                print("title_data:",title_data)
                title_data = title_data.set_index("节")
                title_pre=title_data["代号"][0]

        title_data=title_data.set_index("代号")
        print("title_data:",title_data)
        print(title_pre)
        self.title_data=title_data.loc[title_pre,"题号"]
        # self.title_data=title_data.loc["202312812K","题号"]
        # print("self.title_data:",self.title_data)
        try:
            title=title_pre
            image = Image.open(r"图片题目汇总/%s.jpg" % title)  # open image from path
            print("打开图片：",title)

        except Exception as e:
            print("输入的不是题号：",e.__traceback__.tb_lineno,e)
            title=self.title_data.iloc[self.next_title_number]


            print("title,number:",title,self.next_title_number)
            image = Image.open(r"图片题目汇总/%s.jpg" % title)  # open image from path
        image = image.resize((600,450))  # 调整图像大小
        img = ImageTk.PhotoImage(image)  # read opened image
        # new_image = Image.open("new_image.jpg")
        # new_tk_image = ImageTk.PhotoImage(new_image)
        # label.configure(image=new_tk_image)
        self.label1.config(image=img)#更新
        self.label1.image=img#图片更新要有这一行
        self.title = title

        # title=
        # print("title_data:",title_data)
        #
        # try:
        #     for element in self.prs.slides[self.title_all_number]._element[2][0]:
        #         element.set('advTm', "%s" % (switch_time))  # 设置幻灯片切换时间,一般为几分钟
        #         print("设置时间成功", switch_time)
        # except Exception as e:
        #     print("设置时间不成功：", e)
        # try:
        #     title=self.students.title_message.loc[self.title_all_number, self.students.name]
        #     # title=self.students.title_message.loc[0, self.students.name]
        #     print("title:",title)
        #     self.update(title)
        #     self.title_all_number += 1
        # except Exception as e:
        #     self.gui.destroy()
        #
        #     print("图片更新不成功：",e)
        #
        #     # 插入答案
        #     print("#插入答案")
        #     title_all_number_answer = 1
        #     switch_time = 3000  # 答案切换时间5秒
        #     for title in self.students.title_message.loc[:, self.students.name]:
        #         # print("title:::",title,type(title),type(str(title)))
        #         if str(title) != "nan":
        #
        #             title = "A" + title
        #             if title[4] == "1":
        #                 bixiu = "必修%s" % title[5]
        #             else:
        #                 bixiu = "选修%s" % title[5]
        #
        #             if "Z" in title:
        #                 if "KK" in title:
        #                     chapter = title[7:title.find("Z")]
        #                     # print(chapter)
        #                     section = title[title.find("Z") + 1:title.find("KK") - 1]
        #                     class_hour = title[title.find("KK") - 1]
        #                     title_number = title[title.find("KK") + 2:]
        #                     title_from = "%s第%s章专题%s第%s课时第%s题" % (
        #                         bixiu, chapter, section, class_hour, title_number)
        #                 else:
        #                     chapter = title[7:title.find("Z")]
        #                     # print(chapter)
        #                     section = title[title.find("Z") + 1:title.find("K")]
        #                     title_number = title[title.find("K") + 1:]
        #                     title_from = "%s第%s章专题%s第%s题" % (bixiu, chapter, section, title_number)
        #
        #             else:
        #                 if "KK" in title:
        #                     chapter = title[7:title.find("KK") - 2]
        #                     # print(chapter)
        #                     section = title[title.find("KK") - 2:title.find("KK") - 1]
        #                     class_hour = title[title.find("KK") - 1]
        #                     title_number = title[title.find("KK") + 2:]
        #                     title_from = "%s第%s章第%s节第%s课时第%s题" % (
        #                         bixiu, chapter, section, class_hour, title_number)
        #                 else:
        #                     chapter = title[7:title.find("K") - 1]
        #                     # print(chapter)
        #                     section = title[title.find("K") - 1:title.find("K")]
        #                     title_number = title[title.find("K") + 1:]
        #                     title_from = "%s第%s章第%s节第%s题" % (bixiu, chapter, section, title_number)
        #             print(title_from, title)
        #             # blank_slide_layout = self.prs.slide_layouts[6]
        #             # slide = prs.slides.add_slide(blank_slide_layout)
        #             slide = self.prs.slides[title_all_number_answer+self.title_all_number]
        #
        #             left = Inches(1)
        #             top = Inches(1)
        #             width = Inches(8)
        #             height = Inches(6)
        #
        #
        #
        #             try:
        #                 docx_to_picture(title)
        #                 slide.shapes.add_picture(r"cach/pic/%s.jpg" % title, left, top,width,height)
        #             except Exception as e:
        #                 print("插入图片失败：", e.__traceback__.tb_lineno, e)
        #                 slide.shapes.add_picture(r"图片题目汇总/%s.jpg" % title, left, top,width,height)
        #
        #             left = Inches(1)
        #             top = Inches(0)
        #             width = Inches(8)
        #             height = Inches(2)
        #
        #             # textbox = slide.shapes.add_textbox(left, top, width, height)
        #             # tf = textbox.text_frame
        #             # left = top = Inches(1)
        #             # width = Inches(6)
        #             # height = Inches(2)
        #
        #             textbox = slide.shapes.add_textbox(left, top, width, height)
        #             tf = textbox.text_frame
        #             # tf.text = "This is text inside a textbox"
        #
        #             p = tf.add_paragraph()
        #             p.text = "第%d题答案" % (title_all_number_answer)
        #             p.font.bold = True
        #             p.font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)
        #             p = tf.add_paragraph()
        #             p.text = "题目来源:%s/题库编号:%s" % (title_from, title)
        #
        #             try:
        #                 for element in self.prs.slides[title_all_number_answer+
        #                                                self.title_all_number]._element[2][0]:
        #                     element.set('advTm', "%s" % (switch_time))  # 设置幻灯片切换时间
        #             except Exception as e:
        #                 print("设置时间不成功2：",e,e.__traceback__.tb_lineno)
        #             title_all_number_answer += 1
        #
        #     len_prs_slides = len(self.prs.slides)
        #     for number in range(len_prs_slides, title_all_number_answer+self.title_all_number, -1):
        #         del self.prs.slides._sldIdLst[-1]
        #         # prs._sldIdLst.remove(prs.slides[number])
        #
        #     self.prs.save(r"学生错题word文档\错题重做%s.pptx" % time_now)
        #     os.startfile(r"学生错题word文档\错题重做%s.pptx" % time_now)
        #
        #     return

        # self.gui.destroy()
        # self.update(title)
    def regester_title(self,):
        """将学生易错题目录入题库中"""
        titles=pd.DataFrame(pd.read_excel(r"word题目汇总/题目清单.xlsx"))
    def title_join_test_paper(self):
        """将当前题目加入试卷中"""
        title_have=len(self.title_all)-4
        try:
            self.title
            # print(self.title)
        except:
            self.title="20231271K1"
            self.title_all[1] = {"题号": self.title, "题型": "选择题"}
        title_data=pd.DataFrame(pd.read_excel(r"word题目汇总/题目清单.xlsx")).set_index("题号")
        title_data=title_data.loc[self.title,"题型"]
        try:
            title_type=str(title_data.iloc[-1])
        except:
            title_type=title_data
        print("加入题目：",title_have+1,self.title,title_type)
        self.title_all[title_have+1]={"题号":self.title,"题型":title_type}
        print("self.title_all1:",self.title_all)
        self.title_all[title_type] =self.title_all[title_type]+1
        print("self.title_all2:",self.title_all)
        self.test_paper_show()
    def edit_type(self,number):
        """修改题目的题型，并加入到excel表格中"""
        title=self.title_all[number]["题号"]
        title_data=pd.DataFrame(pd.read_excel(r"word题目汇总/题目清单.xlsx")).set_index("题号")
        print("title_data1:",title_data)
        print("title4:",title)
        self.title_data=title_data.loc[title,"题型"]
        print("self.title_data22:",self.title_data)
    def test_paper_show(self):
        """显示试卷的所有题目"""
        # titlenumbers=len(self.title_all)
        time_all=0
        select_1_type=0
        select_2_type=0
        experiment_type=0
        cal_type=0
        self.title_type={}
        title_all_middle={}
        # title_data = pd.DataFrame(pd.read_excel(r"word题目汇总/题目清单.xlsx")).set_index("题号")
        print("self.title_all22:",self.title_all)
        for number in self.title_all:
            # print("number:",number,self.title_all[number])
            # print("title_data:", title_data)
            title=self.title_all[number]
            # title_data2=title_data.loc[title,"题型"]
            # title_type=str(title_data2.iloc[-1])
            # print("title_data:",title_data,)
            print("title33:",title)
            if type(number)==int:
                title_type=title["题型"]
                if title_type == "单选题":
                    time_all += 2.5
                    select_1_type += 1
                    title_all_middle[select_1_type] = title
                elif title_type == "多选题":
                    time_all += 2.5
                    select_2_type += 1
                    title_all_middle[select_2_type + self.title_all["单选题"]] = title
                elif title_type == "实验题":
                    time_all += 6
                    experiment_type += 1
                    title_all_middle[self.title_all["单选题"]+self.title_all["多选题"]+ experiment_type] = title
                elif title_type == "计算题":
                    time_all += 8
                    cal_type += 1
                    title_all_middle[self.title_all["单选题"]+self.title_all["多选题"]+
                                     self.title_all["实验题"]+ cal_type] = title
                else:
                    time_all += 100
            else:
                title_all_middle[number]=self.title_all[number]

        self.title_all=title_all_middle
        print("self.title_all3:",self.title_all)
        self.label8=tk.Label(self.rollFrame,text="已经录入单选题%s个,多选题%s个，实验题%s个计算题%s个，预计用时%s分钟"
                                %(self.title_all["单选题"],self.title_all["多选题"],self.title_all["实验题"],self.title_all["计算题"],time_all))
        self.label8.grid(row=1,column=1,columnspan=4)

        self.button7 = tk.Button(self.rollFrame, text="生成试卷",
                                 command=self.generate_test_paper)
        self.button7.grid(row=1, column=5, )
        for number in self.title_all:
            if type(self.title_all[number])!=int:
                print("self.title_all[number]:",number,self.title_all[number])
                self.label=tk.Label(self.rollFrame,text="%s、%s%s"%(number,
                    self.title_all[number]["题号"],self.title_all[number]["题型"]))
                self.label.grid(row=2+number*2,column=1)

                self.button6 = tk.Button(self.rollFrame, text="向上移动",
                                         command=lambda up=number:self.change_up(up))
                self.button6.grid(row=2+number*2, column=2,)

                self.button6 = tk.Button(self.rollFrame, text="向下移动",
                                         command=lambda down=number:self.change_down(down))
                self.button6.grid(row=2 + number*2, column=3, )

                self.button6 = tk.Button(self.rollFrame, text="移到最后",
                                         command=lambda down=number:self.change_last(down))
                self.button6.grid(row=2 + number*2, column=4, )
                self.button6 = tk.Button(self.rollFrame, text="删除",
                                         command=lambda down=number:self.delete(down))
                self.button6.grid(row=2 + number*2, column=5, )

                self.button6 = tk.Button(self.rollFrame, text="修改题型",
                                         command=lambda down=number:self.edit_type(down))
                self.button6.grid(row=2 + number*2, column=6, )
                # print("self.title_all[number]:",self.title_all[number])

                image = Image.open(r"图片题目汇总/%s.jpg"%self.title_all[number]["题号"])  # open image from path
                # image = Image.open(r"图片题目汇总/20231283K1.jpg")  # open image from path
                image = image.resize((600,450))  # 调整图像大小
                img = ImageTk.PhotoImage(image)  # read opened image
                print("img:",img)
                self.label7 = tk.Label(self.rollFrame, image=img)  # create a label to insert this image
                self.label7.image = img
                # self.label7.subsample(600,250)
                self.label7.grid(row=2+number*2+1, column=1,columnspan=6)  # set the label in the main window
    def generate_test_paper(self):
        """生成试卷"""
        print("self.students.name:",self.students)
        print("self.students.name:",self.students.name)
        students_list = pd.DataFrame(pd.read_excel(r"sys_files/学生名单.xlsx"))
        print("students_list,self.students.classs:",students_list,self.students.classs)
        students_list = students_list.set_index("班级")
        students_list = students_list.loc[self.students.classs, "姓名"]
        self.names = []
        if type(students_list) == str:
            self.names.append(students_list)
        else:
            for i in students_list:
                self.names.append(i)

        self.test_paper=pd.DataFrame(columns=self.names)
        for number in range(len(self.title_all)-4):
            if type(number)==int:
                print("number+1,self.title_all[number]:",number+1,self.title_all[number+1])
                self.test_paper.loc[number,self.students.name]=str(self.title_all[number+1]["题号"])
        # print("self.test_paper:",self.test_paper)
        self.test_paper=self.test_paper.set_index(self.students.name)
        # self.test_paper.to_excel(r"%s\组卷\%s2023级12班错题.xlsx" %(
        #         student_wrong_path,time_now2,))
        # makedir(r"%s\组卷" % (            student_wrong_path ))
        # self.test_paper.to_excel(r"%s\%s2023级12班错题.xlsx" % (
        #     student_wrong_path, time_now2,))
        try:
            self.test_paper.to_excel(r"%s\组卷\%s2023级12班错题.xlsx" % (
                   student_wrong_path, time_now2,))
        except:
            makedir(r"%s\组卷" % (student_wrong_path))
            self.test_paper.to_excel(r"%s\组卷\%s2023级12班错题.xlsx" % (
                   student_wrong_path, time_now2,))
        title_type=self.title_type

        def class_title_pra_zujuan(student_wrongs, classs, time_now, test_paper, title_wrong_path,
                                   title_wrong_pra_path):
            """#组卷中的错题题目和训练题插入word文档"""
            print("test_paper:", test_paper.__dict__)
            # print("student_wrongs:",student_wrongs)
            # print("student_wrongs['team']",student_wrongs['team'])
            team = student_wrongs['team']
            document = Document()
            # document = Document(r"安装必备文件/sample.docx")
            # section = document.sections[0]
            # # * 设置栏目数为2
            # sectPr = section._sectPr
            # cols = sectPr.xpath('./w:cols')[0]
            # cols.set(qn('w:num'), '2')
            # document.add_heading("%s 的“题不二错”\n时间：%s" % (classs,time_now),1)
            # document = class_title_pra(document, student_wrongs)
            titlewrongs = []

            name_number = 1
            len_student_wrongs = len(student_wrongs)
            # print('student_wrongs',student_wrongs)
            for name in student_wrongs:  # 将学生的错题导入，寻找相应的训练题
                progressbarOne['value'] += 1
                label_2.set("%s:将学生--%s题目添加到word文档，%s/%s" % (team, name, name_number, len(student_wrongs)))
                win.update()
                title_answer_types = {}
                title_option_changes = {}

                print("第%s个学生，共%s个" % (name_number, len_student_wrongs), name)
                student_practices = {}
                student_practices[name] = []
                # if name!='team' and name!='unit' and student_wrongs[name]!=[] and finish_name_all_dict[name][-1]==1:
                if name != 'team' and name != 'unit' and student_wrongs[name] != []:  # 不需要筛选功能，上一行为添加筛选功能
                    # 去除team和unit两个姓名和没有登记错题的学生,以及验证其上一次打印题目是否完成
                    from docx.enum.section import WD_SECTION
                    print("name_number:", name_number)

                    if name_number == 2:
                        new_section = document.sections[0]

                    else:

                        # new_section = document.add_section(WD_SECTION.ODD_PAGE)  # 增加分页，将每个学生隔开
                        new_section = document.add_page_break()  # 增加分页  # 增加分页，将每个学生隔开
                        new_section = document.add_page_break()  # 增加分页  # 增加分页，将每个学生隔开
                        # print("分页成功")
                    print("new_section:", new_section)
                    # if len(student_wrongs)==2:
                    #     document.add_heading("%s" % notice4, 4)
                    # else:
                    # document.add_heading("%s 的“题不二错”%s" % (name,time_now2), 3)
                    if name != "赵安宁":
                        document.add_paragraph("%s 的“错题重做”%s" % (name, time_now2))
                    else:
                        document.add_paragraph("%s的错题重做%s得分_______" % (name, time_now2))
                        title_all = test_paper.title_all
                        # 自动计算分值，并附加到word中
                        # title_all["单选题"]=18
                        # title_all["多选题"]=8
                        # title_all["实验题"]=0
                        # title_all["计算题"]=2
                        single_source = title_all["单选题"] * 3
                        mulitiple_source = title_all["多选题"] * 4
                        exper_source = title_all["实验题"] * 8
                        if title_all["计算题"] > 0:
                            cal_soruce = (100 - single_source - mulitiple_source - exper_source) // title_all["计算题"]
                            print("cal_soruce:", cal_soruce)
                            if cal_soruce > 8:
                                # 将多余的分值分配到计算题
                                excessive = (100 - single_source - mulitiple_source - exper_source) % title_all[
                                    "计算题"]
                                cal_soruce_text = []
                                for number in range(title_all["计算题"]):
                                    if excessive > 0:
                                        cal_soruce_text.append(str(cal_soruce + 1))
                                        excessive -= 1
                                    else:
                                        cal_soruce_text.append(str(cal_soruce))
                                cal_soruce_text.sort()
                                cal_soruce_text = "、".join(cal_soruce_text)
                                print("cal_soruce_text:", cal_soruce_text)
                            else:
                                mulitiple_source = title_all["多选题"] * 3
                                # cal_soruce = (100 - single_source - mulitiple_source - exper_source) // title_all["计算题"]
                                if title_all["计算题"] > 0:
                                    cal_soruce = (100 - single_source - mulitiple_source - exper_source) // title_all[
                                        "计算题"]
                                    if cal_soruce > 8:
                                        # 将多余的分值分配到计算题
                                        excessive = (100 - single_source - mulitiple_source - exper_source) % title_all[
                                            "计算题"]
                                        cal_soruce_text = []
                                        for number in range(title_all["计算题"]):
                                            if excessive > 0:
                                                cal_soruce_text.append(str(cal_soruce + 1))
                                                excessive -= 1
                                            else:
                                                cal_soruce_text.append(str(cal_soruce))
                                        cal_soruce_text.sort()
                                        cal_soruce_text = "、".join(cal_soruce_text)
                        else:
                            cal_soruce_text = "0"
                        print("cal_soruce_text2:", cal_soruce_text)
                        document.add_paragraph("说明：考试时间75分钟，单选题%s个,每个3分，多选题%s个，每个4分，实验题%s个，每个8分，"
                                               "计算题%s个，分值为%s" % (
                                                   title_all["单选题"], title_all["多选题"], title_all["实验题"],
                                                   title_all["计算题"],
                                                   cal_soruce_text))
                    # document.add_heading("%s" % author, 4)
                    number = 0  # 记录该学生的题目打印进度
                    pra_number = 1
                    not_exit_number = 0
                    # wrong_number=0
                    sum_number = len(student_wrongs[name])
                    print(" 将%s的题目（含训练题）输出到word文档中" % name, end=" ")
                    for titlepras in student_wrongs[name]:  # 将每一个学生的题目（含训练题）输出到word文档中
                        # print("1")
                        # print("student_practices[name]",student_wrongs[name])
                        for titlepra in titlepras:
                            # print("2")
                            # print("titlepras",titlepras)
                            if titlepras[titlepra] != []:
                                # print("3")
                                new_section.left_margin = Cm(1.27)
                                new_section.right_margin = Cm(1.27)
                                new_section.top_margin = Cm(1.0)
                                new_section.bottom_margin = Cm(1.0)
                                # print("titlepras[0]:",titlepras)
                                # print("titlepra:",titlepra)
                                # print("titlepras[titlepra]:",titlepras[titlepra])
                                # print("titlepras[‘错题’]:",titlepras["错题"])
                                title_type = findobj(titlepras[titlepra][0], '对象')
                                if titlepra == "错题" and title_type != "类型暂不清楚":  # 如果是错题，则打印题目类型
                                    document.add_heading('%s' % (findobj(titlepras[titlepra][0], '对象')), 2)
                                    # print("title_type:",title_type)
                                for title in titlepras[titlepra]:
                                    # print("4")
                                    new_section.left_margin = Cm(1.27)
                                    new_section.right_margin = Cm(1.27)
                                    new_section.top_margin = Cm(1.0)
                                    new_section.bottom_margin = Cm(1.0)
                                    # heading=document.add_heading('（%d）题库编号：%s 难度：%s' % (pra_number, m,finddif(m)), 3)
                                    # heading=document.add_heading('%d、题库编号：%s' % (pra_number, title), 3)
                                    heading = document.add_paragraph('%d、题库编号：%s' % (pra_number, title))
                                    heading.paragraph_format.line_spacing = 1
                                    heading.paragraph_format.space_before = Pt(0)
                                    heading.paragraph_format.space_after = Pt(0)
                                    # if '错题' in titlepra:
                                    #     # print("5")
                                    #     wrong_number += 1
                                    #     # if wrong_number%3==0:
                                    #     #     print("第%s题"%wrong_number,end=" ")
                                    #     # document.add_heading('%d、题目类别：%s 题号：%s' % (wrong_number,titlepra,m),2)
                                    #     document.add_heading('%d、%s %s' % (wrong_number,m,findobj(m,'对象')),2)
                                    # elif '训练题' in titlepra:
                                    #     # document.add_heading('(%d)%s 题号：%s' % (pra_number,titlepra,m),findobj(m,'对象'),3)
                                    #     document.add_heading('(%d)%s,%s,难度:%s' % (pra_number,m,findobj(m,'对象'),finddif(m)),3)
                                    pra_number += 1

                                    title_option_changes = wordtitle(document, title, title_option_changes)

                    pra_number = 1
                    # wrong_number = 0
                    # print("")
                    # def word_title_answer(student_practices,name):
                    print(" 将%s的题目答案（含训练题）输出到word文档中" % name, end="")
                    # if "模拟" in file_name:
                    #     break

                    try:
                        os.remove(r"cach/excel/select_answer.xlsx")
                        print("name,titlepras:", name, student_wrongs[name])
                    except Exception as e:
                        print("删除文件select answer失败:", e.__traceback__.tb_lineno, e)
                        print("name,titlepras2:", name, student_wrongs[name])
                    for i in range(len(student_wrongs[name])):
                        # for titlepras in student_wrongs[name]:  # 将每一个学生的题目（含训练题）的答案输出到word文档中
                        print("开始输入题目的答案:", name, titlepras, i, )
                        titlepras = student_wrongs[name][i]

                        for titlepra in titlepras:
                            if titlepras[titlepra] != []:
                                # 统计题目答案的类型选项更改的word,未更改的word,试卷的选择题，试卷的图片

                                def wordtitle_answer_find(title):

                                    # print("title_option_changes[title]:", title[1:], title_option_changes)
                                    try:
                                        title2 = title[1:]

                                        # print("title2,title_option_changes:", title2, title_option_changes)
                                        # print(type(title_option_changes[title2]), title_option_changes[title2])
                                        if title_option_changes[title2] != None:
                                            # def find_change_option_answer(title):
                                            #     #查找选择题目原来的答案
                                            #     document_answer = Document(r"word题目汇总\%s.docx" % title)
                                            #     for paragraph in document_answer.paragraphs:
                                            #         answer_options = re.findall(r'([A-Z]+)(正确|错误)', paragraph.text)
                                            #         answer_old = []
                                            #         for answer_option, correctness in answer_options:
                                            #             if correctness == '正确':
                                            #                 answer_old .append( answer_option)
                                            #     print("选择题的答案是：",answer_old,title)
                                            #     return answer_old
                                            # answer_old_change=find_change_option_answer(title)

                                            # 查找选择题目的正确答案
                                            # 更改答案
                                            # print("选项已经更改")
                                            def find_change_option_answer(title):
                                                document_answer = Document(r"word题目汇总\%s.docx" % title)
                                                options = ["A", "B", "C", "D", "E"]
                                                # print("answer:", title_option_changes)
                                                title_answer_option = []  # 题目答案选项，通过答案AB直接得到
                                                for paragraph in document_answer.paragraphs:
                                                    # print(paragraph.text, paragraph.text != None)
                                                    # print(paragraph.runs[0].text)
                                                    if "答案" in paragraph.text and paragraph.runs[0].text[
                                                                                    0:2] == "答案":
                                                        for option in options:
                                                            # 查找选择题的答案，未成功
                                                            if option in paragraph.runs[1].text[0:4]:
                                                                title_answer_option.append(option)
                                                            text2 = paragraph.runs[2].text
                                                            # print("title_answer_option:", option, title_answer_option, text2)
                                                            if (paragraph.runs[2].text[0:4] != "解析") and len(
                                                                    text2) < 5:
                                                                text3 = paragraph.runs[2].text
                                                                if (option in paragraph.runs[2].text[0:4]) and len(
                                                                        text3) < 5:
                                                                    title_answer_option.append(option)
                                                                    text4 = paragraph.runs[3].text
                                                                    if (paragraph.runs[3].text[0:4] != "解析") and len(
                                                                            text4) < 5:
                                                                        if option in paragraph.runs[3].text[0:4]:
                                                                            title_answer_option.append(option)
                                                # print("title_answer_option2:", title_answer_option)

                                                return title_answer_option

                                            title_answer_option = find_change_option_answer(title)
                                            answers = ""

                                            optionss = ""
                                            options = ["A", "B", "C", "D", "E", ]

                                            for i in title_option_changes[title2]:
                                                optionss += options[i]
                                            # print("options:",options)
                                            title_answer_option_number = []
                                            for k in title_answer_option:
                                                for j in range(len(options)):
                                                    if options[j] == k:
                                                        title_answer_option_number.append(j)
                                            # print("title_answer_option_number:",title_answer_option_number)
                                            # print("title_option_changes[title2]:",title_option_changes[title2])
                                            # title_option_changes[title2].sort()
                                            # print("title_option_changes[title2]:", title_option_changes[title2])

                                            for number in title_answer_option_number:
                                                for l in range(len(title_option_changes[title2])):
                                                    # print(number,l,type(number),type(l))
                                                    if number == title_option_changes[title2][l]:
                                                        # print(number, l)
                                                        answers += options[l]
                                            l = list(answers)  # 答案重新排序
                                            l.sort()
                                            answers = "".join(l)
                                            change_text = "新答案：%s 已经将原题目选项%s更改为ABCD" % (answers, optionss)
                                            # print("新答案：%s 已经将原题目选项%s更改为ABCD" % (answers, optionss))
                                        else:
                                            # print("选项未更改")
                                            change_text = ""

                                    except:
                                        change_text = ""

                                    return change_text

                                for title in titlepras[titlepra]:
                                    title = "A" + str(title)
                                    try:
                                        if wordtitle_answer_find(title)[0:3] == "新答案":
                                            title_answer_types[title] = "word_choice"

                                        else:
                                            Document(r"word题目汇总\%s.docx" % title)
                                            title_answer_types[title] = "word_calculation"
                                    except:
                                        try:
                                            if wordtitle_answer_find_excel(title) != "未找到答案":
                                                title_answer_types[title] = "text"
                                            else:
                                                title_answer_types[title] = "img1"
                                        except:
                                            title_answer_types[title] = "img2"

                                    # print("title_answer_types[title]:",title,title_answer_types[title] )
                                # print("titlepras[titlepra]:",titlepras)

                                for m in titlepras[titlepra]:
                                    # print('j,a[i]',i)
                                    m = 'A' + str(m)
                                    # print("开始输入题目的答案",m)
                                    # output(m)
                                    new_section.left_margin = Cm(1.27)
                                    new_section.right_margin = Cm(1.27)
                                    new_section.top_margin = Cm(1.0)
                                    new_section.bottom_margin = Cm(1.0)
                                    # document.add_heading('题目：%s答案' % (wrong_number,titlepra,m),2)

                                    answer_change_text = wordtitle_answer_find(m)
                                    if i > 0:
                                        # print("i:",i,titlepras[i-1])
                                        # 题号前加一个A表示答案Answer
                                        previous_title = "A" + str(student_wrongs[name][i - 1]["错题"][0])
                                        # print("previous_title:", previous_title)
                                        # print("m:", m)
                                        # print(title_answer_types)
                                        # print("类型：", title_answer_types[previous_title])
                                        # print("类型：", title_answer_types[previous_title], title_answer_types[m])
                                        if title_answer_types[previous_title] == "text" and \
                                                title_answer_types[m] == "text":
                                            # print("本题和上题答案都是text类型:", title,wordtitle_answer_find_excel(m))
                                            paragraph.add_run(
                                                '%d、答案：%s' % (pra_number, wordtitle_answer_find_excel(m)))
                                            # print("AA")
                                            pra_number += 1
                                            continue
                                    # paragraph=document.add_paragraph( '%d、答案：' % (pra_number))
                                    paragraph = document.add_paragraph('%d、' % (pra_number))
                                    # if answer_change_text!="":
                                    #     paragraph = document.add_paragraph(text='%s' % ( answer_change_text))
                                    #     paragraph.paragraph_format.space_after = Pt(0)
                                    wordtitle_answer(document, m, title_option_changes, pra_number)
                                    pra_number += 1
                                    # titlepras=student_practices[name]

                            else:
                                titlewrongs.append(titlepras[titlepra])

                name_number += 1

            print("总共有%d题没有对应的训练题分别是%s" % (len(titlewrongs), titlewrongs))

            # document_name = r'%s\%s的“题不二错”%s.docx' % (title_wrong_pra_path, team, time_now)
            def AddFooterNumber(run):
                """在页脚中插入页码"""
                fldChar1 = OxmlElement('w:fldChar')  # creates a new element
                fldChar1.set(qn('w:fldCharType'), 'begin')  # sets attribute on element
                instrText = OxmlElement('w:instrText')
                instrText.set(qn('xml:space'), 'preserve')  # sets attribute on element
                instrText.text = 'Page'
                fldChar2 = OxmlElement('w:fldChar')
                fldChar2.set(qn('w:fldCharType'), 'separate')
                t = OxmlElement('w:t')
                t.text = "Seq"
                fldChar2.append(t)
                fldChar4 = OxmlElement('w:fldChar')
                fldChar4.set(qn('w:fldCharType'), 'end')
                r_element = run._r
                r_element.append(fldChar1)
                # print("r_element1:",r_element.__doc__)

                r_element.append(instrText)
                # print("r_element2:", r_element, type(r_element))
                # print("r_element2:", r_element, type(r_element[0]))
                # print("r_element2:", r_element, len(r_element[0]))

                r_element.append(fldChar2)
                # print("r_element3:", r_element)
                r_element.append(fldChar4)
                # print("r_element4:", r_element)

            def InsertPageNumber(Doc):
                """插入页脚内容"""
                footer = Doc.sections[0].footer  # 获取第一个节的页脚
                footer.is_linked_to_previous = True  # 编号续前一节
                paragraph = footer.paragraphs[0]  # 获取页脚的第一个段落
                paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER  # 页脚居中对齐
                run_footer = paragraph.add_run()  # 添加页脚内容
                AddFooterNumber(run_footer)
                font = run_footer.font
                font.name = 'Times New Roman'  # 新罗马字体
                font.size = Pt(10)  # 10号字体
                font.bold = True  # 加粗

            InsertPageNumber(document)
            print("save file:")
            # document.save(r'%s\%s的“题不二错”%s.docx'%(title_wrong_pra_path,team,time_now))  # 题目登记中某一表格学生的错一练三
            document.save(r'%s/%s错题重做%s.docx' % (title_wrong_pra_path, self.students.name,time_now))  # 题目登记中某一表格学生的错一练三
            # os.startfile(r'%s\%s的“题不二错”%s.docx'%(title_wrong_pra_path,team,time_now))
            os.startfile(r'%s/%s错题重做%s.docx' % (title_wrong_pra_path, self.students.name,time_now))
            # 保存至云盘
            # try:
            #     print("main_path:",main_path)
            #     doc_network_disk=Document(r"安装必备文件\network_disk.docx")
            #     network_disk_path=doc_network_disk.paragraphs[0].text
            #     print("network_disk_path:",network_disk_path)
            #
            #     document.save(r'%s/%s的“题不二错”%s.docx'%(network_disk_path,team,time_now))
            #     tk.messagebox.showinfo(title="保存成功", message=        "文档保存至云盘成功")
            #     # os.startfile(r'%s/%s的“题不二错”%s.docx'%(network_disk_path,team,time_now))
            # except Exception as e:
            #     print("保存文件至云盘失败",e)
            #     yesno=tk.messagebox.askyesno(title="错误提示", message=
            #     "文档保存至云盘失败，是否打开云盘并登陆，账号为1124022，密码为12345678")
            #     if yesno==True:
            #         try:
            #             import subprocess
            #             # 定义要打开的软件路径或名称
            #             software_path = "C:\Program Files\Richtech\EDataStore\EmployeeLogin.exe"
            #             # 调用系统默认程序打开指定的文件或URL
            #             subprocess.Popen(software_path)
            #             yesno2 = tk.messagebox.askyesno(title="提示", message=
            #             "继续保存文档至云盘")
            #             if yesno2:
            #                 document.save(r'U:\%s的“题不二错”%s.docx' % ( team, time_now))
            #                 os.startfile(r'U:\%s的“题不二错”%s.docx' % ( team, time_now))
            #
            #         except FileNotFoundError:
            #             print("未能找到该软件")
            #             tk.messagebox.showinfo(title="错误提示", message=
            #             "文档保存至云盘失败，请前往“%s”查看"%(title_wrong_pra_path))
            #             # os.startfile(r'%s\%s的“题不二错”%s.docx'%(title_wrong_pra_path,team,time_now))

            print(r'前往“%s\%s的“题不二错””查看' % (title_wrong_pra_path, team))
            document_name, doc_path = r'%s的“题不二错”%s.docx' % (team, time_now), r'%s\%s的“题不二错”%s.docx' % (
            title_wrong_path, team, time_now)
            # document_name,title_wrong_pra_path=r'%s\%s的“题不二错”%s.docx'%(team,time_now),title_wrong_path
            # return document_name,doc_path
            return document_name, team, title_wrong_path

        def main_zujuan(test):
            yesno = tk.messagebox.askyesno(title="保存确认",
                                           message="保存word会将录入的题目数据转换为word题目，大约需要几分钟，是否保存")
            if yesno == False:
                return
            # tk.messagebox.showinfo(title="错误提示", message="开始输出word文档，大约需要几分钟")
            main_path = os.path.dirname(os.path.realpath(sys.argv[0]))
            time_now = time.strftime("%Y.%m.%d.%H.%M.%S", time.localtime())
            # for i in range(100):
            #     time.sleep(0.5)
            #     progressbarOne['value'] += 1
            #     win.update()

            # season_wrong()
            # title_not_exit = []  # 所有不存在的题目
            # excel_old='气体 尝试题目分类.xlsx'
            # main_path = os.getcwd()

            # print("main_path:", main_path)
            # title_replace_path = main_path + r"\重复题目列表.xlsx"
            # finish_name_one=excel(main_path+r"\题目筛选\题目完成情况记录表单次.xlsx")
            # finish_name_all=excel(main_path+r"\题目筛选\题目完成情况记录表汇总.xlsx")
            # finish_title_all=excel(main_path+r"\题目筛选\学生已经打印题目记录汇总.xlsx")
            # excel_old = '牛顿定律 尝试题目分类3.xlsx'
            # excel_old = '恒定电流题目分类训练题.xlsx'
            # excel_old = '2牛顿定律 尝试题目分类.xlsx'
            # excel_new=excel_rename(excel_old)
            # wb2 = load_workbook(excel_old)
            # sheet = wb2.get_sheet_by_name("题目分类")
            # print('sheet:',sheet)
            # rowmax = sheet.max_row
            # colmax = sheet.max_column
            # heji = findcol('合计')
            # bianhao = findcol('编号')
            # object = findcol("对象")
            # notice = "注意事项：" \
            #          "1、本题型为实验训练，题型分类有可能不完善，可酌情做题。\n" \
            #          "2、难度：1为容易，2为中等，3为困难\n " \
            #          "3、每一题目的分类依据为该题目最难的知识点。\n " \
            #          "4、尽量记录自己的作答时间，正确率，对题型的熟练度和学习效果\n" \
            #          "5、题目完成后欢迎大家反馈建议。具体如下：\n" \
            #          "(1)题目的针对性如何（选填很好，较好，一般，不好，很不好）\n" \
            #          "(2)每一题目的数量几个合适\n" \
            #          "(3)每一题目的难度设置是否合理，你倾向于什么样的难度\n" \
            #          "(4)你完成了几个题目，正确率为多少\n" \
            #          "(5)答案是否详细，能否看明白，不能看明白的题目比例大约多大\n" \
            #          "(6)你用了多少时间完成这一内容，效果如何（选填很好，较好，一般，不好，很不好）\n" \
            #          "(7)你是否还会再一次复习这些题目,有何计划\n" \
            #          "(8)你的其他建议\n" \
            #          ""
            # notice2 = "难度：1为容易，2为中等，3为困难\n解决一个又一个问题，找到解决问题的方法"
            # notice3 = "难度：从1到10难度逐渐增加，请根据个人基础决定作答数量\n解决一个又一个问题，找到解决问题的方法"
            # time_now2 = time.strftime('%Y{y}%m{m}%d{d}').format(y='年', m='月', d='日')
            # author = "整理人:赵安宁 \n整理时间：%s " % (time_now)
            # notice4 = "%s错题重做 姓名_____" % time_now2
            # 将研究对象的字母替换为文字显示在word文档中
            # content = excel(r'F:\个性化题库20180423\title practice\题目分类方法1.xlsx')
            # content = excel(r'%s\题目分类方法1.xlsx' % main_path)
            # print('content:', content)

            # 查找学生的错题
            student_wrong_path = r'%s\学生错题登记\组卷' % main_path
            # student_wrong_path = r'F:\原题重做\错一练三\题目登记'
            # find_stu_tit_wrong(student_wrong_path)
            # from picture_rename import dir
            def dir(file_dir):
                for root, dirs, files in os.walk(file_dir):
                    return [dirs, files]

            files = dir(r"学生错题登记/组卷")
            print("files[1]:", files[1])
            if len(files[1]) == 0:
                print("列表为空")
                tk.messagebox.showinfo(title="错误提示", message="您还没有登记错题")
            for file_number in range(len(files[1])):  # 对每一个错题登记的文件进行错题word文档编辑，并输出最终word文件
                # tk.Label(win,text="第%s个班级题目处理中，共%s个班级\n" % (file_number + 1, len(files[1]))).pack()
                label_1.set("第%s个班级题目处理中，共%s个班级\n" % (file_number + 1, len(files[1])))
                win.update()
                file_name = files[1][file_number]
                print("第%s个班级信息操作中，共%s个班级\n" % (file_number + 1, len(files[1])))
                # text.insert(tk.INSERT, "第%s个班级作息操作中，共%s个班级\n" % (file_number+1, len(files[1])))
                print("file name i:", file_name)
                # if "错一练三" in i:print("错一练三",i)
                name = find_class_and_unit(file_name)
                classs = name["grade"] + "级" + name['classs'] + folder[1]
                makedir(r'%s\%s\%s' % (main_path, folder[0], classs))
                main_number = '2020'
                print('学生错题登记：', file_name)  # 此处创建各个层级的文件夹
                # unit = '第' + name['unit'] + '单元'
                # makedir(r'%s\%s\%s\%s' % (main_path,folder[0],classs,classs + unit))
                # title_wrong_path = r'%s\%s\%s\%s\%s' % (main_path,folder[0],classs,classs + unit,classs + unit + folder[3])
                title_wrong_path = r'%s\%s\%s' % (main_path, folder[0], classs)
                # makedir(title_wrong_path)
                title_wrong_pra_path = r'%s\%s\%s' % (main_path, folder[0], classs)
                # finish_name_all_dict = name_sift()
                student_wrongs = student_wrong4(student_wrong_path, file_name, classs, main_number)

                progressbarOne['maximum'] = len(student_wrongs[0])
                # 进度值初始值
                progressbarOne['value'] = 1
                win.update()

                # print("student_wrongs:",student_wrongs)
                # print("student_wrongs:",student_wrongs[0])
                # print("len(student_wrongs[0]):",len(student_wrongs[0]))
                for i in range(len(student_wrongs)):
                    print('主程序student_wrongs:', i)
                    # print('主程序student_wrongs:', student_wrongs[i].keys())
                    # tk.Label(win,text="第%d个学生作息操作中，共%s个学生"%(i,len(student_wrongs))).pack

                    # text.insert(tk.INSERT, "%s:,第%个学生作息操作中，共%个学生"%(i["team"],len(i)))
                    if len(student_wrongs[i]) != 1:
                        word_name, team, word_path = class_title_pra_zujuan(
                            student_wrongs[i], classs, time_now, test,
                            title_wrong_path, title_wrong_pra_path)  # 题目登记中某一表格学生的错一练三
                print("总共有错题登记班级数：", len(student_wrongs), student_wrongs)
                for i in student_wrongs:
                    print("总共有错题登记人数：", i["team"], len(i))
                    # if len(i)>2:
                    #     tk.messagebox.showinfo(title="提示",message="%s总共有错题登记人数：%s人"%( i["team"], len(i)-1))
            # tk.Label(win,text="操作完成，请前往“%s学生错题word文档”文件夹查询结果"%main_path)
            # label_2.set("操作完成，请前往“%s\学生错题word文档”文件夹查询结果"%main_path)
            # os.startfile(r"%s/学生错题word文档/2023级12班/%s"%(main_path,word_name))
            win.update()
            return r"%s/学生错题word文档/2023级12班/%s" % (main_path, word_name)

            # for i in files[1]:
            #     print("file name i:",i)
            #     if "错一练三" in i: print("错一练三",i)
            # username = 'esgz201610@126.com'
            # password = 'AUCLOGJRHDYLFBBT'
            # username = 'esgz_zhaoanning@126.com'
            # password = 'NFPZSOMXDEQPIXOJ'
            # print("receive email: ",username)

            # while 1!=0:#开始题目的文档编辑，优先处理邮箱的错题登记，若无邮箱的错题登记，再处理本地电脑上的文件。

            # try:
            #     headers,msg_count=email_get()
            # except:pass
            # if msg_count != 0:
            #     try:
            #         email_message["Word"],email_message["Path"]=word_name,word_path
            #     except:
            #         email_message["Word"],email_message["Path"]="no","no"
            #         print("文件生成不成功，请查找原因")
            #     email_send(email_message)

            # print("email_message:",email_message)
            # if msg_count != 0:
            #     email_delete()

            # email_get()
            # text.insert(tk.INSERT, )

            title_not_exit_print(title_not_exit)
            # name_sift_update()
        main_zujuan(self)
        # clear(self.gui)

        import shutil
        main_path = os.path.dirname(os.path.realpath(sys.argv[0]))
        files = dir(r"%s/学生错题登记/组卷" % (main_path))
        # print("删除登记文件，",self.students.name,files[1])
        for file_number in range(len(files[1])):
            file_name = files[1][file_number]

            source_path = r"%s/学生错题登记/组卷/%s" % (main_path, file_name)

            destination_path = r"%s/学生错题登记/已经打印/%s" % (main_path,file_name)
            makedir(r"学生错题登记/已经打印")
            # print("删除登记文件2，", self.students.name)
            if self.students.name != "赵安宁":
                print("删除登记文件3，", self.students.name,main_path,file_name)
                shutil.move(source_path, destination_path)


        # butt_print()
    def change_up(self,number):
        """改变题目的顺序,上移一个"""
        middle_title=self.title_all[number-1]
        self.title_all[number-1]=self.title_all[number]
        self.title_all[number]=middle_title
        self.test_paper_show()
    def change_down(self,number):
        """改变题目的顺序，下移一个"""
        middle_title=self.title_all[number+1]
        self.title_all[number+1]=self.title_all[number]
        self.title_all[number]=middle_title
        self.test_paper_show()
    def change_last(self,number):
        """改变题目顺序至最后一题"""
        middle_titles={}
        print("change_last_number:",number)
        print("self.title_all:",self.title_all)

        for middle_number in range(1,number):
            print("middle_number1:",middle_number,self.title_all[middle_number])
            middle_titles[middle_number]=self.title_all[middle_number]
        for middle_number in range(number,len(self.title_all)-4):#减去4种题目类型
            print("middle_number2:", middle_number, self.title_all[middle_number])
            middle_titles[middle_number]=self.title_all[middle_number+1]
        for title_type in self.title_all:
            if type(title_type) !=int:
                middle_titles[title_type]=self.title_all[title_type]
        middle_titles[len(self.title_all)-4]=self.title_all[number]
        print("middle_titles22:",middle_titles)
        self.title_all=middle_titles
        print("middle_titles22:",self.title_all)
        self.test_paper_show()
    def delete(self,number):
        """删除题目"""

        middle_number=0
        title_type=self.title_all[number]["题型"]
        middle_dict={}
        del self.title_all[number]
        for number1 in self.title_all:
            if type(number1)==int and number1>number:
                print("number1:",number1)
                middle_dict[number1-1]=self.title_all[number1]
            else:
                middle_dict[number1] = self.title_all[number1]
        middle_dict[title_type]-=1

        print("middle_dict:",middle_dict)



        self.title_all=middle_dict
        self.test_paper_show()

@handle_exceptions(logger)
def zujuan():
    """组卷功能"""
    students=register()
    test_paper=Test_paper(students)

@handle_exceptions(logger)
def set_network_disk():
    """设置网盘路径"""
    try:
        from tkinter import filedialog
        folderpath = tk.filedialog.askdirectory()
        doc_network_disk=Document()
        doc_network_disk.add_paragraph().add_run(folderpath)
        doc_network_disk.save(r"安装必备文件\network_disk.docx")
    except Exception as e:
        print("设置网盘路径失败：",e)
class Change_ppt_time:
    def __init__(self,students,prs,title_all_number):
        self.students=students
        self.prs=prs
        self.title_all_number=title_all_number
        self.gui = tk.Toplevel()  # create gui window
        self.gui.title("设置幻灯片的切换时间")  # set the title of gui
        self.gui.geometry("800x600")  # set the window size of gui
        self.title_all_number = 1
        # print("students.title_message:",students.title_message)
        title=students.title_message.loc[0, students.name]
        # print("title:::", title, type(title), type(str(title)))
        if str(title) != "nan":
            # self.gui.refresh(title)

            self.label3 = tk.Label(self.gui, text="设置本题%s的作答时间：(秒钟)"%title)
            self.label3.pack()
            self.sv3 = tk.StringVar()
            # sv3.trace("w", self.f)
            self.entry_input3 = tk.Entry(self.gui, textvariable=self.sv3)
            self.entry_input3.bind("<Return>", self.refresh)  # 回车响应
            # entry_input3.bind('<Key>', self.update)
            self.entry_input3.pack()

            # button=tk.Button(self.gui,text="下一题",command=Gui)
            # button.pack()

            try:
                image = Image.open(r"cach/pic/%s.jpg" % title)  # open image from path
            except Exception as e:
                print("插入图片失败：",e.__traceback__.tb_lineno,e)
                image = Image.open(r"图片题目汇总/%s.jpg" % title)  # open image from path
            image = image.resize((600, 450))  # 调整图像大小
            img = ImageTk.PhotoImage(image)  # read opened image

            self.label1 = tk.Label(self.gui, image=img)  # create a label to insert this image
            self.label1.pack()  # set the label in the main window

        self.gui.mainloop()  # start mainloop
    def update(self,title,*args):
        #更新题目图片
        print("title,*args:",title,*args)
        # self.label1.destroy()#删除原图片
        self.label3.configure(text="设置本题%s的作答时间："%title)
        len_entry=self.entry_input3.get()
        self.entry_input3.delete(0,len_entry)
        try:
            image = Image.open(r"cach/pic/%s.jpg" % title)  # open image from path
        except Exception as e:
            print("插入图片失败：", e.__traceback__.tb_lineno, e)
            image = Image.open(r"图片题目汇总/%s.jpg" % title)  # open image from path
        image = image.resize((600,450))  # 调整图像大小
        img = ImageTk.PhotoImage(image)  # read opened image
        # new_image = Image.open("new_image.jpg")
        # new_tk_image = ImageTk.PhotoImage(new_image)
        # label.configure(image=new_tk_image)
        self.label1.config(image=img)#更新
        self.label1.image=img#图片更新要有这一行
        # tk.update()
        print("图片更新成功")

    def refresh(self,title,*args):
        switch_time = float(self.entry_input3.get())*1000
        # switch_time = float(2.5)*1000
        # print("switchtime:",switch_time)
        try:
            for element in self.prs.slides[self.title_all_number]._element[2][0]:
                element.set('advTm', "%s" % (switch_time))  # 设置幻灯片切换时间,一般为几分钟
                # print("设置时间成功", switch_time)
        except Exception as e:
            print("设置时间失败：", e,e.__traceback__.tb_lineno)
        try:
            title=self.students.title_message.loc[self.title_all_number, self.students.name]
            # title=self.students.title_message.loc[0, self.students.name]
            print("title22:",title)
            self.update(title)
            self.title_all_number += 1
        except Exception as e:
            self.gui.destroy()
            print("图片更新失败，改为添加答案：",e.__traceback__.tb_lineno,e)

            # 插入答案
            print("#插入答案")
            title_all_number_answer = 1
            switch_time = 50000  # 答案切换时间50秒，手动切换



            #先分页插入每题的详细答案
            for title in self.students.title_message.loc[:, self.students.name]:
                # print("title:::",title,type(title),type(str(title)))
                if str(title) != "nan":

                    title = "A" + title
                    if title[5] == "1":
                        bixiu = "必修%s" % title[6]
                    else:
                        bixiu = "选修%s" % title[6]

                    if "Z" in title:
                        if "KK" in title:
                            chapter = title[7:title.find("Z")]
                            # print(chapter)
                            section = title[title.find("Z") + 1:title.find("KK") - 1]
                            class_hour = title[title.find("KK") - 1]
                            title_number = title[title.find("KK") + 2:]
                            title_from = "%s第%s章专题%s第%s课时第%s题" % (
                                bixiu, chapter, section, class_hour, title_number)
                        else:
                            chapter = title[7:title.find("Z")]
                            # print(chapter)
                            section = title[title.find("Z") + 1:title.find("K")]
                            title_number = title[title.find("K") + 1:]
                            title_from = "%s第%s章专题%s第%s题" % (bixiu, chapter, section, title_number)

                    else:
                        if "KK" in title:
                            chapter = title[7:title.find("KK") - 2]
                            # print(chapter)
                            section = title[title.find("KK") - 2:title.find("KK") - 1]
                            class_hour = title[title.find("KK") - 1]
                            title_number = title[title.find("KK") + 2:]
                            title_from = "%s第%s章第%s节第%s课时第%s题" % (
                                bixiu, chapter, section, class_hour, title_number)
                        else:
                            chapter = title[7:title.find("K") - 1]
                            # print(chapter)
                            section = title[title.find("K") - 1:title.find("K")]
                            title_number = title[title.find("K") + 1:]
                            title_from = "%s第%s章第%s节第%s题" % (bixiu, chapter, section, title_number)
                    print(title_from, title)
                    # blank_slide_layout = self.prs.slide_layouts[6]
                    # slide = prs.slides.add_slide(blank_slide_layout)
                    slide = self.prs.slides[title_all_number_answer+self.title_all_number]

                    left = Inches(0)
                    top = Inches(1)
                    width = Cm(33)
                    height = Cm(21)
                    try:
                        docx_to_picture(title)
                        slide.shapes.add_picture(r"cach/pic/%s.jpg" % title, left, top,width,height)
                    except Exception as e:
                        print("插入图片失败：", e.__traceback__.tb_lineno, e)
                        try:
                            slide.shapes.add_picture(r"图片题目汇总/%s.jpg" % title,
                                                     left, top,width,height)
                        except Exception as e:
                            print("---"*10,)
                            print("插入图片失败：", e.__traceback__.tb_lineno,title, e)
                            slide.shapes.add_picture(r"图片题目汇总/00.jpg",
                                                     left, top, width, height)

                    left = Inches(1)
                    top = Inches(0)
                    width = Inches(8)
                    height = Inches(2)

                    # textbox = slide.shapes.add_textbox(left, top, width, height)
                    # tf = textbox.text_frame
                    # left = top = Inches(1)
                    # width = Inches(6)
                    # height = Inches(2)

                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    tf = textbox.text_frame
                    # tf.text = "This is text inside a textbox"

                    p = tf.add_paragraph()
                    p.text = "第%d题答案" % (title_all_number_answer)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)
                    p = tf.add_paragraph()
                    p.text = "题目来源:%s/题库编号:%s" % (title_from, title)

                    try:
                        for element in self.prs.slides[title_all_number_answer+
                                                       self.title_all_number]._element[2][0]:
                            element.set('advTm', "%s" % (switch_time))  # 设置幻灯片切换时间
                    except Exception as e:
                        print("设置时间不成功2：",e,e.__traceback__.tb_lineno)
                    title_all_number_answer += 1

            # print("xml:", prs.slides[0]._element.xml)
            ##解析pptx结构的xml结构，
            # for element in prs.slides[0]._element[2][0]:
            #     print("element:", help(element))
            #     print("element:", element)
            #     print("element:", element.__len__())
            #     # print("element:",element.__contains__())
            #     print("element:", element.attrib['advTm'])
            #     print("element:", element.set('advTm', "250000"))
            #     print("element:", element.attrib['advTm'])
            #     print("element:", element.tag)
            #     print("element:", element.values)
            #     print("element:", element.keys)

            # print("element:",help(element))
            # print("xml:",prs.slides[0]._element)
            # prs.slides[0]._element.makeelement("p:transition")
            # prs.slides[0]._element.get_or_add_timing()
            # print("time:",help(prs.slides[0]._element.timing))
            # print("time:",prs.slides[0]._element.timing.values)
            # 删除多余的幻灯片

            #再在最后一张ppt里面插入所有选择题的答案，再插入所有计算题的答案
            select_answer_list=pd.DataFrame(pd.read_excel(r"cach/excel/select_answer.xlsx"))
            text=""
            number=1
            slide = self.prs.slides[title_all_number_answer + self.title_all_number]
            for index, row in select_answer_list.iterrows():
                # 通过row可以访问每一行的数据，例如：
                print( row[0], row[1])
                if number%5==0:
                    text += "%2s  %3s  " % (row[0], row[1])
                    print("text34:",text)
                    left = Inches(1)
                    top = Inches(0+(number//5))
                    width = Inches(8)
                    height = Inches(2)

                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    tf = textbox.text_frame
                    # tf.text = "This is text inside a textbox"
                    p = tf.add_paragraph()
                    p.text = text
                    p.font.bold = True
                    p.font.size =Cm(1)
                    p.font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)
                    text=""
                else:
                    text+="%2s  %3s  "%(row[0],row[1])
                number+=1
            #添加最后一行答案
            if number%5!=0:
                top = Inches(0 + (number // 5)+1)
                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                # tf.text = "This is text inside a textbox"
                p = tf.add_paragraph()
                p.text = text
                p.font.bold = True
                p.font.size = Cm(1)
                p.font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)

            len_prs_slides = len(self.prs.slides)
            for number in range(len_prs_slides, title_all_number_answer+self.title_all_number+1, -1):
                del self.prs.slides._sldIdLst[-1]
                # prs._sldIdLst.remove(prs.slides[number])

            self.prs.save(r"学生错题word文档\错题重做%s.pptx" % time_now)
            os.startfile(r"学生错题word文档\错题重做%s.pptx" % time_now)

            return

        # self.gui.destroy()
        # self.update(title)
@handle_exceptions(logger)
def butt_print():
    """打印结果，会删除保存的错题登记excel"""
    def title_pptx():
        """生成ppt文件"""
        students=register()
        # if students.name=="赵安宁":
        # yesno=tk.messagebox.askyesno("询问","是否生成错题重做的pptx文件？")

        # if yesno:
        prs = Presentation("sys_files/sample.pptx")
        left = Inches(1)
        top = Inches(0)
        width = Cm(20)
        height = Cm(2)
        textbox =  prs.slides[0].shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        # tf.text = "This is text inside a textbox"

        p = tf.add_paragraph()
        title_time = "3"
        p.text = "%s错题重做" % (time_now2)
        p.font.bold = True
        # titles = ["202322512KK1", "202322512KK2"]
        title_all_number = 1
        #插入题目
        for title in students.title_message.loc[:, students.name]:
            # print("title:::",title,type(title),type(str(title)))
            if str(title) != "nan":
                if title[4] == "1":
                    bixiu = "必修%s" % title[5]
                else:
                    bixiu = "选修%s" % title[5]

                if "Z" in title:
                    if "KK" in title:
                        chapter = title[6:title.find("Z")]
                        # print(chapter)
                        section = title[title.find("Z") + 1:title.find("KK") - 1]
                        class_hour = title[title.find("KK") - 1]
                        title_number = title[title.find("KK") + 2:]
                        title_from = "%s第%s章专题%s第%s课时第%s题" % (
                            bixiu, chapter, section, class_hour, title_number)
                    else:
                        chapter = title[6:title.find("Z")]
                        # print(chapter)
                        section = title[title.find("Z") + 1:title.find("K")]
                        title_number = title[title.find("K") + 1:]
                        title_from = "%s第%s章专题%s第%s题" % (bixiu, chapter, section, title_number)

                else:
                    if "KK" in title:
                        chapter = title[6:title.find("KK") - 2]
                        # print(chapter)
                        section = title[title.find("KK") - 2:title.find("KK") - 1]
                        class_hour = title[title.find("KK") - 1]
                        title_number = title[title.find("KK") + 2:]
                        title_from = "%s第%s章第%s节第%s课时第%s题" % (
                            bixiu, chapter, section, class_hour, title_number)
                    else:
                        chapter = title[6:title.find("K") - 1]
                        # print(chapter)
                        section = title[title.find("K") - 1:title.find("K")]
                        title_number = title[title.find("K") + 1:]
                        title_from = "%s第%s章第%s节第%s题" % (bixiu, chapter, section, title_number)
                # print(title_from, title)
                left = Inches(1)
                top = Inches(0)
                width = Inches(6)
                height = Inches(2)
                blank_slide_layout = prs.slide_layouts[6]
                # slide = prs.slides.add_slide(blank_slide_layout)
                textbox = prs.slides[title_all_number].shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                p = tf.add_paragraph()
                p.text = "第%d题,题目来源:%s/题库编号:%s" % (title_all_number,title_from, title)

                slide = prs.slides[title_all_number]

                left = Inches(0)
                top = Inches(1)
                width = Cm(33)
                height = Cm(21)

                try:
                    docx_to_picture(title)
                    slide.shapes.add_picture(r"cach/pic/%s.jpg" % title, left, top, width, height)
                except Exception as e:
                    print("插入图片失败：", e.__traceback__.tb_lineno, e)
                    slide.shapes.add_picture(r"图片题目汇总/%s.jpg" % title, left, top, width, height)

                # pic = slide.shapes.add_picture(r"图片题目汇总/%s.jpg" % title, left, top,width,height)



                # textbox = slide.shapes.add_textbox(left, top, width, height)
                # tf = textbox.text_frame
                # left = top = Inches(1)
                # width = Inches(6)
                # height = Inches(2)

                # textbox = slide.shapes.add_textbox(left, top, width, height)
                # tf = textbox.text_frame
                # # tf.text = "This is text inside a textbox"
                # # pptx_window = tk.Toplevel()
                # # pptx_window.title("题目")
                # # # root.withdraw()
                # # image = Image.open(r"图片题目汇总/%s.jpg"%title)  # 将"path_to_your_image"替换为你自己的图片路径
                # # photo = ImageTk.PhotoImage(image)
                # # label = tk.Label(pptx_window, image=photo)
                # # label.pack()
                #
                # # name = simpledialog.askstring("输入框", "请输入您的姓名：")
                # # age = simpledialog.askinteger("输入框", "请输入您的年龄：")
                # # switch_time=tk.simpledialog.askinteger("%s"%title, "请输入本题的作答时间（分钟）")
                # # tk.s
                # p = tf.add_paragraph()
                # p.text = "第%d题/时间：%s分钟" % (title_all_number, switch_time)
                # p.font.bold = True
                # p.font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)
                # p = tf.add_paragraph()
                # p.text = "题目来源:%s/题库编号:%s" % (title_from, title)

                # for element in slide._element:
                #     print("element:",element)
                # for element in prs.slides[title_all_number]._element[2][0]:
                #     element.set('advTm', "%s"%(switch_time*1000*60))#设置幻灯片切换时间
                title_all_number += 1
        #更改时间

        gui=Change_ppt_time(students,prs,title_all_number,)
    student=register()

    if student.name!="赵安宁":
        # butt_print_yesno = tk.messagebox.askyesno(title="保存确认",
        #                                       message="保存word会将录入的数据转换为word题目，是否保存")
        # if butt_print_yesno == False:
        #
        #     return
        # else:

        main(student.name)
    else:
        pass
        main(student.name)
        # title_pptx()

    # os.remove(preview_word)
    import shutil
    files = dir(student_wrong_path)
    print("files[1]:",files[1],student.name)
    for file_number in range(len(files[1])):
        file_name = files[1][file_number]
        main_path = os.path.dirname(os.path.realpath(sys.argv[0]))
        source_path=r"%s/学生错题登记/%s"%(main_path,file_name)

        destination_path=r"学生错题登记/已经打印/%s"%file_name
        makedir(r"学生错题登记/已经打印")
        if student.name!="赵安宁":
            shutil.move(source_path, destination_path)


    delete_title_show(title_show)


@handle_exceptions(logger)
def button_preview_title(title):
    """预览个人登记的题目，不删除登记的题目"""
    pass
    student=register()
    # student = Student(combobox_grade.get(), combobox_class.get(), combobox_name.get())
    doc=Document()
    doc.add_paragraph(student.name)

    for title in student.title_message.loc[:, student.name]:
        print("title1:",title)
        if str(title)!="nan":
            input_doc = Document(r"word题目汇总\%s.docx"%title)
            copy_docx(doc,input_doc)
    doc.save(r"学生错题word文档\%s.docx"%student.name)
    os.startfile(r"学生错题word文档\%s.docx"%student.name)

    # main_path = os.path.dirname(os.path.realpath(sys.argv[0]))
    # time_now = time.strftime("%Y.%m.%d.%H.%M.%S", time.localtime())
    # # time_now2 = time.strftime('%Y{y}%m{m}%d{d}').format(y='年', m='月', d='日')
    # author = "整理人:赵安宁 \n整理时间：%s " % (time_now)
    # student_wrong_path = r'%s\学生错题登记' % main_path
    # from picture_rename import dir
    #
    # files = dir(student_wrong_path)
    # print("files[1]:",files[1])
    # if len(files[1]) == 0:
    #     print("列表为空")
    #     tk.messagebox.showinfo(title="错误提示", message="您还没有登记错题")
    # for file_number in range(len(files[1])):  # 对每一个错题登记的文件进行错题word文档编辑，并输出最终word文件
    #     # tk.Label(win,text="第%s个班级题目处理中，共%s个班级\n" % (file_number + 1, len(files[1]))).pack()
    #     label_1.set("第%s个班级题目处理中，共%s个班级\n" % (file_number + 1, len(files[1])))
    #     win.update()
    #     file_name = files[1][file_number]
    #     print("第%s个班级信息操作中，共%s个班级\n" % (file_number + 1, len(files[1])))
    #     # text.insert(tk.INSERT, "第%s个班级作息操作中，共%s个班级\n" % (file_number+1, len(files[1])))
    #     print("file name i:", file_name)
    #     # if "错一练三" in i:print("错一练三",i)
    #     name = find_class_and_unit(file_name)
    #     classs = name["grade"] + "级" + name['classs'] + folder[1]
    #     makedir(r'%s\%s\%s' % (main_path, folder[0], classs))
    #     main_number = '2020'
    #     print('学生错题登记：', file_name)  # 此处创建各个层级的文件夹
    #     # unit = '第' + name['unit'] + '单元'
    #     # makedir(r'%s\%s\%s\%s' % (main_path,folder[0],classs,classs + unit))
    #     # title_wrong_path = r'%s\%s\%s\%s\%s' % (main_path,folder[0],classs,classs + unit,classs + unit + folder[3])
    #     title_wrong_path = r'%s\%s\%s' % (main_path, folder[0], classs)
    #     # makedir(title_wrong_path)
    #     title_wrong_pra_path = r'%s\%s\%s' % (main_path, folder[0], classs)
    #     # finish_name_all_dict = name_sift()
    #     student_wrongs = student_wrong4(student_wrong_path, file_name, classs, main_number)
    #
    #     progressbarOne['maximum'] = len(student_wrongs[0])
    #     # 进度值初始值
    #     progressbarOne['value'] = 1
    #     win.update()
    #     for i in range(len(student_wrongs)):
    #         print('主程序student_wrongs', i)
    #         print("student_wrongs[i]:",student_wrongs[i])
    #         # tk.Label(win,text="第%d个学生作息操作中，共%s个学生"%(i,len(student_wrongs))).pack
    #
    #
    #         # text.insert(tk.INSERT, "%s:,第%个学生作息操作中，共%个学生"%(i["team"],len(i)))
    #         if len(student_wrongs[i]) != 1:
    #             word_name, team, word_path = class_title_pra(
    #                 student_wrongs[i], classs, time_now, author,
    #                 title_wrong_path, title_wrong_pra_path)  # 题目登记中某一表格学生的错一练三
    #     print("总共有错题登记班级数：", len(student_wrongs))
    #     for i in student_wrongs:
    #         print("总共有错题登记人数：", i["team"], len(i))
    # # tk.Label(win,text="操作完成，请前往“%s学生错题word文档”文件夹查询结果"%main_path)
    # # label_2.set("操作完成，请前往“%s\学生错题word文档”文件夹查询结果"%main_path)
    # os.startfile(r"%s/学生错题word文档/2023级12班/%s"%(main_path,word_name))
    # win.update()
    # return r"%s/学生错题word文档/2023级12班/%s"%(main_path,word_name)
    #
    # title_not_exit_print(title_not_exit)
    # # name_sift_update()
class Title():
    """构建一个题目的完整信息"""
    def __init__(self,title):
        self.title=title
        titles=pd.DataFrame(pd.read_excel(r"word题目汇总/题目清单.xlsx"))
def clear(window1):
    "关闭弹窗"
    # print("window",window)
    student = Student(combobox_grade.get(), combobox_class.get(), combobox_name.get())
    # btn1=list[1]
    # btn1.deselect()
    # window1=list[0]
    window1.destroy()
def register_ti_fast(*args):
    """快速录入题目，通过多选框勾选"""

    def input_title_test3(title):
        """输入题目，对多选框题目是否存在进行判断，存在则将内容显示在屏幕上"""
        # all_titles=pd.read_excel(r"word题目汇总/题目清单.xlsx")
        # 将表格中的所有数据全部转换为str类型
        # print("title::::", title)
        all_titles = pd.DataFrame(pd.read_excel(r"word题目汇总/题目清单.xlsx")).astype(str)
        # print("title:",title)
        # print("all_titles:",all_titles,type(entry_input.get()),all_titles.iloc[0,0],type(all_titles.iloc[0,0]))
        # print(entry_input.get(),all_titles['题号'])
        # print(all_titles['题号'].isin([int(title)]))
        if (all_titles['题号'].isin([title])).any():  # 判断一个数据是否在表格中
            # print("题号正确3")
            # word_to_pic(title)
            # 预览题目
            title_doc = Document(r"word题目汇总/%s.docx" % title)
            label_title_text = tk.Label(win, text="%s信息如下" % title)
            text = ""
            # title_show.delete(1)
            text = "你已经登记的题号："
            students = register()
            for ti in students.title_message.loc[:, students.name]:
                if str(ti) != "nan":
                    text += str(ti)
                    text += "  "
            # print("text:", text)
            title_show.insert(tk.INSERT, text)
            delete_title_show(title_show)
            image = Image.open(r"%s\图片题目汇总\A%s.jpg" % (main_path, title))
            # image = Image.open(
            #     r"F:\个性化题库20180423\题不二错1.8\题不二错1.8测试版\图片题目汇总\A202311343KK4.jpg")  # 将"path_to_your_image"替换为你自己的图片路径
            photo = ImageTk.PhotoImage(image)
            label = tk.Label(register_window, image=photo)
            label.grid(row=0, column=5, rowspan=15)

            text += "\n你即将添加的题目：\n"

            for paragrah in title_doc.paragraphs:
                # print(paragrah.text)
                text += paragrah.text
            # print("title_doc.text:", text)
            title_show.insert(tk.INSERT, text)

            student = Student(combobox_grade.get(), combobox_class.get(), combobox_name.get()
                              ).select_title(title)
            #
            # entry_input.delete(0, tk.END)
            # title_show.update()


        else:
            print("输入题号错误4", title)
            delete_title_show(title_show)
            tk.messagebox.showinfo(title="错误提示", message=
            "您登记的错题%s题库中暂时还没有收录，请检查是否录入错误，或者与管理员联系" % title)
    from tkinter import Toplevel
    register_window = tk.Toplevel()
    # register_window.geometry()
    register_window.title("快速录入题目")
    # register_window.geometry("1000x600+400+200")  # 设置弹窗大小及位置（这里是300x200像素，在屏幕上居中）
    students=register()
    # print("students:",students.__dict__)
    label=tk.Label(register_window,text="你要输入的题目的章节是:%s%s"
                   %(combobox_chapter.get(),combobox_section.get()))
    #布局，按行、列布局，不考虑宽度
    # label.grid(row=0, column=0)
    title_data = pd.DataFrame(pd.read_excel(r"word题目汇总/题目清单.xlsx"))
    # title_data = pd.DataFrame(pd.read_excel(r"word题目汇总/题目清单.xlsx")).set_index("录入来源")
    #设置书
    # title_data=title_data.loc[combobox_book.get(),"章节"]
    if "考试" in combobox_book.get():
        print("combobox_book.get()2:", combobox_book.get())
        year = int(combobox_chapter.get()[0:4])
        month = int(combobox_chapter.get()[5:7])
        kaoshitime=str(year)+"年"
        my_date = datetime.date(year, month, 1)  # 2023年10月1日

        # 将日期转换为整数，表示自1970年1月1日以来的总天数
        date_as_int = (my_date - datetime.date(1970, 1, 1)).days+45658-20089  # 计算天数差

        # title_data = title_data.set_index("章节")
        title_data=title_data[(title_data["录入来源"]==combobox_book.get() )&
                               (title_data["章节"] == date_as_int) &( title_data["节"] == combobox_section.get() )]

        # print("title_data:",title_data)
        title_data = title_data.set_index("节")
        print("title_data:",title_data["题目"])

        row=0

        for title in title_data["题号"]:
            # title=str(title_pre)+"K"+str(i)
            #button带参数传递，响应函数
            if  title[0]!="A":
                print("title,title[10:]:",title,title[10:])
                btn = tk.Checkbutton(register_window, text="第%s题"%(title[13:]),
                                     command=lambda c=title: input_title_test2(c))
                btn.deselect()
                btn.grid(row=row, column=0)
                row+=1
        save = tk.Button(register_window, text="保存", bg='#BEBEBE',
                         command=lambda c=register_window: clear(c))
        save.grid(row=row, column=1)
    else:
        chapter= combobox_chapter.get()
        section= combobox_section.get()
        print("chapter,section:",chapter,section)
        if "第" in  chapter and "章" in chapter:
            chapter=int(chapter[1:len(chapter)-1])
        if "第" in section and "节" in section:
            section = str(section[1: -1])
        print("chapter,section:",chapter,section)
        title_data = title_data[(title_data["录入来源"] == combobox_book.get()) &
                                (title_data["章节"] == chapter) &                                (title_data["节"] == section)]
        # title_data = title_data[(title_data["录入来源"] == combobox_book.get()) &
        #                         (title_data["章节"] == chapter) ]
        # title_data = title_data[(title_data["录入来源"] == combobox_book.get()) ]
        title_data = title_data.set_index("节")
        print("title_data:",title_data["题目"])
        row=0
        titles_list=[]
        for title in title_data["题号"]:
            # title=str(title_pre)+"K"+str(i)
            #button带参数传递，响应函数
            if  title[0]!="A" and (title not in titles_list):
                titles_list.append(title)
                print("title,title[10:]:",title,title[10:])
                print("len(title_data[代号][0]):",len(title_data["代号"][0]))
                title2=title[len(title_data["代号"][0]):]
                if "L" in title2:
                    title2="例"+title2[1:]
                elif "Z" in title:
                    title2="针对训练"+title2[1:]
                elif "K" in title:
                    title2="课时作业"+title2[1:]
                btn = tk.Checkbutton(register_window, text=title2, command=lambda c=title: input_title_test2(c))
                btn.deselect()
                btn.grid(row=row, column=0)
                row+=1
    # for book in range(6):
    #     if combobox_book.get()=="步步高2023级必修%d"%book:
    #         title_pre="20231"+str(book)#必修
    #         break
    #     elif combobox_book.get()=="步步高2023级选择必修%d"%book:
    #         title_pre="20232"+str(book)#选择选修
    #         break
    # #设置章和节
    # if "考试" in combobox_book.get():
    #     # print("combobox_book.get():", combobox_book.get())
    #     year = int(combobox_chapter.get()[0:4])
    #     month = int(combobox_chapter.get()[5:7])
    #     title_pre+="K%s%s"%(year,month)
    #
    # else:
    #     title_pre += "%s" % combobox_chapter.get()[1,-1]
    # #设置节
    #
    #
    #
    #
    #
    # for section in range(7):
    #     if combobox_section.get() == "第%d节"%section:
    #         title_pre += str(section)  # +第1节
    #         break
    # for special_subject in range(20):
    #     if combobox_section.get() == "专题%d"%special_subject:
    #         title_pre += str("Z%d"%special_subject)  # +专题1
    #         break
    #
    #
    #
    # print("title_pre:",title_pre)
    # # register_ti(title_pre)
    # #如果一节内容有两个课时，则双排展示
    # print("--------------------你需要手动设置具有两个课时的章节列表------------------")
    # two_train_in_one_section=["20231254","20231262","2023125Z10","20231281","2023127Z10",
    #                           "20231131","20231134","20231135","20231112"
    #                           ,"20231113","20231114","2023114Z14","2023114Z15","20231145",
    #                           ]#"2023221Z5"选择必修2专题5第二课时难度较大直接删除
    # three_in_one_section=["20231145"]
    # if title_pre in two_train_in_one_section:
    #     # label=tk.Label(register_window,text="第一课时")
    #
    #     for i in range(1,15):
    #         title=str(title_pre)+"1KK"+str(i)
    #         #button带参数传递，响应函数
    #         print("title：str(title_pre)+1KK+str(i):",title)
    #         btn_var = tk.BooleanVar()
    #         btn = tk.Checkbutton(register_window, text="第一课时（训练一）第%d题"%i,
    #              variable= btn_var,  command=lambda c=title: input_title_test2(c))
    #         btn.deselect()
    #         btn.grid(row=i, column=0)
    #     for i in range(1, 7):
    #         title = str(title_pre) + "1KL" + str(i)
    #         # button带参数传递，响应函数
    #         btn = tk.Checkbutton(register_window, text="第一课时（训练一）例%d" % i,
    #                              command=lambda c=title: input_title_test2(c))
    #         btn.deselect()
    #         btn.grid(row=i + 14, column=0)
    #     for i in range(1, 6):
    #         title = str(title_pre) + "1KZ" + str(i)
    #         # button带参数传递，响应函数
    #         btn = tk.Checkbutton(register_window, text="第一课时（训练一）针对训练%d" % i,
    #                              command=lambda c=title: input_title_test2(c))
    #         btn.deselect()
    #         btn.grid(row=i + 20, column=0)
    #     for i in range(1,15):
    #         title=str(title_pre)+"2KK"+str(i)
    #         #button带参数传递，响应函数
    #         btn = tk.Checkbutton(register_window, text="第二课时(训练二)第%d题"%i,
    #                              command=lambda c=title: input_title_test2(c))
    #         btn.deselect()
    #         btn.grid(row=i, column=1)
    #     for i in range(1, 7):
    #         title = str(title_pre) + "2KL" + str(i)
    #         # button带参数传递，响应函数
    #         btn = tk.Checkbutton(register_window, text="第二课时(训练二)例%d" % i,
    #                              command=lambda c=title: input_title_test2(c))
    #         btn.deselect()
    #         btn.grid(row=i + 14, column=1)
    #     for i in range(1, 6):
    #         title = str(title_pre) + "2KZ" + str(i)
    #         # button带参数传递，响应函数
    #         btn = tk.Checkbutton(register_window, text="第二课时(训练二)针对训练%d" % i,
    #                              command=lambda c=title: input_title_test2(c))
    #         btn.deselect()
    #         btn.grid(row=i + 20, column=1)
    #     if title_pre in three_in_one_section:
    #         for i in range(1, 15):
    #             title = str(title_pre) + "3KK" + str(i)
    #             # button带参数传递，响应函数
    #             btn = tk.Checkbutton(register_window, text="第三课时(训练三)第%d题" % i,
    #                                  command=lambda c=title: input_title_test2(c))
    #
    #             btn.deselect()
    #             btn.grid(row=i, column=2)
    #         for i in range(1, 7):
    #             title = str(title_pre) + "3KL" + str(i)
    #             # button带参数传递，响应函数
    #             btn = tk.Checkbutton(register_window, text="第三课时(训练三)例%d" % i,
    #                                  command=lambda c=title: input_title_test2(c))
    #             btn.deselect()
    #             btn.grid(row=i + 14, column=2)
    #         for i in range(1, 6):
    #             title = str(title_pre) + "3KZ" + str(i)
    #             # button带参数传递，响应函数
    #             btn = tk.Checkbutton(register_window, text="第三课时(训练三)针对训练%d" % i,
    #                                  command=lambda c=title: input_title_test2(c))
    #             btn.deselect()
    #             btn.grid(row=i + 20, column=2)
    # else:
    #     for i in range(1,15):
    #         title=str(title_pre)+"K"+str(i)
    #         #button带参数传递，响应函数
    #         btn = tk.Checkbutton(register_window, text="第%d题"%i,
    #                              command=lambda c=title: input_title_test2(c))
    #         btn.deselect()
    #         btn.grid(row=i, column=0)
    #     for i in range(1,7):
    #         title=str(title_pre)+"L"+str(i)
    #         #button带参数传递，响应函数
    #         btn = tk.Checkbutton(register_window, text="例%d"%i,
    #                              command=lambda c=title: input_title_test2(c))
    #         btn.deselect()
    #         btn.grid(row=i+14, column=0)
    #     for i in range(1,6):
    #         title=str(title_pre)+"Z"+str(i)
    #         #button带参数传递，响应函数
    #         btn = tk.Checkbutton(register_window, text="针对训练%d"%i,
    #                              command=lambda c=title: input_title_test2(c))
    #         btn.deselect()
    #         btn.grid(row=i+20, column=0)
    # save=tk.Button(register_window, text="保存", bg='#BEBEBE',
    #                command=lambda c=register_window: clear(c))
    # save.grid(row=i+21,column=1)
    text = "你已经登记的题号："
    students = register()
    for title in students.title_message.loc[:, students.name]:
        if str(title) != "nan":
            text += str(title)
            text += "  "
    print(text)
    register_window.mainloop()

@handle_exceptions(logger)
def edit_title():
    """学生修改登记的错题"""

    def delete_title(title):
        """删除题目"""

        yesno = tk.messagebox.askyesno(title="删除错题询问", message="是否删除题目：%s" % title)
        students = register()
        print("yesno:", yesno)
        if yesno:
            students.delete_title(title)
            title_show.delete("0.0", tk.END)
            text = "你已经登记的题号："
            print("self.title_message.loc[self.name]:", students.title_message)
            students.title_message=students.title_message.reset_index()
            # print("self.title_message.loc[self.name]:", students.title_message.loc[:, students.name])
            tk.messagebox.showinfo(title="成功提示",message="删除错题%s成功："%title)
            for title in students.title_message.loc[:, students.name]:
                # print("title:::",title,type(title),type(str(title)))
                if str(title) != "nan":
                    text += str(title)
                    text += "  "
            print(text)
            title_show.insert(tk.INSERT, text)
            # label = tk.Label(edit_window, text="你已经删除题目：%s"%title, )

            label.grid(row=0, column=0, sticky="w",columnspan=6)



    edit_window=tk.Tk()
    edit_window.geometry('600x500')
    edit_window.title("修改题目")
    students=register()

    print("students:",students.__dict__)
    label=tk.Label(edit_window,text="以下是你已经录入的题目，点击题号删除不需要的题目吧")
    #布局，按行、列布局，不考虑宽度
    label.grid(row=1,column=0,columnspan=6,)#行列跨越
    # students = register()
    i=8
    print("students.title_message.loc[:, students.name]:",students.title_message.loc[:, students.name])
    for title in students.title_message.loc[:, students.name]:
        # button带参数传递，响应函数
        if str(title) != "nan":
            # button_val = tk.StringVar()
            #使用验证参数 validata,参数值为 focusout 当失去焦点的时候，验证输入框内容是否正确

            # entry_input = tk.Entry(win, textvariable=Py_String, validate="focusout", )
            # entry_input = tk.Entry(win, textvariable=Py_String, validate="focusout", validatecommand=input_title)

            btn = tk.Button(edit_window, text=title,command=lambda c=title: delete_title(c))
            btn.grid(row=i//4, column=i%4,sticky="w")
            i+=1
            # print("i//4:",title,i,i//4,i%4,i//4+1,i%4-1)
# def clear_all(win):
#     for widget in win.winfo_children():
#         widget.destroy()
# def wrong_season(name):
#     """录入所有错题及其错误原因"""
#     student=register()
#     clear_all(win)
#     import tkinter as tk
#
#
#     # win.title("录入所有错题及其错误原因")
#     label=tk.Label(win,text="录入你所有的错题，并选择错误原因，")
#     label2=tk.Label(win,text="选择开始章节：")
#
#     combobox_chapter_wrong = ttk.Combobox(win)
#     combobox_chapter_wrong .bind('<<ComboboxSelected>>')
#     combobox_chapter_wrong ['values'] = ("第5章","第6章","第7章","第8章",)  # 设置下拉框的选项列表
#     combobox_chapter_wrong .current(0)
#
#     combobox_section_wrong = ttk.Combobox(win)
#     combobox_section_wrong['values'] = ("第1节","第2节","第3节","第4节","第5节","第6节",
#                                   "专题1","专题2","专题3","专题4","专题5","专题6",)  # 设置下拉框的选项列表
#     combobox_section_wrong.current(0)
#     combobox_section_wrong.bind('<<ComboboxSelected>>', )
#
#     def title_find():
#         for chapter in range(15):
#             if combobox_chapter_wrong.get() == "第%d章" % chapter:
#                 title_pre = "202312" + str(chapter)  # 必修2第5章
#                 break
#         for section in range(7):
#             if combobox_section_wrong.get() == "第%d节" % section:
#                 title_pre += str(section)  # +第1节
#                 break
#         for special_subject in range(10):
#             if combobox_section_wrong.get() == "专题%d" % special_subject:
#                 title_pre += str("Z%d" % special_subject)  # +专题1
#                 break
#         return title_pre
#         # if combobox_chapter_wrong.get()=="第五章":
#         #     chapter=5
#         # elif combobox_chapter_wrong.get()=="第六章":
#         #     chapter=6
#         # elif combobox_chapter_wrong.get() == "第七章":
#         #     chapter = 7
#         # elif combobox_chapter_wrong.get() == "第八章":
#         #     chapter = 8
#         # if combobox_section_wrong.get()=="第1节":
#         #     section=1
#         # elif combobox_chapter_wrong.get()== "第2节":
#         #     section=2
#         # elif combobox_chapter_wrong.get()== "第3节":
#         #     section=3
#         # elif combobox_chapter_wrong.get()== "第4节":
#         #     section=4
#         # elif combobox_chapter_wrong.get()== "第5节":
#         #     section=5
#         # elif combobox_chapter_wrong.get()== "第6节":
#         #     section=6
#         # elif combobox_chapter_wrong.get()== "第7节":
#         #     section=7
#         # elif combobox_chapter_wrong.get()== "专题1":
#         #     section=2
#         # elif combobox_chapter_wrong.get()== "专题1":
#         #     section=2
#         # elif combobox_chapter_wrong.get()== "专题1":
#         #     section=2
#         # elif combobox_chapter_wrong.get()== "专题1":
#         #     section=2
#     title_pre=title_find()
#     # wrong_title_text.set(1)
#
#     label3 = tk.Label(win, textvariable=wrong_title_text)
#     title_all=pd.DataFrame(pd.read_excel(r"word题目汇总\题目清单.xlsx"))
#
#     # title_all=title_all.loc[:,"题号"]
#     title_all=title_all[title_all["题号"].str.contains(title_pre)]["题号"]#pandas查找
#     student.title_all=title_all
#
#     # print("title_all:",title_all)
#     i=0
#     title=student.title_all.iloc[i]
#     # title=
#     wrong_title_text.set("题目：第%d题,题库编号:%s"%(i+1,title))
#     label4=tk.Label(win, text="请选择你的错误原因：")
#     combobox_wrong_season = ttk.Combobox(win)
#     combobox_wrong_season['values'] = ("审题不清楚","研究对象不明确","受力分析错误","运动过程不清晰",
#                                        "临界状态判断错误","隐含条件未找到","不知道用什么规律",
#                                        "公式记忆错误","计算错误","没有思路","没有时间作答","其他"
#                                        )  # 设置下拉框的选项列表
#     # combobox_wrong_season.current(0)
#     wrong_sea=[title,combobox_wrong_season.get(),i]
#     combobox_wrong_season.bind('<<ComboboxSelected>>',student.wrong_season_next)
#     # wrong_title_text=tk.StringVar()
#     # wrong_title_s=test()
#
#
#     # print("wrong season:",wrong_sea)
#
#     # button_wrong_pre=tk.Button(win,text="上一题",command=lambda: student.wrong_season(wrong_sea))
#     # button_wrong_next=tk.Button(win,text="下一题",command=lambda: student.wrong_season_next(wrong_sea))
#
#     label.grid(row=0,column=0,columnspan=1)
#     label2.grid(row=1,column=0,columnspan=1)
#     combobox_chapter_wrong.grid(row=1,column=1)
#     combobox_section_wrong.grid(row=1,column=2)
#     label3.grid(row=2,column=0)
#     label4.grid(row=3,column=0,)
#     combobox_wrong_season.grid(row=3,column=1)
#     # button_wrong_pre.grid(row=10,column=2,)
#     button_wrong_next.grid(row=10,column=3,)
#     win.mainloop()


def main(test):
    yesno=tk.messagebox.askyesno( title="保存确认",
                                              message="保存word会将录入的题目数据转换为word题目，大约需要几分钟，是否保存")
    if yesno==False:
        return
    # tk.messagebox.showinfo(title="错误提示", message="开始输出word文档，大约需要几分钟")
    main_path = os.path.dirname(os.path.realpath(sys.argv[0]))
    time_now = time.strftime("%Y.%m.%d.%H.%M.%S", time.localtime())
    # for i in range(100):
    #     time.sleep(0.5)
    #     progressbarOne['value'] += 1
    #     win.update()

    # season_wrong()
    # title_not_exit = []  # 所有不存在的题目
    # excel_old='气体 尝试题目分类.xlsx'
    # main_path = os.getcwd()

    # print("main_path:", main_path)
    # title_replace_path = main_path + r"\重复题目列表.xlsx"
    # finish_name_one=excel(main_path+r"\题目筛选\题目完成情况记录表单次.xlsx")
    # finish_name_all=excel(main_path+r"\题目筛选\题目完成情况记录表汇总.xlsx")
    # finish_title_all=excel(main_path+r"\题目筛选\学生已经打印题目记录汇总.xlsx")
    # excel_old = '牛顿定律 尝试题目分类3.xlsx'
    # excel_old = '恒定电流题目分类训练题.xlsx'
    # excel_old = '2牛顿定律 尝试题目分类.xlsx'
    # excel_new=excel_rename(excel_old)
    # wb2 = load_workbook(excel_old)
    # sheet = wb2.get_sheet_by_name("题目分类")
    # print('sheet:',sheet)
    # rowmax = sheet.max_row
    # colmax = sheet.max_column
    # heji = findcol('合计')
    # bianhao = findcol('编号')
    # object = findcol("对象")
    # notice = "注意事项：" \
    #          "1、本题型为实验训练，题型分类有可能不完善，可酌情做题。\n" \
    #          "2、难度：1为容易，2为中等，3为困难\n " \
    #          "3、每一题目的分类依据为该题目最难的知识点。\n " \
    #          "4、尽量记录自己的作答时间，正确率，对题型的熟练度和学习效果\n" \
    #          "5、题目完成后欢迎大家反馈建议。具体如下：\n" \
    #          "(1)题目的针对性如何（选填很好，较好，一般，不好，很不好）\n" \
    #          "(2)每一题目的数量几个合适\n" \
    #          "(3)每一题目的难度设置是否合理，你倾向于什么样的难度\n" \
    #          "(4)你完成了几个题目，正确率为多少\n" \
    #          "(5)答案是否详细，能否看明白，不能看明白的题目比例大约多大\n" \
    #          "(6)你用了多少时间完成这一内容，效果如何（选填很好，较好，一般，不好，很不好）\n" \
    #          "(7)你是否还会再一次复习这些题目,有何计划\n" \
    #          "(8)你的其他建议\n" \
    #          ""
    # notice2 = "难度：1为容易，2为中等，3为困难\n解决一个又一个问题，找到解决问题的方法"
    # notice3 = "难度：从1到10难度逐渐增加，请根据个人基础决定作答数量\n解决一个又一个问题，找到解决问题的方法"
    # time_now2 = time.strftime('%Y{y}%m{m}%d{d}').format(y='年', m='月', d='日')
    # author = "整理人:赵安宁 \n整理时间：%s " % (time_now)
    # notice4 = "%s错题重做 姓名_____" % time_now2
    # 将研究对象的字母替换为文字显示在word文档中
    # content = excel(r'F:\个性化题库20180423\title practice\题目分类方法1.xlsx')
    # content = excel(r'%s\题目分类方法1.xlsx' % main_path)
    # print('content:', content)

    # 查找学生的错题
    student_wrong_path = r'%s\学生错题登记' % main_path
    # student_wrong_path = r'F:\原题重做\错一练三\题目登记'
    # find_stu_tit_wrong(student_wrong_path)
    # from picture_rename import dir
    def dir(file_dir):
        for root, dirs, files in os.walk(file_dir):
            return [dirs, files]
    files = dir(student_wrong_path)
    print("files[1]:",files[1])
    if len(files[1]) == 0:
        print("列表为空")
        tk.messagebox.showinfo(title="错误提示", message="您还没有登记错题")
    for file_number in range(len(files[1])):  # 对每一个错题登记的文件进行错题word文档编辑，并输出最终word文件
        # tk.Label(win,text="第%s个班级题目处理中，共%s个班级\n" % (file_number + 1, len(files[1]))).pack()
        label_1.set("第%s个班级题目处理中，共%s个班级\n" % (file_number + 1, len(files[1])))
        win.update()
        file_name = files[1][file_number]
        print("第%s个班级信息操作中，共%s个班级\n" % (file_number + 1, len(files[1])))
        # text.insert(tk.INSERT, "第%s个班级作息操作中，共%s个班级\n" % (file_number+1, len(files[1])))
        print("file name i:", file_name)
        # if "错一练三" in i:print("错一练三",i)
        name = find_class_and_unit(file_name)
        classs = name["grade"] + "级" + name['classs'] + folder[1]
        makedir(r'%s\%s\%s' % (main_path, folder[0], classs))
        main_number = '2020'
        print('学生错题登记：', file_name)  # 此处创建各个层级的文件夹
        # unit = '第' + name['unit'] + '单元'
        # makedir(r'%s\%s\%s\%s' % (main_path,folder[0],classs,classs + unit))
        # title_wrong_path = r'%s\%s\%s\%s\%s' % (main_path,folder[0],classs,classs + unit,classs + unit + folder[3])
        title_wrong_path = r'%s\%s\%s' % (main_path, folder[0], classs)
        # makedir(title_wrong_path)
        title_wrong_pra_path = r'%s\%s\%s' % (main_path, folder[0], classs)
        # finish_name_all_dict = name_sift()
        student_wrongs = student_wrong4(student_wrong_path, file_name, classs, main_number)

        progressbarOne['maximum'] = len(student_wrongs[0])
        # 进度值初始值
        progressbarOne['value'] = 1
        win.update()


        # print("student_wrongs:",student_wrongs)
        # print("student_wrongs:",student_wrongs[0])
        # print("len(student_wrongs[0]):",len(student_wrongs[0]))
        for i in range(len(student_wrongs)):
            print('主程序student_wrongs:', i)
            # print('主程序student_wrongs:', student_wrongs[i].keys())
            # tk.Label(win,text="第%d个学生作息操作中，共%s个学生"%(i,len(student_wrongs))).pack


            # text.insert(tk.INSERT, "%s:,第%个学生作息操作中，共%个学生"%(i["team"],len(i)))
            if len(student_wrongs[i]) != 1 :
                word_name, team, word_path = class_title_pra(
                    student_wrongs[i], classs, time_now, test,
                    title_wrong_path, title_wrong_pra_path)  # 题目登记中某一表格学生的错一练三
        print("总共有错题登记班级数：", len(student_wrongs),student_wrongs)
        for i in student_wrongs:
            print("总共有错题登记人数：", i["team"], len(i))
            # if len(i)>2:
            #     tk.messagebox.showinfo(title="提示",message="%s总共有错题登记人数：%s人"%( i["team"], len(i)-1))
    # tk.Label(win,text="操作完成，请前往“%s学生错题word文档”文件夹查询结果"%main_path)
    # label_2.set("操作完成，请前往“%s\学生错题word文档”文件夹查询结果"%main_path)
    # os.startfile(r"%s/学生错题word文档/2023级12班/%s"%(main_path,word_name))
    win.update()
    return r"%s/学生错题word文档/2023级12班/%s"%(main_path,word_name)

    # for i in files[1]:
    #     print("file name i:",i)
    #     if "错一练三" in i: print("错一练三",i)
    # username = 'esgz201610@126.com'
    # password = 'AUCLOGJRHDYLFBBT'
    # username = 'esgz_zhaoanning@126.com'
    # password = 'NFPZSOMXDEQPIXOJ'
    # print("receive email: ",username)

    # while 1!=0:#开始题目的文档编辑，优先处理邮箱的错题登记，若无邮箱的错题登记，再处理本地电脑上的文件。

    # try:
    #     headers,msg_count=email_get()
    # except:pass
    # if msg_count != 0:
    #     try:
    #         email_message["Word"],email_message["Path"]=word_name,word_path
    #     except:
    #         email_message["Word"],email_message["Path"]="no","no"
    #         print("文件生成不成功，请查找原因")
    #     email_send(email_message)

    # print("email_message:",email_message)
    # if msg_count != 0:
    #     email_delete()

    # email_get()
    # text.insert(tk.INSERT, )

    title_not_exit_print(title_not_exit)
    # name_sift_update()





if __name__== '__main__':
    import tkinter as tk
    from tkinter import ttk
    import time

    # 创建主窗口
    root = tk.Tk()
    root.title("启动界面")
    root.geometry("300x200")

    # 创建一个标签显示启动信息
    label = ttk.Label(root, text="正在启动...", font=("Arial", 16))
    label.pack(pady=50)

    # 创建一个进度条
    progress = ttk.Progressbar(root, orient="horizontal", length=200, mode="determinate")
    progress.pack(pady=20)


    # 模拟启动过程
    def start_up():
        for i in range(101):
            progress['value'] = i
            root.update_idletasks()
            time.sleep(0.05)  # 模拟加载时间
        root.destroy()  # 关闭启动界面


    # 延迟1秒后开始启动
    root.after(5000, start_up)

    # 运行主循环
    root.mainloop()

    # 启动界面关闭后，继续执行主程序
    print("启动完成，进入主程序...")
    # 这里可以继续写你的主程序代码
    # season_wrong()
    title_not_exit = []  # 所有不存在的题目
    # excel_old='气体 尝试题目分类.xlsx'
    # main_path = os.getcwd()
    main_path = os.path.dirname(os.path.realpath(sys.argv[0]))
    print("main_path:", main_path)
    title_replace_path = main_path + r"\重复题目列表.xlsx"
    # finish_name_one=excel(main_path+r"\题目筛选\题目完成情况记录表单次.xlsx")
    # finish_name_all=excel(main_path+r"\题目筛选\题目完成情况记录表汇总.xlsx")
    # finish_title_all=excel(main_path+r"\题目筛选\学生已经打印题目记录汇总.xlsx")
    # excel_old = '牛顿定律 尝试题目分类3.xlsx'
    # excel_old = '恒定电流题目分类训练题.xlsx'
    excel_old = 'sys_files/2牛顿定律 尝试题目分类.xlsx'
    # excel_new=excel_rename(excel_old)
    wb2 = load_workbook(excel_old)
    sheet = wb2.get_sheet_by_name("题目分类")
    # print('sheet:',sheet)
    rowmax = sheet.max_row
    colmax = sheet.max_column
    heji = findcol('合计')
    bianhao = findcol('编号')
    object = findcol("对象")
    notice = "注意事项：" \
             "1、本题型为实验训练，题型分类有可能不完善，可酌情做题。\n" \
             "2、难度：1为容易，2为中等，3为困难\n " \
             "3、每一题目的分类依据为该题目最难的知识点。\n " \
             "4、尽量记录自己的作答时间，正确率，对题型的熟练度和学习效果\n" \
             "5、题目完成后欢迎大家反馈建议。具体如下：\n" \
             "(1)题目的针对性如何（选填很好，较好，一般，不好，很不好）\n" \
             "(2)每一题目的数量几个合适\n" \
             "(3)每一题目的难度设置是否合理，你倾向于什么样的难度\n" \
             "(4)你完成了几个题目，正确率为多少\n" \
             "(5)答案是否详细，能否看明白，不能看明白的题目比例大约多大\n" \
             "(6)你用了多少时间完成这一内容，效果如何（选填很好，较好，一般，不好，很不好）\n" \
             "(7)你是否还会再一次复习这些题目,有何计划\n" \
             "(8)你的其他建议\n" \
             ""
    notice2 = "难度：1为容易，2为中等，3为困难\n解决一个又一个问题，找到解决问题的方法"
    notice3 = "难度：从1到10难度逐渐增加，请根据个人基础决定作答数量\n解决一个又一个问题，找到解决问题的方法"
    time_now = time.strftime("%Y.%m.%d.%H.%M.%S", time.localtime())
    time_now2 = time.strftime('%Y{y}%m{m}%d{d}').format(y='年', m='月', d='日')
    author = "整理时间：%s " % (time_now)
    # author = "整理人:赵安宁 \n整理时间：%s " % (time_now)
    # notice4 = "%s错题重做 姓名_____" % time_now2
    # 将研究对象的字母替换为文字显示在word文档中
    # content = excel(r'F:\个性化题库20180423\title practice\题目分类方法1.xlsx')
    content = excel(r'%s\sys_files\题目分类方法1.xlsx' % main_path)
    # print('content:', content)

    # 查找学生的错题
    student_wrong_path = r'%s\学生错题登记' % main_path
    # student_wrong_path = r'F:\原题重做\错一练三\题目登记'
    # find_stu_tit_wrong(student_wrong_path)
    # from 安装必备文件/picture_rename import dir


    def dir(file_dir):
        for root, dirs, files in os.walk(file_dir):
            return [dirs, files]

    # student_wrong_path=r'F:\原题重做\错一练三\题目登记'
    files = dir(student_wrong_path)
    # print('files:', files)

    # 调用Tk()创建主窗口

    # INSERT 光标处插入；END 末尾处插入

    folder = ['学生错题word文档', '班', '第x单元', '原题重做', '错一练三']
    # time_now = time.strftime("%Y.%m.%d.%H.%M.%S", time.localtime())
    # author="整理人:赵安宁() 时间：%s "%(time_now)
    # notice="注意事项：" \
    #        "1、本题型为实验训练，题型分类有可能不完善，可酌情做题。\n" \
    #        "2、难度：1为容易，2为中等，3为困难\n " \
    #        "3、每一题目的分类依据为该题目最难的知识点。\n " \
    #        "4、尽量记录自己的作答时间，正确率，对题型的熟练度和学习效果\n" \
    #        "5、题目完成后欢迎大家反馈建议。具体如下：\n" \
    #        "(1)题目的针对性如何（选填很好，较好，一般，不好，很不好）\n" \
    #        "(2)每一题目的数量几个合适\n" \
    #        "(3)每一题目的难度设置是否合理，你倾向于什么样的难度\n" \
    #        "(4)你完成了几个题目，正确率为多少\n" \
    #        "(5)答案是否详细，能否看明白，不能看明白的题目比例大约多大\n" \
    #        "(6)你用了多少时间完成这一内容，效果如何（选填很好，较好，一般，不好，很不好）\n" \
    #        "(7)你是否还会再一次复习这些题目,有何计划\n" \
    #        "(8)你的其他建议\n" \
    #        ""
    # notice2="难度：1为容易，2为中等，3为困难\n解决一个又一个问题，找到解决问题的方法"

    makedir(r'%s\%s' % (main_path, folder[0]))
    main_path = os.getcwd()
    win = tk.Tk()
    # 给主窗口起一个名字，也就是窗口的名字
    win.title('题不二错_测试版')
    # 开启主循环，让窗口处于显示状态
    win.iconbitmap('%s\sys_files\pikaqiu.ico' % main_path)
    win.geometry('600x500')
    # 创建一个文本控件
    # width 一行可见的字符数；height 显示的行数
    # text = tk.Text(win, width=60, height=30, undo=True, autoseparators=False)
    # text.insert(tk.INSERT, "点击开始按钮，开始执行程序，完成后会有提示\n")
    # 适用 pack(fill=X) 可以设置文本域的填充模式。比如 X表示沿水平方向填充，Y表示沿垂直方向填充，BOTH表示沿水平、垂直方向填充
    # text.pack()
    progressbarOne = tk.ttk.Progressbar(win,length=400)
    progressbarOne.pack(side=tk.TOP)

    # confirm_button = tkinter.Button(root, text="确定", command=show)
    # confirm_button.pack(side=tk.TOP)
    label_1 = tk.StringVar()
    # label_1.set("点击开始，执行程序")
    tk.Label(win, textvariable=label_1,).pack()
    label_2 = tk.StringVar()
    # label_2.set("点击开始，执行程序")
    tk.Label(win, textvariable=label_2,).pack()

    tk.Label(win, text="249365523@qq.com",).pack(side="bottom")
    tk.Label(win, text="制作：恩施高中 赵安宁",).pack(side="bottom")




    from tkinter import ttk

    combobox_grade= ttk.Combobox(win)
    combobox_grade['values'] = ("2023级")
    combobox_grade.current(0)

    combobox_class = ttk.Combobox(win)
    combobox_class.bind('<<ComboboxSelected>>', class_register)
    # combobox_class['values'] = ("12班")
    # combobox_class.current(0)
    combobox_class['values'] = ("1班","2班","3班","4班","5班","6班","7班",
        "8班","9班","10班","11班","12班","13班","14班","15班","16班","17班",
        "18班","19班","20班","21班","22班","23班","24班","25班","26班")  # 设置下拉框的选项列表
    combobox_class.current(11)

    # combobox_name = ttk.Combobox(win,postcommand=register)
    combobox_name = ttk.Combobox(win)
    combobox_name['values'] = ("赵安宁","朱佳琦")  # 设置下拉框的选项列表


    combobox_name.current(1)
    combobox_name.bind('<<ComboboxSelected>>', register)
    # combobox.bind('<<ComboboxSelected>>', on_select)  # 绑定事件处理函数到下拉框的选择事件上
    label = tk.Label(win, text="输入自己的个人信息：")
    label.pack()
    combobox_grade.pack()
    combobox_class.pack()
    combobox_name.pack()
    # button_regist_ti_fast.pack()


    label = tk.Label(win, text="选择章节，批量录入题目：")
    label.pack()

    combobox_book = ttk.Combobox(win)
    combobox_book.bind('<<ComboboxSelected>>')
    # combobox_class['values'] = ("12班")
    # combobox_class.current(0)
    def chapter_select(*args):
        """选择书籍，相应的刷新章节"""
        title_data = pd.DataFrame(pd.read_excel(r"word题目汇总/题目清单.xlsx")).set_index("章节")
        if "考试" in combobox_book.get():
            # print("combobox_book.get():", combobox_book.get())
            year=int(combobox_chapter.get()[0:4])
            month=int(combobox_chapter.get()[5:7])
            my_date = datetime.date(year,month,1)  # 2023年10月1日

            # 将日期转换为整数，表示自1970年1月1日以来的总天数
            print("year,month:",year,month)
            date_as_int = (my_date - datetime.date(1970, 1, 1)).days +25569 # 计算天数差
            print("date_as_int:",date_as_int)
            title_data = title_data.loc[date_as_int, "节"]
            chapter_sections = []
            for chapter_section in title_data:
                print("book_chapter:", chapter_section)
                try:
                    chapter_section = str(int(chapter_section))
                except:
                    pass
                if chapter_section not in chapter_sections:
                    chapter_sections.append(chapter_section)
            chapter_sections2 = []
            for chapter_section in chapter_sections:
                if chapter_section not in chapter_sections2:
                    chapter_sections2.append(chapter_section)
            # print("title_data", title_data)
            print("book_chapters2:", chapter_sections2,chapter_sections)
        else:
            title_data = title_data.loc[int(combobox_chapter.get()[1:-1]), "节"]
            chapter_sections=[]
            for chapter_section in  title_data:
                print("book_chapter_步步高:",chapter_section)
                try:
                    chapter_section=str(int(chapter_section))
                except:
                    pass
                if chapter_section not in chapter_sections:
                    chapter_sections.append(chapter_section)
            print("chapter_sections:",chapter_sections)
            chapter_sections2 = []
            for chapter_section in chapter_sections:
                try:
                    chapter_sections2.append("第"+str(int(chapter_section))+"节")
                except:
                    chapter_sections2.append(chapter_section)
            # print("title_data", title_data)
            print("book_chapters3:",chapter_sections2)
        chapter_sections2.sort()
        combobox_section.config(values=chapter_sections2)#重新排序
        combobox_section.current(0)
    def book_select(*args):
        """选择书籍，相应的刷新章节"""
        title_data = pd.DataFrame(pd.read_excel(r"word题目汇总/题目清单.xlsx")).set_index("录入来源")
        print("combobox_book.get():", combobox_book.get())
        title_data = title_data.loc[combobox_book.get(), "章节"]
        book_chapters=[]
        for book_chapter in  title_data:
            print("book_chapter:",book_chapter)
            try:
                book_chapter=str(int(book_chapter))
            except:
                pass
            if book_chapter not in book_chapters:
                book_chapters.append(book_chapter )
        book_chapters_2=[]
        for book_chapter in book_chapters:
            if "考试" not in combobox_book.get():
                book_chapters_2.append("第"+book_chapter+"章")
            else:
                print("2345")
                # delta = pd.Timedelta(str(int(book_chapter)) + 'days')
                # time = pd.to_datetime('1899-12-30') + delta
                # 将数字 41957 转换为 Python 的日期时间对象
                import datetime
                python_datetime = datetime.datetime(1899, 12, 30) + datetime.timedelta(days=int(book_chapter))
                # time.strftime('%Y{y}%m{m}%d{d}').format(y='年', m='月', d='日')
                python_datetime = python_datetime.strftime('%Y{y}%m{m}').format(y='年', m='月')
                book_chapters_2.append(python_datetime)
        # print("title_data", title_data)
        print("book_chapters:",book_chapters_2)
        combobox_chapter.config(values=book_chapters_2)
        combobox_chapter.current(0)
        chapter_select()
    combobox_book['values'] = ("2023级考试","步步高2023级必修1","步步高2023级必修2","步步高2023级必修3",
                               "步步高2023级选择必修1","步步高2023级选择必修2","步步高2023级选择必修3",)  # 设置下拉框的选项列表
    combobox_book.bind('<<ComboboxSelected>>', book_select)
    combobox_book.current(1)

    combobox_chapter = ttk.Combobox(win)
    combobox_chapter.bind('<<ComboboxSelected>>',chapter_select)
    # combobox_class['values'] = ("12班")
    # combobox_class.current(0)
    combobox_chapter['values'] = ("第1章","第2章","第3章","第4章",
                                  "第5章","第6章","第7章","第8章", "第9章", "第10章", "第11章", "第12章", "第13章")  # 设置下拉框的选项列表
    combobox_chapter.current(0)

    # combobox_name = ttk.Combobox(win,postcommand=register)
    combobox_section = ttk.Combobox(win)
    combobox_section['values'] = ("第1节","第2节","第3节","第4节","第5节","第6节","章末素养提升",
        '专题1', '专题2', '专题3', '专题4', '专题5', '专题6', '专题7',
        '专题8', '专题9', '专题10', '专题11', '专题12', '专题13',
        '专题14', '专题15', '专题16', '专题17', '专题18', '专题19',
    )  # 设置下拉框的选项列表
    combobox_section.current(0)
    combobox_section.bind('<<ComboboxSelected>>', register_ti_fast)


    combobox_book.pack()
    combobox_chapter.pack()
    combobox_section.pack()



    # 新建
    # labe1 = Label(win, text="账号：")
    # labe2 = Label(win, text="密码：")
    # labe1.grid(row=0)
    # labe2.grid(row=1)
    # 创建动字符串
    label = tk.Label(win, text="输入要登记的题目：")
    label.pack()
    sv = tk.StringVar()
    sv.trace("w", input_title_test)
    entry_input = tk.Entry(win, textvariable=sv)
    entry_input.bind("<Return>",input_title_save)#回车响应
    entry_input.bind('<Key>', input_title_test)

    combobox_wrong_season = ttk.Combobox(win)
    combobox_wrong_season['values'] = ("选择错误原因","审题不清楚","研究对象不明确","受力分析错误","运动过程不清晰",
                                       "临界状态判断错误","隐含条件未找到","不知道用什么规律",
                                       "公式记忆错误","计算错误","没有思路","没有时间作答","其他"
                                       )  # 设置下拉框的选项列表
    combobox_wrong_season.current(0)
    combobox_wrong_season.bind('<<ComboboxSelected>>',register_wrong_season)
    # Py_String = tk.StringVar()
    # 使用验证参数 validata,参数值为 focusout 当失去焦点的时候，验证输入框内容是否正确

    # entry_input = tk.Entry(win, textvariable=Py_String, validate="focusout", )
    # entry_input = tk.Entry(win, textvariable=Py_String, validate="focusout", validatecommand=input_title)

    label_title_text = tk.Label(win, text="题目信息如下：")
    title_show = tk.Text(win, wrap=tk.WORD)
    # entry2 = tk.Entry(win)

    # entry=tk.Entry()
    entry_input.pack()
    combobox_wrong_season.pack()
    label_title_text.pack()

    #图片显示题目预览
    # image = Image.open(
    #     r"F:\个性化题库20180423\题不二错1.8\题不二错1.8测试版\图片题目汇总\A202311343KK4.jpg")  # 将"path_to_your_image"替换为你自己的图片路径
    # image=image.resize((100,100))
    # photo = ImageTk.PhotoImage(image)
    # photo=photo
    # import time
    # label = tk.Label(win, image=photo)
    # label.pack()
    # button_main=tk.Button(win, text="开始",command=main)
    # button_question_type_train=tk.Button(win, text="题型专练",command=question_type_train)
    button_network_disk=tk.Button(win, text="设置网盘",command=set_network_disk)
    button_print=tk.Button(win, text="保存word文档",command=butt_print)
    name=combobox_name.get()
    button_preview=tk.Button(win, text="预览", command=lambda: button_preview_title(name))
    button_edit=tk.Button(win, text="修改",command=edit_title)
    # button_wrong=tk.Button(win, text="录入错因",command=lambda: wrong_season(name))
    # wrong_title_text=tk.StringVar()
    button_admin=tk.Button(win, text="管理员",command=make_word_question_bank)
    button_zujuan=tk.Button(win, text="组卷",command=zujuan)
    # button_regist_ti_fast=tk.Button(win, text="快速录入题目",command=register_ti_fast)
    button_zujuan.pack(side="right")
    button_admin.pack(side="right")
    # button_question_type_train.pack(side="right")
    button_network_disk.pack(side="right")
    button_print.pack(side="right")
    button_preview.pack(side="right")
    button_edit.pack(side="right")
    # button_wrong.pack(side="right")
    # button_main.pack()
    title_show.pack()
    register()

    win.mainloop()
